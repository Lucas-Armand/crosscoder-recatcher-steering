#!/usr/bin/env python3
from pathlib import Path
from zipfile import ZipFile
import io,re
import matplotlib.pyplot as p
from matplotlib.patches import FancyBboxPatch,Rectangle
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE
from PIL import Image
R=Path("/home/lucas/crosscoder-recatcher-steering");SRC=R/"MODEL DIFFS_TARGET_VS_RANDOM_V5_2026-09-02.pptx";D=R/"MODEL DIFFS_INTERNAL_TABLE1_AUDIT_V6_2026-09-03.pptx";PDF=D.with_suffix(".pdf");O=R/"presentation/generated/table1_audit_v6_20260903";O.mkdir(parents=True,exist_ok=True)
BG="#F7F4EE";N="#10233F";P="#7257C7";OR="#EE8A2D";T="#159C9C";M="#687386";W="#FFFFFF";G="#D8D9D8";RED="#C93948";GR="#338A63";GO="#D4A72C";PR="#F8E5E7"
def setup(t,s):
 f=p.figure(figsize=(13.333,7.5),facecolor=BG);a=f.add_axes([0,0,1,1]);a.axis("off");a.text(.045,.955,"INTERNAL APPENDIX - TABLE 1 AUDIT",size=8.4,color=RED,weight="bold",va="top");a.text(.045,.885,t,size=21,color=N,weight="bold",va="top");a.text(.045,.815,s,size=9.7,color=M,va="top");a.text(.95,.035,"Internal audit appendix - 2026-09-03",ha="right",size=7.3,color=M);return f,a
def card(a,x,y,w,h,t,b,c,fs=8.2,face=W):
 a.add_patch(FancyBboxPatch((x,y),w,h,boxstyle="round,pad=.007",facecolor=face,edgecolor=G));a.add_patch(Rectangle((x,y+h-.009),w,.009,color=c));a.text(x+.015,y+h-.037,t,size=10,color=N,weight="bold",va="top");a.text(x+.015,y+h-.085,b,size=fs,color=M,va="top",linespacing=1.25)
def save(f,n):q=O/f"{n}.png";f.savefig(q,dpi=180,facecolor=BG);p.close(f);return q
im={}
f,a=setup("THE ORIGINAL CASE-STUDY RATIONALE MIXED BENCHMARKS","The improvement/regression framing matches published HumanEval+, not published BigCodeBench")
for x,title,ds,cl in [(.055,"PUBLISHED TABLE 1",[14.94,9.32],[-18.72,6.71]),(.525,"CANONICAL V4 STUDY",[21.34,11.93],[-23.17,-25.18])]:
 card(a,x,.37,.42,.38,title,"",P);a.text(x+.20,.64,"DeepSeek",ha="center",size=9,color=P,weight="bold");a.text(x+.34,.64,"CodeLlama",ha="center",size=9,color=OR,weight="bold")
 for i,r in enumerate(["HumanEval+","BigCodeBench"]):
  y=.55-i*.12;a.text(x+.02,y,r,size=8.5,color=N,weight="bold")
  for xx,v in [(x+.20,ds[i]),(x+.34,cl[i])]:a.text(xx,y,f"{v:+.2f} p.p.",ha="center",size=13,color=GR if v>0 else RED,weight="bold")
card(a,.055,.15,.89,.135,"WHAT WE MISREAD","We compared our CodeLlama BigCodeBench regression with the HumanEval+ regression in Table 1, creating an incorrect impression of agreement.",RED,8.8,PR);im[33]=save(f,33)
f,a=setup("THE DISCREPANCY MOTIVATED A CONTROLLED RE-EVALUATION","No regeneration: both paths receive the same archived BigCodeBench programs")
card(a,.05,.43,.25,.28,"SAME ARCHIVED OUTPUTS","1,140 tasks x 10 generations\nper model\n\n45,600 samples",N,9);card(a,.37,.53,.23,.18,"ORIGINAL PATH","Archived ReCatcher labels",OR,9);card(a,.37,.29,.23,.18,"AUDITED PATH","Extraction v4\nBigCodeBench 0.1.5",T,9)
a.annotate("",(.36,.60),(.30,.57),arrowprops=dict(arrowstyle="-|>",color=G,lw=2));a.annotate("",(.36,.38),(.30,.55),arrowprops=dict(arrowstyle="-|>",color=G,lw=2))
card(a,.67,.31,.28,.40,"DELTAS","DeepSeek\nTable 1  +9.32\narchive   +4.58\nv4           +1.11\n\nCodeLlama\nTable 1   +6.71\narchive    +6.71\nv4          -24.28",P,9.2)
a.text(.05,.18,"ALIGNMENT GATE: 45,600/45,600 prompts and programs match byte-for-byte by task_id + generation.",size=9.3,color=GR,weight="bold");a.text(.05,.11,"The comparison isolates extraction and evaluation effects from generation.",size=9.2,color=N,weight="bold");im[34]=save(f,34)
f,a=setup("TWO REPRODUCIBILITY DEFECTS WERE CONFIRMED","One changes correctness labels; the other misreports the open-weight generation limit")
card(a,.05,.28,.43,.48,"1  EXCEPTIONS COULD BECOME PASS",'except Exception as e:\n    q.put(str(e))\n\nif "FAILED" in test_output:\n    return False\nreturn True\n\nBigCodeBench/12  -  exp_0\nSyntaxError: unmatched ")"\narchived label = PASS',RED,8.6,PR)
card(a,.515,.28,.435,.48,"2  CONFIG DOES NOT MATCH EXECUTION","config.json: max_length = 256\n\nActual open-weight call:\nmax_new_tokens = 1024\n(max_length is commented out)\n\n2,184 DS-base and 6,178 CL-merged\ncontinuations exceed 256 tokens.\nCL-merged median is approx. 1,020.",GO,8.6)
a.text(.05,.19,"INVALID PROGRAMS LABELED PASS",size=9.2,color=N,weight="bold");a.text(.30,.19,"DS base 33  -  DS FT 425  -  CL base 194  -  CL merged 3,701",size=9.4,color=RED,weight="bold");a.text(.05,.105,"No systematic v4 defect was found in this audit; residual valid-code disagreements remain a separate review item.",size=9,color=GR,weight="bold");im[35]=save(f,35)
f,a=setup("RECOMMENDATION: INDEPENDENTLY RE-AUDIT ALL TABLE 1 CELLS","These artifacts reproduce the issues and support a complete review")
card(a,.045,.43,.285,.33,"SOURCE DATA","Zenodo 14997627\nGenerations + testing results\n\nzenodo.org/records/14997627",N,8.5);card(a,.357,.43,.285,.33,"ORIGINAL IMPLEMENTATION","ReCatcher/utils.py\ncode_generation/generate_code.py\ncode_generation/constants.py\n\ngithub.com/AltafAllahAbbassi/ReCatcher",OR,7.9);card(a,.67,.43,.285,.33,"AUDITED V4 PIPELINE","Audit report + CSV/JSON\ncode_extraction.py\nevaluate_bigcodebench_subset.py\npostprocessing_v4_validation.md\n\ngithub.com/Lucas-Armand/\ncrosscoder-recatcher-steering",T,7.5)
for i,t in enumerate(["Reconstruct every Table 1 cell from archived labels","Compile-gate every stored candidate before execution","Re-evaluate identical outputs with a pinned official harness","Separate extraction, execution, timeout, and dependency failures","Publish per-task/generation labels and effective configs"]):
 y=.315-i*.052;a.text(.065,y,str(i+1),ha="center",va="center",size=9,color=W,weight="bold",bbox=dict(boxstyle="circle,pad=.25",facecolor=RED,edgecolor=RED));a.text(.095,y,t,size=8.7,color=N,weight="bold",va="center")
a.text(.60,.16,"Decision needed: corrected Table 1 values\nand a single canonical evaluation pipeline.",size=9,color=RED,weight="bold",va="top");im[36]=save(f,36)
prs=Presentation(SRC);blank=prs.slide_layouts[6];links={35:[(.05,.28,.43,.48,"https://github.com/AltafAllahAbbassi/ReCatcher/blob/main/ReCatcher/utils.py#L258-L301"),(.515,.28,.435,.48,"https://github.com/AltafAllahAbbassi/ReCatcher/blob/main/code_generation/generate_code.py#L14-L27")],36:[(.045,.43,.285,.33,"https://zenodo.org/records/14997627"),(.357,.43,.285,.33,"https://github.com/AltafAllahAbbassi/ReCatcher"),(.67,.43,.285,.33,"https://github.com/Lucas-Armand/crosscoder-recatcher-steering/tree/presentation-v1/reports/recatcher_table1_v4_reconstruction_20260902")]}
for n in range(33,37):
 sl=prs.slides.add_slide(blank);sl.shapes.add_picture(str(im[n]),0,0,width=prs.slide_width,height=prs.slide_height)
 for x,y,w,h,u in links.get(n,[]):
  sh=sl.shapes.add_shape(MSO_SHAPE.RECTANGLE,int(x*prs.slide_width),int((1-y-h)*prs.slide_height),int(w*prs.slide_width),int(h*prs.slide_height));sh.fill.background();sh.line.fill.background();sh.click_action.hyperlink.address=u
prs.save(D)
with ZipFile(D) as z:
 pages=[]
 for n in range(1,len(prs.slides)+1):
  if n>=33:pages.append(Image.open(im[n]).convert("RGB"));continue
  rel=z.read(f"ppt/slides/_rels/slide{n}.xml.rels").decode();c=[]
  for t in re.findall(r'Target="([^"]+)"',rel):
   if t.startswith("../media/"):b=z.read("ppt/"+t[3:]);c.append((len(b),b))
  pages.append(Image.open(io.BytesIO(max(c)[1])).convert("RGB"))
pages[0].save(PDF,save_all=True,append_images=pages[1:],resolution=180)
print(D);print(PDF);print(len(prs.slides))
