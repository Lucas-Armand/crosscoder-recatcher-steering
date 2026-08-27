#!/usr/bin/env python3
from pathlib import Path
import io, zipfile, re, json, textwrap
import matplotlib.pyplot as plt
import numpy as np
from matplotlib.patches import FancyBboxPatch, Rectangle
from pptx import Presentation
from PIL import Image

REPO=Path('/home/lucas/crosscoder-recatcher-steering')
SRC=REPO/'MODEL DIFFS_ERROR_FOCUSED_ALPHA3_PARTIAL_2026-08-26.pptx'
DEST=REPO/'MODEL DIFFS_CAUSAL_STEERING_RESULTS_2026-08-27.pptx'
PDF=REPO/'MODEL DIFFS_CAUSAL_STEERING_RESULTS_2026-08-27.pdf'
OUT=REPO/'presentation/generated/causal_steering_results_20260827'
BG='#F7F4EE'; NAVY='#10233F'; PURPLE='#7257C7'; ORANGE='#EE8A2D'; TEAL='#159C9C'; MUTED='#687386'; WHITE='#FFFFFF'; GRID='#D8D9D8'; RED='#C93948'; GREEN='#338A63'; GOLD='#D4A72C'
OUT.mkdir(parents=True,exist_ok=True)

DS={2468:[3,4,7,7,7],2621:[3,4,6,5,8],15235:[2,3,6,5,4],2913:[3,6,5,6,11]}
DS_CHANGED={2468:[13,25,34,37,44],2621:[12,22,28,31,40],15235:[14,15,23,25,26],2913:[14,24,29,34,45]}
CL={7692:[0,1,1,1,1],10818:[0,0,1,1,1],5642:[0,0,1,1,1],11596:[0,1,1,1,1],4309:[0,0,1,2,2]}
CL_CHANGED={7692:[7,15,22,26,28],10818:[5,14,20,25,30],5642:[10,16,23,30,32],11596:[13,19,24,27,31],4309:[8,14,19,24,29]}

def setup(title,sub='',section='CAUSAL STEERING · CANONICAL SEED'):
 f=plt.figure(figsize=(13.333,7.5),facecolor=BG); a=f.add_axes([0,0,1,1]);a.axis('off')
 a.text(.045,.955,section,fontsize=8.3,color=TEAL,weight='bold',va='top')
 a.text(.045,.885,title,fontsize=22,color=NAVY,weight='bold',va='top')
 if sub:a.text(.045,.815,sub,fontsize=10,color=MUTED,va='top')
 a.text(.95,.035,'CrossCoder model diffing · partial validated readout · 2026-08-27',ha='right',fontsize=7.2,color=MUTED)
 return f,a

def card(a,x,y,w,h,title,body,accent,ts=10,bs=8.2):
 a.add_patch(FancyBboxPatch((x,y),w,h,boxstyle='round,pad=.007,rounding_size=.011',facecolor=WHITE,edgecolor='#E1DED7'))
 a.add_patch(Rectangle((x,y+h-.01),w,.01,color=accent,lw=0));a.text(x+.016,y+h-.04,title,fontsize=ts,color=NAVY,weight='bold',va='top')
 a.text(x+.016,y+h-.09,body,fontsize=bs,color=MUTED,va='top',linespacing=1.28)

def save(f,n):
 p=OUT/f'{n:02d}.png';f.savefig(p,dpi=180,facecolor=BG);plt.close(f);return p

imgs=[]
# 19
f,a=setup('THE α=3 SWEEP REVEALS BOTH MODEL AND TASK SENSITIVITY','Error-focused candidates · continuous last-token steering · official BigCodeBench 0.1.5')
card(a,.055,.37,.40,.36,'DEEPSEEK · CONTAMINATION IMPROVEMENTS','25/35 features evaluated\n83 feature-task transitions\n14/80 unique tasks corrected\n\nSeveral features correct multiple tasks.',PURPLE,12,9.2)
card(a,.49,.37,.40,.36,'CODELLAMA · LOGIC/RUNTIME REGRESSIONS','30/32 features evaluated\n6 feature-task transitions\n4/50 unique tasks corrected\n\nSingle-feature corrections are sparse.',ORANGE,12,9.2)
card(a,.055,.16,.835,.13,'READING','The sweep reduces 16,384 latents to causal candidates, but repeated correction of the same task also reveals task susceptibility. Dose curves and controls are required next.',TEAL,10,8.8)
imgs.append(save(f,19))

# 20
f,a=setup('FROM α=3 SCREENING TO CONTROLLED DOSE–RESPONSE','The confirmatory design separates trajectory control, outcome control, directionality, and specificity')
steps=[('1','α=0 GATE','byte-exact canonical reproduction'),('2','TARGET CURVE','|α| = 1 · 2 · 3 · 4 · 5'),('3','REVERSE TEST','opposite model and direction'),('4','CONTROLS','3 random latents + 3 orthogonal shams'),('5','TWO ENDPOINTS','evaluated code change + official pass')]
for i,(num,t,b) in enumerate(steps):
 x=.045+i*.19; card(a,x,.42,.17,.29,num+'  '+t,b,[PURPLE,TEAL,ORANGE,GOLD,GREEN][i],9.2,8)
card(a,.055,.18,.835,.14,'INTERVENTION','Layer 16 · selected decoder side · last residual position · added at every autoregressive generation step · temperature 0.2 · top-p 0.95 · max 512 · seed = 1000 + task_idx × 100',NAVY,9.5,8.5)
a.text(.055,.11,'α=0 passed: DeepSeek 80/80 byte-exact · CodeLlama 50/50 byte-exact.',fontsize=9.5,color=GREEN,weight='bold')
imgs.append(save(f,20))

# 21
f,a=setup('STEERING CHANGES TRAJECTORIES—AND DEEPSEEK OUTCOMES','Complete direct target curves; α=3 integrated from the canonical sweep')
ax=f.add_axes([.07,.30,.40,.43]);
for fid,v in DS.items():ax.plot([1,2,3,4,5],v,marker='o',lw=2,label=str(fid))
ax.set_title('DeepSeek · official fail→pass / 80',loc='left',color=NAVY,weight='bold');ax.set_xlabel('|α|');ax.set_ylabel('passes');ax.set_xticks([1,2,3,4,5]);ax.grid(color=GRID,alpha=.6);ax.spines[['top','right']].set_visible(False);ax.legend(frameon=False,ncol=2,fontsize=7)
ax=f.add_axes([.54,.30,.40,.43]);
for fid,v in CL.items():ax.plot([1,2,3,4,5],v,marker='o',lw=2,label=str(fid))
ax.set_title('CodeLlama · official fail→pass / 50',loc='left',color=NAVY,weight='bold');ax.set_xlabel('|α|');ax.set_ylabel('passes');ax.set_xticks([1,2,3,4,5]);ax.set_ylim(-.1,2.4);ax.grid(color=GRID,alpha=.6);ax.spines[['top','right']].set_visible(False);ax.legend(frameon=False,ncol=2,fontsize=7)
a.text(.07,.19,'DEEPSEEK',fontsize=9,color=PURPLE,weight='bold');a.text(.16,.19,'16 unique tasks corrected; peak 11/80.',fontsize=8.5,color=NAVY)
a.text(.54,.19,'CODELLAMA',fontsize=9,color=ORANGE,weight='bold');a.text(.65,.19,'5 primary tasks corrected; peak 2/50.',fontsize=8.5,color=NAVY)
a.text(.07,.115,'Specificity versus the new random/sham controls remains pending; causal trajectory and outcome effects are already directly measured.',fontsize=8.7,color=RED,weight='bold')
imgs.append(save(f,21))

# 22
f,a=setup('FEATURE 2468: THE LEADING DEEPSEEK MECHANISTIC CANDIDATE','Screening strength and stable causal response align around the contamination boundary')
ax=f.add_axes([.08,.30,.44,.43]);ax.plot([1,2,3,4,5],DS[2468],marker='o',lw=3,color=PURPLE,label='official pass');ax.plot([1,2,3,4,5],DS_CHANGED[2468],marker='s',lw=2,color='#B8A7E5',label='output changed');ax.set_xticks([1,2,3,4,5]);ax.set_xlabel('|α|, negative on base');ax.set_ylabel('tasks / 80');ax.grid(color=GRID,alpha=.6);ax.spines[['top','right']].set_visible(False);ax.legend(frameon=False)
card(a,.57,.49,.36,.24,'SCREENING','paired local · rank #1\nbase-enriched · max\nE/V = −7.29\nsupport: 62/80',PURPLE,10,8.5)
card(a,.57,.25,.36,.19,'CAUSAL READOUT','7/80 passes at α=−3, −4 and −5\n34 → 44 outputs changed',GREEN,10,8.5)
a.text(.57,.15,'Meaning ↔ failure mode ↔ activation timing',fontsize=10,color=NAVY,weight='bold')
a.text(.57,.105,'Interpretation: a stable lever over late post-solution behavior.',fontsize=8.7,color=MUTED)
imgs.append(save(f,22))

# 23
f,a=setup('DEEPSEEK /7: STEERING REMOVES POST-SOLUTION TEST CONTAMINATION','Feature 2468 · α=−4 · official fail→pass')
card(a,.055,.30,.42,.43,'BASE · FAIL','return top_product\n\n# test_task.py\nimport unittest\nfrom task import task_func\n\nclass TestTask(unittest.TestCase):\n    ...',RED,11,8.7)
card(a,.52,.30,.42,.43,'STEERED · PASS','return top_product\n\n# task_func("path/to/sales.csv")\n# task_func("sales.csv")\n\n[no executable test harness]',GREEN,11,8.7)
card(a,.055,.13,.885,.105,'MECHANISTIC READING','The useful function is preserved; the principal intervention occurs after the solution boundary and removes executable test/import contamination. This is a clean outcome-level example, not proof that every 2468 correction uses the same mechanism.',PURPLE,9.5,8.2)
imgs.append(save(f,23))

# 24
f,a=setup('FEATURE 2468 ALSO MAKES COHERENT CHANGES THAT DO NOT REPAIR THE TASK','Semantic control is not equivalent to sufficient correction')
card(a,.055,.39,.27,.34,'/97 · STILL FAILS','Only expected floating-point constants inside generated tests change.\n\nTrajectory changed late, but the irrelevant test harness remains.',PURPLE,10,8.2)
card(a,.365,.39,.27,.34,'/823 · STILL FAILS','Generated test precision and comments change near the end.\n\nThe intervention affects post-solution material without repairing the evaluated function.',PURPLE,10,8.2)
card(a,.675,.39,.27,.34,'WHY THIS MATTERS','The feature is behaviorally coherent, but cleanup is neither universal nor sufficient. Some tasks still contain the original logical failure.',ORANGE,10,8.2)
a.text(.055,.20,'Observed:',fontsize=9,color=GREEN,weight='bold');a.text(.13,.20,'directional control over late generated content.',fontsize=8.8,color=NAVY)
a.text(.055,.14,'Not implied:',fontsize=9,color=RED,weight='bold');a.text(.14,.14,'a single-feature explanation for all 80 failures.',fontsize=8.8,color=NAVY)
imgs.append(save(f,24))

# 25
f,a=setup('FEATURE 4309: CODELLAMA’S STRONGEST CURRENT OUTCOME CANDIDATE','Trajectory effects grow with dose; official repairs remain rare')
ax=f.add_axes([.08,.30,.44,.43]);ax.plot([1,2,3,4,5],CL[4309],marker='o',lw=3,color=ORANGE,label='official pass');ax.plot([1,2,3,4,5],CL_CHANGED[4309],marker='s',lw=2,color='#F4BE84',label='output changed');ax.set_xticks([1,2,3,4,5]);ax.set_xlabel('|α|, negative on merged');ax.set_ylabel('tasks / 50');ax.grid(color=GRID,alpha=.6);ax.spines[['top','right']].set_visible(False);ax.legend(frameon=False)
card(a,.57,.49,.36,.24,'OUTCOME','0 · 0 · 1 · 2 · 2 passes\n8 → 29 outputs changed\nprimary successes include /183 and /490',ORANGE,10,8.5)
card(a,.57,.25,.36,.19,'BOUNDARY CONDITION','More trajectory changes do not translate proportionally into correct programs.',RED,10,8.5)
a.text(.57,.145,'Current interpretation: distributed regression mechanisms.',fontsize=9,color=NAVY,weight='bold')
imgs.append(save(f,25))

# 26
f,a=setup('CODELLAMA /490: STEERING REPAIRS PARSING, WRITING, AND RETURN','Feature 4309 · α=−5 · official fail→pass')
card(a,.055,.30,.42,.43,'MERGED · FAIL','return json.loads(s)\n\n[does not parse XML correctly]\n[does not write the requested JSON]\n[one generated function returns nothing]',RED,11,8.7)
card(a,.52,.30,.42,.43,'STEERED · PASS','result = xmltodict.parse(s)\nwith open(file_path, "w") as f:\n    json.dump(result, f)\nreturn result',GREEN,11,8.7)
card(a,.055,.13,.885,.105,'MECHANISTIC READING','The intervention produces a semantically aligned multi-part repair. Because CodeLlama successes are sparse, this is evidence for selected-task causal control—not a broad single-feature account of merged-model regressions.',ORANGE,9.5,8.2)
imgs.append(save(f,26))

# 27
f,a=setup('CODELLAMA CHANGES CODE FREQUENTLY WITHOUT REPAIRING IT','Feature 4309 exposes trajectory control and limited outcome specificity')
card(a,.055,.39,.27,.34,'/738 · STILL FAILS','The intervention removes a closing parenthesis in multiple doctest examples.\n\nThe code changes consistently, but in a harmful or irrelevant way.',ORANGE,10,8.2)
card(a,.365,.39,.27,.34,'DOSE RESPONSE','Outputs changed:\n8 · 14 · 19 · 24 · 29\n\nOfficial passes:\n0 · 0 · 1 · 2 · 2',ORANGE,10,9)
card(a,.675,.39,.27,.34,'INTERPRETATION','The latent can move the trajectory, but the 50-task logic/runtime cohort contains heterogeneous failure mechanisms.',RED,10,8.2)
a.text(.055,.16,'Causal effect: supported.',fontsize=9.5,color=GREEN,weight='bold');a.text(.25,.16,'Broad repair mechanism: not supported.',fontsize=9.5,color=RED,weight='bold')
imgs.append(save(f,27))

# 28
f,a=setup('THE EVIDENCE CHAIN MAKES EACH CLAIM AUDITABLE','Association, interpretation, causality, and specificity remain separate layers',section='SYNTHESIS')
labels=['TRANSITION','ERROR-FOCUSED\nSCREEN','TOKENS +\nTIMING','DIRECTION','α=0 GATE','DOSE CURVE','SEMANTIC\nCHANGE','FAIL→PASS','CONTROLS']
cols=[MUTED,PURPLE,PURPLE,TEAL,GREEN,TEAL,ORANGE,GREEN,RED]
for i,(lab,col) in enumerate(zip(labels,cols)):
 x=.035+i*.104;a.add_patch(FancyBboxPatch((x,.46),.09,.16,boxstyle='round,pad=.006',facecolor=WHITE,edgecolor=col,lw=2));a.text(x+.045,.54,lab,ha='center',va='center',fontsize=7.6,color=NAVY,weight='bold')
 if i<8:a.annotate('',xy=(x+.105,.54),xytext=(x+.093,.54),arrowprops=dict(arrowstyle='->',color=MUTED))
card(a,.055,.20,.40,.16,'ASSOCIATIVE → INTERPRETATIVE','transition · screening · tokens/timing',PURPLE,9.5,8.3)
card(a,.50,.20,.40,.16,'CAUSAL → SPECIFICITY','gate · intervention · code/outcome · controls',GREEN,9.5,8.3)
imgs.append(save(f,28))

# 29
f,a=setup('WHAT THE CURRENT RESULTS SUPPORT—AND WHAT THEY DO NOT','Claims are tiered by evidential strength',section='SYNTHESIS')
card(a,.055,.27,.42,.46,'SUPPORTED NOW','• screening compresses 16,384 latents\n• α=0 reproduces canonical generations\n• selected features causally alter code\n• some interventions yield repeated official passes\n• DeepSeek has stronger single-feature outcome control\n• CodeLlama is a useful boundary condition',GREEN,11,8.8)
card(a,.52,.27,.42,.46,'NOT YET SUPPORTED','• every transition is feature-caused\n• every target beats new random/sham controls\n• one feature explains heterogeneous CodeLlama errors\n• 14175 has a complete predicted-direction curve\n• generalization across new seeds/tasks',RED,11,8.8)
a.text(.055,.15,'Specificity is a result to be demonstrated—not inferred from screening or one successful task.',fontsize=10,color=NAVY,weight='bold')
imgs.append(save(f,29))

# 30
f,a=setup('A REUSABLE METHODOLOGY FOR CAUSAL MODEL DIFFING','The output is a falsifiable mechanism hypothesis—not merely a feature ranking',section='SYNTHESIS')
items=['Paired behavioral transitions','Error taxonomy + focused cohort','Same-text sparse representation','Conditioned E/V screening','Meaning + timing + failure mode','Directional hypothesis','Canonical α=0 gate','Dose + reverse steering','Official outcome + controls']
for i,t in enumerate(items):
 r=i//3;c=i%3;x=.055+c*.305;y=.63-r*.19;card(a,x,y,.27,.14,f'{i+1}',t,[PURPLE,TEAL,ORANGE][c],10,8.2)
a.text(.055,.10,'Every stage can reject, refine, or narrow the hypothesis before a broad mechanistic claim is made.',fontsize=9.4,color=NAVY,weight='bold')
imgs.append(save(f,30))

# 31
f,a=setup('MODEL DIFFERENCES CAN BECOME TESTABLE CAUSAL HYPOTHESES','A shared representation connects behavioral contrast to controlled intervention',section='CONCLUSION')
a.text(.075,.62,'Behavioral differences',fontsize=18,color=NAVY,weight='bold');a.text(.39,.62,'→',fontsize=24,color=TEAL,weight='bold');a.text(.47,.62,'shared sparse features',fontsize=18,color=PURPLE,weight='bold');a.text(.78,.62,'→',fontsize=24,color=TEAL,weight='bold')
a.text(.075,.47,'interpretable timing + semantics',fontsize=18,color=ORANGE,weight='bold');a.text(.50,.47,'→',fontsize=24,color=TEAL,weight='bold');a.text(.58,.47,'causal steering',fontsize=18,color=GREEN,weight='bold')
card(a,.075,.20,.80,.16,'SUPPORTED CONCLUSION','Shared CrossCoder features can localize model differences, generate directional causal hypotheses, and—in selected cases—produce reproducible behavioral transitions.',NAVY,11,10)
a.text(.075,.12,'DeepSeek: more localized outcome control · CodeLlama: causal trajectory control, but more distributed failures.',fontsize=9.3,color=MUTED)
imgs.append(save(f,31))

# Replace slides 19 onward in the latest deck.
prs=Presentation(SRC)
while len(prs.slides)>18:
 sid=prs.slides._sldIdLst[-1]; prs.part.drop_rel(sid.rId); prs.slides._sldIdLst.remove(sid)
for p in imgs:
 s=prs.slides.add_slide(prs.slide_layouts[6]);s.shapes.add_picture(str(p),0,0,width=prs.slide_width,height=prs.slide_height)
prs.save(DEST)

# Full PDF: recover the largest raster from each preserved slide, then append generated slides.
pages=[]
with zipfile.ZipFile(DEST) as z:
 for n in range(1,19):
  xml=z.read(f'ppt/slides/slide{n}.xml').decode(); rel=z.read(f'ppt/slides/_rels/slide{n}.xml.rels').decode()
  candidates=[]
  for rid,target in re.findall(r'Id="(rId\d+)"[^>]*Target="([^"]+)"',rel):
   if '/media/' not in target and not target.startswith('../media/'): continue
   name='ppt/'+re.sub(r'^\.\./','',target); blob=z.read(name); candidates.append((len(blob),blob))
  if not candidates: raise RuntimeError(f'no raster for slide {n}')
  im=Image.open(io.BytesIO(max(candidates)[1])).convert('RGB');pages.append(im)
for p in imgs:pages.append(Image.open(p).convert('RGB'))
pages[0].save(PDF,save_all=True,append_images=pages[1:],resolution=150)
print(DEST);print(PDF);print('slides',len(prs.slides),'pdf_pages',len(pages))
