#!/usr/bin/env python3
"""Build compact partial alpha=3 results and insert them before superseded slides."""
from pathlib import Path
import argparse,csv,glob,json,os,re,shutil,statistics
from collections import Counter
import matplotlib.pyplot as plt
from matplotlib.backends.backend_pdf import PdfPages
from matplotlib.patches import FancyBboxPatch,Rectangle
from pptx import Presentation

BG="#F7F4EE";NAVY="#10233F";PURPLE="#7257C7";ORANGE="#EE8A2D";TEAL="#159C9C";MUTED="#687386";WHITE="#FFFFFF";GRID="#D8D9D8";RED="#C93948";GREEN="#338A63"

def setup(title,sub=""):
 f=plt.figure(figsize=(13.333,7.5),facecolor=BG);a=f.add_axes([0,0,1,1]);a.axis("off")
 a.text(.045,.955,"CANONICAL ERROR-FOCUSED α=3 SWEEP · PARTIAL READOUT",fontsize=8.5,color=TEAL,weight="bold",va="top")
 a.text(.045,.885,title,fontsize=23,color=NAVY,weight="bold",va="top")
 if sub:a.text(.045,.815,sub,fontsize=10.2,color=MUTED,va="top")
 a.text(.95,.035,"CrossCoder model diffing · canonical seed · 2026-08-26",ha="right",fontsize=7.5,color=MUTED)
 return f,a

def card(a,x,y,w,h,title,body,accent,ts=10,bs=8):
 a.add_patch(FancyBboxPatch((x,y),w,h,boxstyle="round,pad=.007,rounding_size=.011",facecolor=WHITE,edgecolor="#E1DED7"))
 a.add_patch(Rectangle((x,y+h-.01),w,.01,color=accent,lw=0));a.text(x+.016,y+h-.04,title,fontsize=ts,color=NAVY,weight="bold",va="top")
 a.text(x+.016,y+h-.085,body,fontsize=bs,color=MUTED,va="top",linespacing=1.3)

def read_family(name,run,eval_dir,screen,n,total):
 run=Path(run);base={r["task_id"]:r for r in map(json.loads,open(run/"generations/bigcodebench__baseline_alpha0_results.jsonl"))}
 arm={int(r["feature_id"]):r for r in csv.DictReader(open(run/"ARM_MANIFEST.csv"))};ranks={}
 for p in Path(screen).glob("*_absolute.csv"):
  for r in csv.DictReader(open(p)):ranks.setdefault(int(r["feature_id"]),[]).append((int(r["rank"]),p.stem,float(r["abs_ev"])))
 pat=re.compile(r"(?im)^\s*(#\s*)?(tests?/|test_|import\s+(pytest|unittest)|from\s+(pytest|unittest))")
 rows=[];tasks=Counter()
 for ep in glob.glob(str(Path(eval_dir)/"*_eval.json")):
  m=re.search(r"__f(\d+)_",Path(ep).name)
  if not m:continue
  fid=int(m.group(1));ev=json.load(open(ep));passed={t for t,v in ev["eval"].items() if v[0].get("correct")};tasks.update(passed)
  gp=next(run.glob(f"generations/bigcodebench__f{fid}_*_results.jsonl"));gr=list(map(json.loads,gp.open()))
  if len(gr)!=n:continue
  changed=clean=0;positions=[]
  for r in gr:
   before=base[r["task_id"]]["completion"];after=r["completion"]
   if before==after:continue
   changed+=1;i=0
   while i<min(len(before),len(after)) and before[i]==after[i]:i+=1
   positions.append(i/max(1,len(before)))
   if name=="DeepSeek" and pat.search(before) and not pat.search(after):clean+=1
  best=min(ranks[fid]);rows.append({"model":name,"feature":fid,"passes":len(passed),"changed":changed,"cleanup":clean,
   "mean_first_divergence_changed_pct":100*statistics.mean(positions) if positions else 100,
   "median_first_divergence_changed_pct":100*statistics.median(positions) if positions else 100,
   "ev":float(arm[fid]["ev"]),"abs_ev":float(arm[fid]["abs_ev"]),"orientation":arm[fid]["orientation"],
   "summary":arm[fid]["summary"],"best_screen":best[1].replace("_absolute",""),"best_screen_rank":best[0],"pass_tasks":sorted(passed)})
 rows.sort(key=lambda r:(-r["passes"],-r["abs_ev"],r["feature"]))
 return {"model":name,"n":n,"features_done":len(rows),"features_total":total,"baseline_passes":0,"rows":rows,
  "feature_task_transitions":sum(r["passes"] for r in rows),"unique_tasks":len(tasks),"task_counts":tasks}

def write_report(data,out):
 out.mkdir(parents=True,exist_ok=True)
 fields=["model","feature","passes","changed","cleanup","mean_first_divergence_changed_pct","median_first_divergence_changed_pct","ev","abs_ev","orientation","summary","best_screen","best_screen_rank","pass_tasks"]
 with (out/"feature_results.csv").open("w",newline="") as f:
  w=csv.DictWriter(f,fields);w.writeheader()
  for d in data:
   for r in d["rows"]:w.writerow({**r,"pass_tasks":";".join(r["pass_tasks"])})
 with (out/"task_susceptibility.csv").open("w",newline="") as f:
  w=csv.writer(f);w.writerow(["model","task_id","features_correcting"])
  for d in data:
   for t,n in d["task_counts"].most_common():w.writerow([d["model"],t,n])
 compact=[{k:v for k,v in d.items() if k!="task_counts"}|{"task_counts":dict(d["task_counts"])} for d in data]
 (out/"summary.json").write_text(json.dumps(compact,indent=2)+"\n")
 lines=["# Canonical error-focused alpha=3 partial results",""]
 for d in data:
  lines+= [f"## {d['model']}",f"- Features evaluated: {d['features_done']}/{d['features_total']}",f"- Baseline: 0/{d['n']}",f"- Feature-task pass transitions: {d['feature_task_transitions']}",f"- Unique corrected tasks: {d['unique_tasks']}/{d['n']}",""]
 (out/"README.md").write_text("\n".join(lines)+"\n")

def build(data,out):
 ds,cl=data;out.mkdir(parents=True,exist_ok=True);imgs=[]
 with PdfPages(out/"alpha3_partial_insert.pdf") as pdf:
  f,a=setup("THE PARTIAL RESULT SEPARATES A STRONGER AND A WEAKER REGIME","Official BigCodeBench 0.1.5 evaluation; α=3; continuous last-token steering")
  for x,d,col in [(.06,ds,PURPLE),(.53,cl,ORANGE)]:
   card(a,x,.38,.41,.35,d["model"].upper(),f"Features evaluated      {d['features_done']}/{d['features_total']}\nBaseline                 0/{d['n']}\nFeature-task transitions {d['feature_task_transitions']}\nUnique corrected tasks   {d['unique_tasks']}/{d['n']}\n\nAll results use the canonical per-task seed rule.",col,12,10)
  a.text(.06,.25,"PROMISING",fontsize=11,color=PURPLE,weight="bold");a.text(.06,.20,"DeepSeek contains several multi-task causal candidates.",fontsize=9,color=NAVY)
  a.text(.53,.25,"WEAKER / DIFFUSE",fontsize=11,color=ORANGE,weight="bold");a.text(.53,.20,"CodeLlama yields sparse, task-specific corrections.",fontsize=9,color=NAVY)
  a.text(.06,.105,"PROVISIONAL · DeepSeek is paused at 28/35; CodeLlama is complete at 32/32.",fontsize=10,color=RED,weight="bold")
  p=out/"13_partial_overview.png";f.savefig(p,dpi=175,facecolor=BG);pdf.savefig(f);plt.close(f);imgs.append(p)

  f,a=setup("DEEPSEEK: THE BEST FEATURES PRODUCE MULTIPLE OFFICIAL CORRECTIONS","Causal outcome first; screening E/V breaks ties")
  top=ds["rows"][:10];ys=list(range(len(top)))[::-1]
  ax=f.add_axes([.08,.18,.42,.55]);ax.barh(ys,[r["changed"] for r in top],color="#D9D0F1",label="output changed");ax.barh(ys,[r["passes"] for r in top],color=PURPLE,label="official pass")
  ax.set_yticks(ys,[str(r["feature"]) for r in top]);ax.set_xlabel("tasks / 80");ax.spines[["top","right"]].set_visible(False);ax.legend(frameon=False,fontsize=8)
  card(a,.55,.22,.39,.51,"TOP 5", "\n".join(f"{i+1}. {r['feature']}   {r['passes']}/80 pass   {r['changed']}/80 changed   E/V {r['ev']:+.2f}" for i,r in enumerate(ds["rows"][:5])),PURPLE,11,8.4)
  a.text(.55,.15,"Changing many outputs is not sufficient:",fontsize=9,color=NAVY,weight="bold");a.text(.55,.105,"14295 changed 35/80 but produced only 2 official passes.",fontsize=8.8,color=RED)
  p=out/"14_deepseek_features.png";f.savefig(p,dpi=175,facecolor=BG);pdf.savefig(f);plt.close(f);imgs.append(p)

  f,a=setup("DEEPSEEK: 93 TRANSITIONS COLLAPSE TO 14 SUSCEPTIBLE TASKS","Feature-task outcomes are not independent causal successes")
  top_tasks=ds["task_counts"].most_common(14);ax=f.add_axes([.08,.19,.55,.57]);y=list(range(len(top_tasks)))[::-1];ax.barh(y,[n for _,n in top_tasks],color=PURPLE);ax.set_yticks(y,[t.split("/")[-1] for t,_ in top_tasks]);ax.set_xlabel("features that corrected the task");ax.spines[["top","right"]].set_visible(False)
  card(a,.68,.42,.27,.30,"CONCENTRATION",f"/435   {ds['task_counts']['BigCodeBench/435']} features\n/316   {ds['task_counts']['BigCodeBench/316']} features\n/166   {ds['task_counts']['BigCodeBench/166']} features\n/496   {ds['task_counts']['BigCodeBench/496']} features",PURPLE,10,9)
  card(a,.68,.18,.27,.17,"INTERPRETATION","The screen finds causal levers, but some tasks are broadly steering-susceptible. Shams and matched controls remain necessary.",RED,9.5,7.8)
  p=out/"15_deepseek_concentration.png";f.savefig(p,dpi=175,facecolor=BG);pdf.savefig(f);plt.close(f);imgs.append(p)

  r=next(x for x in ds["rows"] if x["feature"]==2468)
  f,a=setup("FEATURE 2468 IS THE LEADING DEEPSEEK CANDIDATE","Meaning ↔ failure mode ↔ activation timing must remain distinct")
  card(a,.055,.39,.28,.34,"SCREENING","paired local · rank 1\nE/V "+f"{r['ev']:+.2f}"+f"\nsummary: {r['summary']}\nselected around the late divergence window",PURPLE,11,9)
  card(a,.36,.39,.28,.34,"CAUSAL OUTCOME",f"7/80 official passes\n34/80 outputs changed\n13/80 contamination cleanups\nbaseline: 0/80",GREEN,11,9)
  card(a,.665,.39,.28,.34,"CAUSED TIMING",f"Among changed outputs:\nmean first divergence: {r['mean_first_divergence_changed_pct']:.1f}%\nmedian: {r['median_first_divergence_changed_pct']:.1f}%\n\nContinuous steering can redirect earlier than the natural activation window.",ORANGE,11,8.3)
  a.text(.06,.24,"SUPPORTED",fontsize=10,color=GREEN,weight="bold");a.text(.16,.24,"local screening nominated a feature with repeated official causal effects.",fontsize=9,color=NAVY)
  a.text(.06,.17,"NOT YET SUPPORTED",fontsize=10,color=RED,weight="bold");a.text(.22,.17,"feature-specificity versus sham/matched latent controls.",fontsize=9,color=NAVY)
  p=out/"16_feature2468.png";f.savefig(p,dpi=175,facecolor=BG);pdf.savefig(f);plt.close(f);imgs.append(p)

  f,a=setup("CODELLAMA: SINGLE-FEATURE CONTROL IS WEAK AND TASK-SPECIFIC","Complete 32-feature sweep; baseline 0/50")
  passes=[r for r in cl["rows"] if r["passes"]];card(a,.055,.37,.43,.37,"SIX FEATURES · ONE PASS EACH","\n".join(f"{r['feature']}  →  {r['pass_tasks'][0]}" for r in passes),ORANGE,11,8.5)
  card(a,.53,.37,.42,.37,"FOUR UNIQUE TASKS","/119 corrected by 3 features\n/823 corrected by 1\n/490 corrected by 1\n/630 corrected by 1\n\n26/32 features produced no official correction.",ORANGE,11,9)
  a.text(.055,.24,"Trajectory changes occur earlier and more diffusely than in the contamination-focused DeepSeek cohort.",fontsize=9.2,color=NAVY,weight="bold")
  a.text(.055,.17,"Boundary condition: correlation and trajectory control do not imply that one latent can repair a heterogeneous merged-model regression.",fontsize=9,color=RED)
  p=out/"17_codellama_results.png";f.savefig(p,dpi=175,facecolor=BG);pdf.savefig(f);plt.close(f);imgs.append(p)

  f,a=setup("NEXT TARGETS COMBINE CAUSAL EFFECT WITH SCREENING STRENGTH","Ordering: official passes first; |E/V| breaks ties")
  for x,d,col in [(.055,ds,PURPLE),(.525,cl,ORANGE)]:
   a.text(x,.75,d["model"].upper(),fontsize=12,color=col,weight="bold")
   cols=[("RANK",.055),("FEATURE",.13),("PASS",.235),("BEST SCREEN",.31),("|E/V|",.445)]
   for label,dx in cols:a.text(x+dx,.69,label,fontsize=7,color=MUTED,weight="bold")
   for i,r in enumerate(d["rows"][:5]):
    y=.625-i*.085;a.add_patch(Rectangle((x,y-.025),.41,.067,facecolor=WHITE,edgecolor=GRID))
    vals=[str(i+1),str(r["feature"]),f"{r['passes']}/{d['n']}",f"{r['best_screen'].replace('_',' ')} · #{r['best_screen_rank']}",f"{r['abs_ev']:.2f}"]
    for val,(_,dx) in zip(vals,cols):a.text(x+dx,y,val,fontsize=7.8,color=NAVY,weight="bold" if i==0 else "normal")
  a.text(.055,.11,"These are candidates for full dose curves plus sham and randomly selected/matched latent controls—not confirmatory winners yet.",fontsize=9.5,color=RED,weight="bold")
  p=out/"18_next_top5.png";f.savefig(p,dpi=175,facecolor=BG);pdf.savefig(f);plt.close(f);imgs.append(p)
 return imgs

def insert(prs,index,p):
 s=prs.slides.add_slide(prs.slide_layouts[6]);s.shapes.add_picture(str(p),0,0,width=prs.slide_width,height=prs.slide_height)
 sid=prs.slides._sldIdLst[-1];prs.slides._sldIdLst.remove(sid);prs.slides._sldIdLst.insert(index,sid)

def main():
 ap=argparse.ArgumentParser();ap.add_argument("--source",type=Path,required=True);ap.add_argument("--dest",type=Path,required=True);ap.add_argument("--output-dir",type=Path,required=True);ap.add_argument("--report-dir",type=Path,required=True);a=ap.parse_args()
 ds=read_family("DeepSeek","runs/focused_subtype_dstk100_alpha3_canonical_v1","/tmp/focused_alpha3_deck_snapshot_20260826/dstk/eval","reports/focused_subtype_screening_dstk100_contamination_v1",80,35)
 cl=read_family("CodeLlama","runs/focused_subtype_codellama_alpha3_canonical_v1","runs/focused_subtype_codellama_alpha3_canonical_v1/evaluations","reports/focused_subtype_screening_codellama_wrong_logic_v1",50,32)
 data=[ds,cl];write_report(data,a.report_dir);imgs=build(data,a.output_dir);shutil.copy2(a.source,a.dest);prs=Presentation(a.dest)
 for i,p in enumerate(imgs):insert(prs,12+i,p)
 prs.save(a.dest);print(a.dest,len(prs.slides))

if __name__=="__main__":main()
