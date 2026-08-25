#!/usr/bin/env python3
"""Build the clean 11-slide CrossCoder team-story draft."""
from pathlib import Path
import argparse,csv,json,re
import numpy as np
import matplotlib.pyplot as plt
from matplotlib.backends.backend_pdf import PdfPages
from matplotlib.patches import Rectangle,FancyBboxPatch
from pptx import Presentation
from pptx.util import Inches

BG="#F7F4EE"; NAVY="#10233F"; PURPLE="#7257C7"; ORANGE="#EE8A2D"; TEAL="#159C9C"
MUTED="#687386"; WHITE="#FFFFFF"; GRID="#D8D9D8"; GREY="#D7D6D2"; RED="#D85A56"; GOLD="#DDA531"

def setup(title,sub="",section="CROSSCODER MODEL DIFFING"):
    f=plt.figure(figsize=(13.333,7.5),facecolor=BG);a=f.add_axes([0,0,1,1]);a.axis("off")
    a.text(.045,.955,section,fontsize=8.5,color=TEAL,weight="bold",va="top")
    a.text(.045,.885,title,fontsize=24,color=NAVY,weight="bold",va="top")
    if sub:a.text(.045,.815,sub,fontsize=10.5,color=MUTED,va="top")
    a.text(.95,.035,"CrossCoder model diffing · team discussion draft · 2026-08-25",ha="right",fontsize=7.5,color=MUTED)
    return f,a

def card(a,x,y,w,h,title,body="",accent=TEAL,face=WHITE,ts=11,bs=8.7):
    a.add_patch(FancyBboxPatch((x,y),w,h,boxstyle="round,pad=0.007,rounding_size=.011",facecolor=face,edgecolor="#E1DED7"))
    a.add_patch(Rectangle((x,y+h-.01),w,.01,color=accent,linewidth=0))
    a.text(x+.018,y+h-.045,title,fontsize=ts,color=NAVY,weight="bold",va="top")
    if body:a.text(x+.018,y+h-.092,body,fontsize=bs,color=MUTED,va="top",linespacing=1.35)

def arrow(a,x1,y1,x2,y2,color=TEAL):
    a.annotate("",xy=(x2,y2),xytext=(x1,y1),arrowprops=dict(arrowstyle="->",lw=2,color=color))

def save(f,p,pdf):
    f.savefig(p,dpi=170,facecolor=BG);pdf.savefig(f,facecolor=BG);plt.close(f)

def table(a,data,heads,widths,x,y,w,h,color=NAVY,hi=None,fs=7.8):
    rh=h/(len(data)+1);xx=[x]
    for q in widths:xx.append(xx[-1]+w*q)
    a.add_patch(Rectangle((x,y+h-rh),w,rh,color=color))
    for j,v in enumerate(heads):a.text(xx[j]+.005,y+h-rh/2,v,color=WHITE,fontsize=fs,weight="bold",va="center")
    for i,row in enumerate(data):
        yy=y+h-(i+2)*rh;a.add_patch(Rectangle((x,yy),w,rh,color=WHITE if i%2==0 else "#EFEDE8"))
        for j,v in enumerate(row):a.text(xx[j]+.005,yy+rh/2,str(v),color=TEAL if j==hi else NAVY,fontsize=fs,weight="bold" if j==hi else "normal",va="center")

def read10(p):
    with open(p,newline="") as f:return list(csv.DictReader(f))[:10]

def eval_counts(p):
    out=[]
    for q in p.glob("evaluations/bigcodebench__f*_eval.json"):
        if re.search(r"__f\d+_(?:pos|neg)3_eval",q.name):out.append(json.load(open(q))["passed"])
    return out

def build(repo,out):
    out.mkdir(parents=True,exist_ok=True);imgs=[];pdf_path=out/"MODEL_DIFFS_NEW_STORY_DRAFT_2026-08-25.pdf"
    with PdfPages(pdf_path) as pdf:
        # 1
        f=plt.figure(figsize=(13.333,7.5),facecolor=BG);a=f.add_axes([0,0,1,1]);a.axis("off")
        a.add_patch(Rectangle((0,0),.026,1,color=TEAL));a.text(.07,.85,"CROSSCODER MODEL DIFFING",fontsize=10,color=TEAL,weight="bold")
        a.text(.07,.68,"FROM MODEL DIFFS",fontsize=34,color=NAVY,weight="bold");a.text(.07,.57,"TO CAUSAL CONTROL",fontsize=34,color=NAVY,weight="bold")
        a.text(.07,.43,"A shared sparse representation for explaining and testing\nbehavioral differences between code models",fontsize=16,color=MUTED,linespacing=1.35)
        a.add_patch(FancyBboxPatch((.70,.25),.20,.44,boxstyle="round,pad=.02,rounding_size=.025",facecolor=WHITE,edgecolor=GRID))
        for i,(t,c) in enumerate([("BEHAVIOR",PURPLE),("FEATURES",TEAL),("STEERING",ORANGE)]):
            yy=.59-i*.12;a.text(.80,yy,t,ha="center",fontsize=11,color=c,weight="bold")
            if i<2:arrow(a,.80,yy-.035,.80,yy-.085,c)
        a.text(.94,.04,"1",ha="right",fontsize=8,color=MUTED)
        p=out/"01_cover.png";save(f,p,pdf);imgs.append(p)

        # 2
        f,a=setup("RESEARCH QUESTION","Evaluation tells us what changed. CrossCoder asks which internal differences can explain—and control—it.")
        card(a,.06,.27,.27,.43,"BEHAVIOR","Which tasks changed?\n\npass → fail\nfail → pass\n\nWhat failure mode changed?",PURPLE)
        card(a,.365,.27,.27,.43,"REPRESENTATION","Paired layer-16 residuals\n\nSame-text alignment\n\nShared sparse features",TEAL)
        card(a,.67,.27,.27,.43,"CAUSAL TEST","Choose candidate directions\n\nIntervene during generation\n\nCompare against controls",ORANGE)
        arrow(a,.335,.49,.36,.49);arrow(a,.64,.49,.665,.49)
        a.text(.06,.15,"Behavior localizes the model difference. Sparse features turn it into a testable hypothesis.",fontsize=11,color=NAVY,weight="bold")
        p=out/"02_question.png";save(f,p,pdf);imgs.append(p)

        # 3
        f,a=setup("ONE ANALYSIS FRAMEWORK, TWO VERY DIFFERENT MODEL CHANGES","The same same-text CrossCoder design spans fine-tuning and model merging")
        for yy,name,left,right,col,note in [
          (.48,"DEEPSEEK","BASE","FINE-TUNED",PURPLE,"specialization through fine-tuning"),
          (.18,"CODELLAMA","BASE","MERGED",ORANGE,"specialization through model merging")]:
            a.text(.055,yy+.18,name,fontsize=12,color=col,weight="bold");a.text(.055,yy+.135,note,fontsize=8.5,color=MUTED)
            for x,t in [(.25,left),(.25,right)]:
                pass
            card(a,.23,yy,.16,.18,left+"\nlayer 16","same evaluated token IDs",col,WHITE,10,8)
            card(a,.43,yy,.16,.18,right+"\nlayer 16","same evaluated token IDs",col,WHITE,10,8)
            arrow(a,.39,yy+.09,.425,yy+.09,col);arrow(a,.59,yy+.09,.635,yy+.09,col)
            card(a,.64,yy,.24,.18,"CROSSCODER","16,384 latents · ReLU\nexact TopK-100",col,WHITE,10,8)
        a.text(.055,.105,"Different specialization mechanisms; identical analytical question: which shared latent differences track—and control—behavior?",fontsize=9.5,color=NAVY,weight="bold")
        p=out/"03_two_models.png";save(f,p,pdf);imgs.append(p)

        # 4
        f,a=setup("THE TWO SPECIALIZATIONS MOVE PERFORMANCE IN OPPOSITE DIRECTIONS","Canonical extraction-v4 labels · same raw generations, audited reprocessing")
        def result(x,y,title,base,var,total,col,label):
            card(a,x,y,.40,.39,title,"",col)
            a.text(x+.03,y+.255,f"{base}/{total}",fontsize=24,color=NAVY,weight="bold");a.text(x+.03,y+.215,"BASE",fontsize=8,color=MUTED)
            a.text(x+.22,y+.255,f"{var}/{total}",fontsize=24,color=col,weight="bold");a.text(x+.22,y+.215,label,fontsize=8,color=MUTED)
            pp=(var-base)/total*100;a.text(x+.03,y+.105,f"{pp:+.2f} pp",fontsize=20,color=col,weight="bold")
            a.text(x+.22,y+.11,"BigCodeBench\npass-rate change",fontsize=8.5,color=MUTED)
        result(.06,.34,"DEEPSEEK · BIGCODEBENCH",268,404,1140,PURPLE,"FINE-TUNED")
        result(.54,.34,"CODELLAMA · BIGCODEBENCH",314,27,1140,ORANGE,"MERGED")
        a.text(.06,.24,"DeepSeek: 215 improvements · 79 regressions",fontsize=9.5,color=PURPLE,weight="bold")
        a.text(.54,.24,"CodeLlama: 4 improvements · 291 regressions",fontsize=9.5,color=ORANGE,weight="bold")
        a.add_patch(FancyBboxPatch((.06,.09),.88,.105,boxstyle="round,pad=.01",facecolor=WHITE,edgecolor=GRID))
        a.text(.08,.155,"RECATCHER CONTEXT",fontsize=8.5,color=TEAL,weight="bold")
        a.text(.08,.115,"Directionally consistent motivation: specialization can improve aggregate capability while still creating localized regressions.",fontsize=9,color=NAVY)
        a.text(.08,.085,"Not a direct numerical replication: ReCatcher used 10 generations/prompt and different regression definitions.",fontsize=8,color=MUTED)
        p=out/"04_behavior.png";save(f,p,pdf);imgs.append(p)

        # 5
        f,a=setup("FEATURE SCREENING","We rank stable behavioral contrasts—not raw activation alone")
        stages=[
          ("4 SUMMARIES","max · early_max\nmean · active fraction"),
          ("Δ PER TASK","aggregate(specialized text)\n− aggregate(base text)"),
          ("OBSERVED EFFECT","mean Δ positives\n− mean Δ controls"),
          ("200 PERMUTATIONS","shuffle labels\nrecompute effect"),
          ("E/V","effect ÷ null SD") ]
        for i,(h,b) in enumerate(stages):
            x=.045+i*.19;card(a,x,.36,.16,.27,h,b,TEAL if i<4 else GOLD,WHITE,9.5,8.2)
            if i<4:arrow(a,x+.162,.495,x+.187,.495)
        a.text(.06,.25,"Minimum support",fontsize=9,color=MUTED);a.text(.06,.205,"max(3 tasks, 10% of positives)",fontsize=10,color=NAVY,weight="bold")
        a.text(.48,.25,"Interpretation",fontsize=9,color=MUTED);a.text(.48,.205,"E/V is permutation signal-to-noise—not a z-score or p-value.",fontsize=10,color=NAVY,weight="bold")
        a.text(.06,.11,"Positive and negative orientations are retained: either enrichment direction may support behavior-aligned steering.",fontsize=9.5,color=TEAL,weight="bold")
        p=out/"05_screen_method.png";save(f,p,pdf);imgs.append(p)

        # 6
        f,a=setup("EIGHT COMPLEMENTARY SCREENING CELLS PER MODEL","Transition type × comparison design × temporal scope")
        for px,name,col in [(.045,"DEEPSEEK",PURPLE),(.515,"CODELLAMA",ORANGE)]:
            a.text(px,.755,name,fontsize=13,color=col,weight="bold")
            for trans,yy in [("REGRESSION",.46),("IMPROVEMENT",.18)]:
                off=name=="CODELLAMA" and trans=="IMPROVEMENT"
                card(a,px,yy,.43,.245,trans,"",GREY if off else col,"#E7E6E2" if off else WHITE,10)
                for j,(u,v) in enumerate([("ASSOCIATION","global · local"),("PAIRED MODEL Δ","global · local")]):
                    x=px+.018+j*.205;a.add_patch(Rectangle((x,yy+.035),.19,.105,facecolor="#F0EFEC" if off else BG,edgecolor=GRID))
                    a.text(x+.012,yy+.105,u,fontsize=8,color=MUTED if off else NAVY,weight="bold");a.text(x+.012,yy+.065,v,fontsize=8,color=MUTED)
                if off:a.text(px+.215,yy+.188,"ONLY 4 POSITIVES · EXCLUDED",ha="center",fontsize=8,color=RED,weight="bold")
        a.text(.045,.095,"Local = ±10% around the first normalized divergence—not a semantic error category.",fontsize=9,color=NAVY,weight="bold")
        p=out/"06_cells.png";save(f,p,pdf);imgs.append(p)

        # 7
        ds=read10(repo/"reports/eight_cell_screening_dstk100_v1/improvement_association_global_base_enriched.csv")
        cl=read10(repo/"reports/eight_cell_screening_codellama_base_merged_v1/regression_association_global_absolute.csv")
        f,a=setup("WHAT THE SCREENING ACTUALLY RETURNS","Representative valid cells; E/V is the primary ranking signal")
        for x,title,rows,col in [(.045,"DEEPSEEK · IMPROVEMENT ASSOCIATION",ds,PURPLE),(.515,"CODELLAMA · REGRESSION ASSOCIATION",cl,ORANGE)]:
            a.text(x,.75,title,fontsize=10,color=col,weight="bold")
            data=[[r["rank"],r["feature_id"],f'{float(r["effect"]):+.3f}',f'{float(r["ev"]):+.2f}'] for r in rows]
            table(a,data,["Rank","Feature","Effect","E/V"],[.18,.28,.27,.27],x,.18,.43,.52,col,3,8.2)
        a.text(.045,.105,"Examples—not a single global ranking. Candidate pools are formed by union and deduplication across valid cells.",fontsize=9,color=NAVY,weight="bold")
        p=out/"07_top10.png";save(f,p,pdf);imgs.append(p)

        # 8
        f,a=setup("SCREEN FIRST; TEST CAUSALLY AT A COMMON OPERATING POINT","The statistical shortlist reduces the intervention search by ~99.8%")
        for x,big,small,col in [(.10,"16,384","latents / CrossCoder",NAVY),(.40,"≈30–40","candidates / model",TEAL),(.70,"≈0.2%","of latent space",RED)]:
            a.text(x,.63,big,fontsize=36,color=col,weight="bold",ha="center");a.text(x,.56,small,fontsize=9,color=MUTED,ha="center")
        arrow(a,.18,.61,.30,.61);arrow(a,.48,.61,.60,.61)
        card(a,.07,.20,.39,.25,"WHY |α| = 3?","Earlier dose-response work suggested a practical exploratory point: visible causal changes below the most disruptive high-dose regime.",GOLD)
        card(a,.54,.20,.39,.25,"WHY STANDARDIZE?","One magnitude makes dozens of candidates comparable before full dose curves, reverse-model tests, and controls.",TEAL)
        a.text(.07,.11,"|α|=3 is an operating point—not an independently validated optimum.",fontsize=9.5,color=RED,weight="bold")
        p=out/"08_funnel.png";save(f,p,pdf);imgs.append(p)

        # 9
        f,a=setup("OUR CURRENT INTERVENTION IS CONTINUOUS STEERING","The selected decoder direction is added at every autoregressive generation step")
        xs=np.linspace(.13,.83,6)
        a.text(.07,.61,"PROMPT",fontsize=10,color=NAVY,weight="bold")
        arrow(a,.105,.61,xs[0]-.035,.61,NAVY)
        for i,x in enumerate(xs):
            a.add_patch(FancyBboxPatch((x-.03,.56),.06,.10,boxstyle="round,pad=.006",facecolor=WHITE,edgecolor=PURPLE))
            a.text(x,.61,f"t{i+1}",ha="center",va="center",fontsize=10,color=NAVY,weight="bold")
            a.text(x,.72,"+ αdⱼ",ha="center",fontsize=10,color=ORANGE,weight="bold");arrow(a,x,.70,x,.665,ORANGE)
            if i<5:arrow(a,x+.035,.61,xs[i+1]-.035,.61,NAVY)
        card(a,.07,.20,.40,.23,"CURRENT PROTOCOL","Layer 16 · last residual position\nfeature decoder on intervened model side\napplied through the full generation",PURPLE)
        card(a,.53,.20,.40,.23,"WHAT IT IS NOT","Not a one-token intervention\nnot activation-gated or clamped\nnot conditional on natural activation",ORANGE)
        a.text(.07,.105,"This tests continuous control of the generation trajectory—not whether one naturally active token alone caused the error.",fontsize=9.5,color=NAVY,weight="bold")
        p=out/"09_steering_protocol.png";save(f,p,pdf);imgs.append(p)

        # 10
        ds_c=eval_counts(repo/"runs/semantic_top10_dstk100_alpha3_v1");cl_c=eval_counts(repo/"runs/semantic_top10_codellama_alpha3_v1")
        f,a=setup("THE α=3 SWEEP REVEALS BOTH MODEL AND TASK SENSITIVITY","Official pass changes motivate—but do not replace—random and sham controls")
        for rect,vals,base,title,col in [([.055,.31,.39,.40],ds_c,15,"DeepSeek · 30 features",PURPLE),([.475,.31,.39,.40],cl_c,0,"CodeLlama · 37 features",ORANGE)]:
            q=f.add_axes(rect,facecolor=WHITE);d=[v-base for v in vals];lo,hi=min(d),max(d)
            q.hist(d,bins=np.arange(lo-.5,hi+1.5),color=col,edgecolor=WHITE,rwidth=.86);q.axvline(0,color=NAVY,ls="--",lw=1.3)
            q.set_xticks(range(lo,hi+1));q.set_xlabel("Net official pass change",fontsize=8);q.set_ylabel("Features",fontsize=8);q.set_title(title,loc="left",fontsize=10,color=NAVY,weight="bold");q.spines[["top","right"]].set_visible(False);q.grid(axis="y",color=GRID,alpha=.5)
        card(a,.875,.49,.10,.22,"TASK\nSENSITIVITY","DS /1030\n29/30\n\nCL /119\n10/11",RED,WHITE,9,7.5)
        card(a,.055,.12,.27,.12,"TARGET","screen-selected latent",TEAL,WHITE,8.5,7.5)
        card(a,.355,.12,.27,.12,"RANDOM LATENTS","not selected by screen/α3",GOLD,WHITE,8.5,7.5)
        card(a,.655,.12,.27,.12,"ORTHOGONAL SHAMS","norm-controlled directions",MUTED,WHITE,8.5,7.5)
        a.text(.055,.075,"Controls are always good practice; concentrated successes make their necessity especially visible here.",fontsize=8.7,color=NAVY,weight="bold")
        p=out/"10_alpha3.png";save(f,p,pdf);imgs.append(p)

        # 11
        f,a=setup("TEN CANDIDATES ADVANCE TO BBASV","Causal outcome at |α|=3 first; applicable screening rank breaks ties")
        data=[
          ["DeepSeek","3048","−","FT","+4","19/80"],["DeepSeek","13801","−","FT","+3","18/80"],["DeepSeek","7828","−","FT","+3","18/80"],["DeepSeek","15669","−","FT","+3","18/80"],["DeepSeek","13191","−","FT","+2","17/80"],
          ["CodeLlama","13147","−","base","+1","1/50"],["CodeLlama","12253","+","base","+1","1/50"],["CodeLlama","10570","−","base","+1","1/50"],["CodeLlama","14359","+","base","+1","1/50"],["CodeLlama","2310","+","base","+1","1/50"]]
        table(a,data,["Model","Feature","Repair sign","Reverse model","Net Δ","Passes"],[.20,.14,.16,.20,.14,.16],.045,.20,.54,.52,NAVY,None,7.6)
        card(a,.62,.48,.33,.24,"BBASV · DIRECT","Move the worse model toward the behavior-associated direction.\n\nFull dose curve for each Top-5 feature.",TEAL)
        card(a,.62,.22,.33,.22,"BBASV · REVERSE + CONTROLS","Invert direction on the better model.\n\nCompare random latents and orthogonal shams.",ORANGE)
        a.text(.045,.12,"DeepSeek colors: purple · CodeLlama colors: orange",fontsize=8.5,color=MUTED)
        a.text(.62,.12,"Bidirectional Behavior-Aligned Steering Validation",fontsize=9.5,color=NAVY,weight="bold")
        p=out/"11_bbasv.png";save(f,p,pdf);imgs.append(p)
    return imgs,pdf_path

def main():
    ap=argparse.ArgumentParser();ap.add_argument("--repo",type=Path,default=Path("."));ap.add_argument("--dest",type=Path,required=True);ap.add_argument("--output-dir",type=Path,required=True)
    x=ap.parse_args();imgs,pdf=build(x.repo,x.output_dir);prs=Presentation();prs.slide_width=Inches(13.333);prs.slide_height=Inches(7.5);blank=prs.slide_layouts[6]
    while len(prs.slides):pass
    for p in imgs:
        s=prs.slides.add_slide(blank);s.shapes.add_picture(str(p),0,0,width=prs.slide_width,height=prs.slide_height)
    prs.save(x.dest);print(pdf)

if __name__=="__main__":main()
