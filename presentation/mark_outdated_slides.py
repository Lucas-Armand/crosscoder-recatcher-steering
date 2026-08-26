#!/usr/bin/env python3
"""Mark superseded result slides with a red tint and warning banner."""
import argparse
from pathlib import Path
from lxml import etree
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches, Pt

def main():
 ap=argparse.ArgumentParser()
 ap.add_argument("--deck",type=Path,required=True)
 ap.add_argument("--first-slide",type=int,default=13)
 x=ap.parse_args()
 prs=Presentation(x.deck)
 for number,slide in enumerate(prs.slides,1):
  if number < x.first_slide: continue
  overlay=slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,0,0,prs.slide_width,prs.slide_height)
  overlay.fill.solid();overlay.fill.fore_color.rgb=RGBColor(185,22,38);overlay.line.fill.background()
  color=overlay.fill._xPr.solidFill[0]
  alpha=etree.SubElement(color,"{http://schemas.openxmlformats.org/drawingml/2006/main}alpha")
  alpha.set("val","36000")
  banner=slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,0,0,prs.slide_width,Inches(.58))
  banner.fill.solid();banner.fill.fore_color.rgb=RGBColor(145,0,20);banner.line.fill.background()
  p=banner.text_frame.paragraphs[0]
  p.text="OUTDATED · superseded candidate pool and/or non-canonical seed"
  p.font.name="Aptos Display";p.font.size=Pt(17);p.font.bold=True;p.font.color.rgb=RGBColor(255,255,255)
  p.alignment=1
 prs.save(x.deck)
 print(f"marked slides {x.first_slide}-{len(prs.slides)} in {x.deck}")

if __name__=="__main__":
 main()
