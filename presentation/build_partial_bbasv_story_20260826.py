#!/usr/bin/env python3
"""Append provisional BBASV result/story slides to the clean 11-slide deck."""
from pathlib import Path
import argparse, csv, shutil, subprocess
from collections import Counter
import matplotlib.pyplot as plt
import numpy as np
from matplotlib.backends.backend_pdf import PdfPages
from matplotlib.patches import FancyBboxPatch, Rectangle
from pptx import Presentation

BG="#F7F4EE"; NAVY="#10233F"; PURPLE="#7257C7"; ORANGE="#EE8A2D"; TEAL="#159C9C"
MUTED="#687386"; WHITE="#FFFFFF"; GRID="#D8D9D8"; RED="#D85A56"; GOLD="#DDA531"; GREEN="#338A63"

def setup(title, sub="", section="BBASV · PROVISIONAL RESULTS"):
    f=plt.figure(figsize=(13.333,7.5),facecolor=BG); a=f.add_axes([0,0,1,1]); a.axis("off")
    a.text(.045,.955,section,fontsize=8.5,color=TEAL,weight="bold",va="top")
    a.text(.045,.885,title,fontsize=23,color=NAVY,weight="bold",va="top")
    if sub:a.text(.045,.815,sub,fontsize=10.2,color=MUTED,va="top")
    a.text(.95,.035,"PROVISIONAL · evaluated arms as of 2026-08-26",ha="right",fontsize=7.5,color=RED,weight="bold")
    return f,a

def card(a,x,y,w,h,title,body="",accent=TEAL,face=WHITE,ts=10.5,bs=8.3):
    a.add_patch(FancyBboxPatch((x,y),w,h,boxstyle="round,pad=.007,rounding_size=.011",facecolor=face,edgecolor="#E1DED7"))
    a.add_patch(Rectangle((x,y+h-.01),w,.01,color=accent,linewidth=0))
    a.text(x+.016,y+h-.04,title,fontsize=ts,color=NAVY,weight="bold",va="top")
    if body:a.text(x+.016,y+h-.085,body,fontsize=bs,color=MUTED,va="top",linespacing=1.32)

def save(f,p,pdf):
    f.savefig(p,dpi=175,facecolor=BG); pdf.savefig(f,facecolor=BG); plt.close(f)

def rows(path):
    with open(path,newline="") as f:return list(csv.DictReader(f))

def data(repo,project):
    root=repo/("runs/dstk100_bbasv_v1" if project=="deepseek" else "runs/codellama_base_merged_bbasv_v1")/"analysis_at_stop"/"analysis"
    return rows(root/"official_transition_summary.csv"), rows(root/"arm_change_summary.csv"), rows(root/"task_change_details.csv")

def select(rs,kind="target",side="direct"):
    return [r for r in rs if r["kind"]==kind and r["side"]==side]

def feature_curves(ax,rs,features,side,color,metric="net_pass",pending=True):
    for i,ft in enumerate(features):
        q=sorted([r for r in rs if int(r["feature"])==ft and r["side"]==side and r["kind"]=="target"],key=lambda z:abs(float(z["alpha"])))
        if q:
            x=[abs(float(r["alpha"])) for r in q]; y=[float(r[metric]) for r in q]
            ax.plot(x,y,marker="o",lw=2,alpha=.9,label=str(ft),color=plt.cm.Purples(.42+i*.1) if color==PURPLE else plt.cm.Oranges(.42+i*.1))
            if pending and max(x)<5:
                ax.plot([max(x),5],[y[-1],y[-1]],ls=":",lw=1,color=RED,alpha=.55)
    ax.axhline(0,color=NAVY,lw=1,ls="--");ax.set_xticks([1,2,3,4,5]);ax.grid(color=GRID,alpha=.55);ax.spines[["top","right"]].set_visible(False)

def build(repo,out):
    out.mkdir(parents=True,exist_ok=True); imgs=[]
    ds,ds_change,ds_tasks=data(repo,"deepseek"); cl,cl_change,cl_tasks=data(repo,"codellama")
    dsf=[3048,13801,7828,15669,13191]; clf=[13147,12253,10570,14359,2310]
    with PdfPages(out/"MODEL_DIFFS_PARTIAL_BBASV_STORY_2026-08-26.pdf") as pdf:
        # 11
        f,a=setup("BBASV TURNS ASSOCIATIONS INTO CONTROLLED CAUSAL TESTS","Bidirectional Behavior-Aligned Steering Validation")
        stages=[("TOP-5 / MODEL","screen + α=3",PURPLE),("DIRECT","worse model\nbehavior-aligned",TEAL),("REVERSE","better model\nopposite direction",ORANGE),("DOSE","|α| = 1…5",GOLD),("CONTROLS","3 random + 3 sham",RED),("OFFICIAL EVAL","paired by task",GREEN)]
        for i,(t,b,c) in enumerate(stages):
            x=.045+i*.155;card(a,x,.42,.135,.24,t,b,c,WHITE,9,7.6)
            if i<5:a.annotate("",xy=(x+.153,.54),xytext=(x+.137,.54),arrowprops=dict(arrowstyle="->",color=MUTED,lw=1.5))
        card(a,.055,.18,.27,.14,"CODELLAMA","5/5 target curves · 1 random\n3 random + 3 sham still pending",ORANGE,WHITE,9,7.6)
        card(a,.365,.18,.27,.14,"DEEPSEEK","2 partial/full target curves + α=3 anchors\nremaining targets and all controls pending",PURPLE,WHITE,9,7.6)
        card(a,.675,.18,.27,.14,"READING RULE","Normal color = observed\nRed = missing / provisional",RED,WHITE,9,7.6)
        p=out/"11_bbasv_design.png";save(f,p,pdf);imgs.append(p)

        # 12
        f,a=setup("STEERING CHANGES TRAJECTORIES MORE OFTEN THAN OUTCOMES","Two aligned endpoints: evaluated-code change and official fail↔pass")
        for rect,title,rr,cc,features,col in [([.06,.44,.40,.28],"DEEPSEEK · DIRECT",ds,ds_change,dsf,PURPLE),([.54,.44,.40,.28],"CODELLAMA · DIRECT",cl,cl_change,clf,ORANGE)]:
            ax=f.add_axes(rect,facecolor=WHITE)
            # code-change counts from arm change table
            for i,ft in enumerate(features):
                q=sorted([r for r in cc if r.get("kind")=="target" and r.get("side")=="direct" and int(r["feature_id"])==ft],key=lambda z:abs(float(z["alpha"])))
                if q:ax.plot([abs(float(r["alpha"])) for r in q],[100*float(r["changed"])/(80 if col==PURPLE else 50) for r in q],marker="o",alpha=.75,label=str(ft))
            ax.set_title(title+" · code changed",loc="left",fontsize=9.5,color=NAVY,weight="bold");ax.set_ylabel("% evaluated code changed",fontsize=8);ax.set_xticks([1,2,3,4,5]);ax.grid(color=GRID,alpha=.5);ax.spines[["top","right"]].set_visible(False)
        for rect,title,rr,features,col in [([.06,.12,.40,.23],"Official net pass change",ds,dsf,PURPLE),([.54,.12,.40,.23],"Official net pass change",cl,clf,ORANGE)]:
            ax=f.add_axes(rect,facecolor=WHITE);feature_curves(ax,rr,features,"direct",col);ax.set_title(title,loc="left",fontsize=9.5,color=NAVY,weight="bold");ax.set_ylabel("fail→pass − pass→fail",fontsize=8);ax.set_xlabel("|α|",fontsize=8)
        a.text(.50,.755,"TRAJECTORY EFFECT",ha="center",fontsize=8,color=TEAL,weight="bold")
        a.text(.50,.385,"TASK OUTCOME",ha="center",fontsize=8,color=GREEN,weight="bold")
        p=out/"12_change_and_outcome_curves.png";save(f,p,pdf);imgs.append(p)

        # 13
        f,a=setup("DEEPSEEK: EARLY RESULTS SHOW POSITIVE BUT TASK-CONCENTRATED CONTROL","Direct baseline 15/80 · reverse baseline 57/80 · 25/45 non-baseline target arms evaluated")
        ax=f.add_axes([.06,.38,.43,.34],facecolor=WHITE);feature_curves(ax,ds,dsf,"direct",PURPLE);ax.set_title("Direct · steer base toward FT-associated behavior",loc="left",fontsize=10,color=NAVY,weight="bold");ax.set_ylabel("net official pass change");ax.set_xlabel("|α|");ax.legend(ncol=2,fontsize=7,frameon=False)
        ax=f.add_axes([.55,.38,.39,.34],facecolor=WHITE);feature_curves(ax,ds,dsf,"reverse",PURPLE);ax.set_title("Reverse · steer FT oppositely",loc="left",fontsize=10,color=NAVY,weight="bold");ax.set_ylabel("net official pass change");ax.set_xlabel("|α|");ax.legend(ncol=2,fontsize=7,frameon=False)
        card(a,.06,.15,.27,.15,"FAVORS THE STORY","13801: +2,+2,+3,+3,+3 direct\nand −1/−2 net in reverse",GREEN,WHITE,9,7.6)
        card(a,.365,.15,.27,.15,"QUALIFIES IT","/1030 and /706 recur across targets:\npart of the gain is task susceptibility",GOLD,WHITE,9,7.6)
        card(a,.67,.15,.27,.15,"PENDING","3048 and 7828 dose curves\nall random and sham controls",RED,WHITE,9,7.6)
        p=out/"13_deepseek_official.png";save(f,p,pdf);imgs.append(p)

        # 14
        f,a=setup("DEEPSEEK FEATURE 3048: A COHERENT TEST/IMPORT-BOUNDARY CANDIDATE","Semantic evidence and one clean α=−3 fail→pass example")
        toks=[("_",15296),("test",3088),("Func",3084),("from",2467),("Task",2306),("import",1907),("def",646)]
        ax=f.add_axes([.06,.43,.34,.30],facecolor=WHITE);ax.barh([x[0] for x in toks][::-1],[x[1] for x in toks][::-1],color=PURPLE);ax.set_title("Activation-weighted top tokens",loc="left",fontsize=10,color=NAVY,weight="bold");ax.spines[["top","right"]].set_visible(False);ax.grid(axis="x",color=GRID,alpha=.5)
        card(a,.44,.43,.50,.30,"BIGCODEBENCH/762 · FAIL → PASS","BASE: broad try/except + appended pytest suite\n\nSTEERED α=−3: direct implementation + return zip name;\npost-solution tests removed\n\nOfficial transition: fail → pass",GREEN,WHITE,10,8.2)
        card(a,.06,.17,.27,.16,"MEANING","test/import/function-boundary contexts",PURPLE,WHITE,9,7.5)
        card(a,.365,.17,.27,.16,"FAILURE MODE","post-solution contamination and\nover-defensive implementation",ORANGE,WHITE,9,7.5)
        card(a,.67,.17,.27,.16,"CAUSAL CHANGE","removes tests and simplifies body\nin a task that becomes correct",TEAL,WHITE,9,7.5)
        a.text(.06,.105,"Interpretation is coherent, but 3048 specificity versus controls remains PENDING.",fontsize=9.2,color=RED,weight="bold")
        p=out/"14_deepseek_f3048_success.png";save(f,p,pdf);imgs.append(p)

        # 15
        f,a=setup("FEATURE 3048 OFTEN CHANGES THE RIGHT REGION WITHOUT SOLVING THE TASK","At α=−3, 23/80 evaluated codes changed; only 4 became official passes")
        q=[r for r in ds_tasks if int(r["feature_id"])==3048 and r["side"]=="direct" and abs(float(r["alpha"]))==3 and r["changed"]=="1"]
        c=Counter(r["category"] for r in q); labels=[x for x,_ in c.most_common()]; vals=[c[x] for x in labels]
        ax=f.add_axes([.06,.36,.43,.36],facecolor=WHITE);ax.barh(labels[::-1],vals[::-1],color=PURPLE);ax.set_xlabel("changed tasks");ax.set_title("First-change taxonomy",loc="left",fontsize=10,color=NAVY,weight="bold");ax.spines[["top","right"]].set_visible(False);ax.grid(axis="x",color=GRID,alpha=.5)
        card(a,.55,.49,.39,.23,"WHAT CHANGED","imports/tests/comments/function structure\nfrequently move—the expected semantic neighborhood",TEAL)
        card(a,.55,.24,.39,.20,"WHY MOST STILL FAIL","Removing contamination cannot repair every\nunderlying logic/API/runtime defect",GOLD)
        a.text(.06,.14,"23 code changes → 4 fail→pass",fontsize=19,color=NAVY,weight="bold")
        a.text(.55,.14,"Causal scope ≠ complete task solution",fontsize=11,color=RED,weight="bold")
        p=out/"15_deepseek_f3048_limits.png";save(f,p,pdf);imgs.append(p)

        # 16
        f,a=setup("CODELLAMA: COMPLETE TARGET CURVES SHOW WEAK OUTCOME SPECIFICITY","All 50 target direct/reverse arms evaluated · one random latent evaluated · shams pending")
        ax=f.add_axes([.06,.39,.42,.34],facecolor=WHITE);feature_curves(ax,cl,clf,"direct",ORANGE);feature_curves(ax,cl,[13879],"direct",ORANGE,pending=False);ax.set_title("Direct · merged model (baseline 0/50)",loc="left",fontsize=10,color=NAVY,weight="bold");ax.set_ylabel("net official pass change");ax.set_xlabel("|α|");ax.legend(ncol=2,fontsize=7,frameon=False)
        ax=f.add_axes([.54,.39,.40,.34],facecolor=WHITE);feature_curves(ax,cl,clf,"reverse",ORANGE);feature_curves(ax,cl,[13879],"reverse",ORANGE,pending=False);ax.set_title("Reverse · base model (baseline 36/50)",loc="left",fontsize=10,color=NAVY,weight="bold");ax.set_ylabel("net official pass change");ax.set_xlabel("|α|");ax.legend(ncol=2,fontsize=7,frameon=False)
        card(a,.06,.15,.27,.15,"TARGET RESULT","mostly 0 or +1; 10570 reaches +2\nat |α|=5",ORANGE,WHITE,9,7.5)
        card(a,.365,.15,.27,.15,"CONTROL RESULT","random 13879 reaches +1 direct\nand +2 reverse",RED,WHITE,9,7.5)
        card(a,.67,.15,.27,.15,"CURRENT READING","trajectory control exists; outcome\nspecificity is not established",GOLD,WHITE,9,7.5)
        p=out/"16_codellama_official.png";save(f,p,pdf);imgs.append(p)

        # 17
        f,a=setup("CODELLAMA FEATURE 13147: A CLEAN RETURN-VALUE REPAIR","Task /490 is corrected from |α|=2 through 5 and differs from the common /119 response")
        toks=[("spaces",42408),("newline",13735),("#",3302),("…",1364),("TODO",329),("return",272),("pass",269)]
        ax=f.add_axes([.06,.43,.34,.30],facecolor=WHITE);ax.barh([x[0] for x in toks][::-1],[x[1] for x in toks][::-1],color=ORANGE);ax.set_title("Activation-weighted top tokens",loc="left",fontsize=10,color=NAVY,weight="bold");ax.spines[["top","right"]].set_visible(False);ax.grid(axis="x",color=GRID,alpha=.5)
        card(a,.44,.43,.50,.30,"BIGCODEBENCH/490 · FAIL → PASS","BASE: parses XML and writes JSON—but returns nothing\n\nSTEERED: adds `return result` to both generated functions\n\nStable at α=−2,−3,−4,−5",GREEN,WHITE,10,8.2)
        card(a,.06,.17,.27,.16,"MEANING","incomplete function / placeholder /\nreturn-boundary contexts",ORANGE,WHITE,9,7.5)
        card(a,.365,.17,.27,.16,"FAILURE MODE","required return value omitted",PURPLE,WHITE,9,7.5)
        card(a,.67,.17,.27,.16,"REVERSE TEST","base loses 1 pass only at |α|=4–5",TEAL,WHITE,9,7.5)
        a.text(.06,.105,"Promising mechanistic example; matched random/sham comparison for /490 remains PENDING.",fontsize=9.2,color=RED,weight="bold")
        p=out/"17_codellama_f13147_success.png";save(f,p,pdf);imgs.append(p)

        # 18
        f,a=setup("CODELLAMA ALSO EXPOSES A NON-SPECIFICALLY STEERABLE TASK","/119 is corrected by four targets—and by the random feature 13879")
        direct=[r for r in cl if r["side"]=="direct"]
        names=["10570","12253","14359","2310","random 13879"]
        # all known stable /119 patterns; plot pass at each alpha from official totals (0/1; 10570 gets another task at 5)
        mat=[]
        for ft in [10570,12253,14359,2310,13879]:
            q=sorted([r for r in direct if int(r["feature"])==ft],key=lambda z:abs(float(z["alpha"])))
            mat.append([int(r["fail_to_pass"]) for r in q])
        ax=f.add_axes([.07,.34,.50,.37],facecolor=WHITE);im=ax.imshow(mat,aspect="auto",cmap="Oranges",vmin=0,vmax=max(max(x) for x in mat));ax.set_yticks(range(len(names)),names);ax.set_xticks(range(5),["1","2","3","4","5"]);ax.set_xlabel("|α|");ax.set_title("Official fail→pass count",loc="left",fontsize=10,color=NAVY,weight="bold")
        for i,row in enumerate(mat):
            for j,v in enumerate(row):ax.text(j,i,str(v),ha="center",va="center",color=NAVY,weight="bold")
        card(a,.62,.49,.32,.22,"AGAINST SPECIFICITY","A random latent reproduces the same\n/119 transition at |α|=3–5",RED)
        card(a,.62,.24,.32,.20,"FOR THE METHOD","The control catches exactly the\noverclaim BBASV was designed to detect",GREEN)
        a.text(.07,.17,"Causal change? YES",fontsize=11,color=TEAL,weight="bold");a.text(.30,.17,"Feature-specific mechanism? NOT YET",fontsize=11,color=RED,weight="bold")
        p=out/"18_codellama_task_sensitivity.png";save(f,p,pdf);imgs.append(p)

        # 19
        f,a=setup("THE EVIDENCE CHAIN MAKES EACH CLAIM AUDITABLE","Association, interpretation, causality, and specificity are different evidential layers",section="SYNTHESIS")
        stages=[("TRANSITION","observed",PURPLE),("SCREEN","association",PURPLE),("TOKENS + TIMING","interpretation",ORANGE),("DIRECTION","hypothesis",ORANGE),("STEERING","causality",TEAL),("SEMANTIC CHANGE","mechanism",TEAL),("FAIL↔PASS","task outcome",GREEN),("CONTROLS","specificity",RED)]
        for i,(t,b,c) in enumerate(stages):
            x=.035+i*.12;card(a,x,.42,.105,.22,t,b,c,WHITE,8.4,7.1)
            if i<7:a.annotate("",xy=(x+.119,.53),xytext=(x+.106,.53),arrowprops=dict(arrowstyle="->",color=MUTED,lw=1.2))
        a.text(.055,.27,"DeepSeek 3048",fontsize=10,color=PURPLE,weight="bold");a.text(.20,.27,"✓  ✓  ✓  ✓  ✓  ✓  ✓",fontsize=13,color=GREEN,weight="bold");a.text(.56,.27,"controls pending",fontsize=9,color=RED,weight="bold")
        a.text(.055,.20,"CodeLlama 13147",fontsize=10,color=ORANGE,weight="bold");a.text(.20,.20,"✓  ✓  ✓  ✓  ✓  ✓  ✓",fontsize=13,color=GREEN,weight="bold");a.text(.56,.20,"partial controls",fontsize=9,color=RED,weight="bold")
        a.text(.055,.11,"The methodology does not treat every fail→pass as feature-specific; it tests where the evidence chain breaks.",fontsize=10,color=NAVY,weight="bold")
        p=out/"19_evidence_chain.png";save(f,p,pdf);imgs.append(p)

        # 20
        f,a=setup("WHAT THE PARTIAL RESULTS SUPPORT—AND WHAT THEY DO NOT","Current claims are deliberately tiered",section="SYNTHESIS")
        card(a,.055,.39,.42,.33,"SUPPORTED NOW","• screening compresses 16,384 latents to testable candidates\n• selected features causally alter generation trajectories\n• some interventions create official fail→pass transitions\n• 3048 and 13147 have semantically coherent examples\n• controls identify task-level susceptibility",GREEN,WHITE,10,8.3)
        card(a,.53,.39,.42,.33,"NOT YET SUPPORTED","• all model regressions/improvements are feature-caused\n• every successful target is feature-specific\n• CodeLlama targets outperform generic perturbations\n• generalization beyond selected cohorts\n• robustness across seeds",RED,WHITE,10,8.3)
        a.text(.055,.24,"Best current causal claim",fontsize=9,color=MUTED,weight="bold")
        a.text(.055,.17,"Selected CrossCoder features can exert directional, semantically coherent control in selected tasks.",fontsize=13,color=NAVY,weight="bold")
        a.text(.055,.105,"Specificity is a result to be demonstrated—not assumed from screening or a single correction.",fontsize=10,color=RED,weight="bold")
        p=out/"20_claims.png";save(f,p,pdf);imgs.append(p)

        # 21
        f,a=setup("A REUSABLE METHODOLOGY FOR CAUSAL MODEL DIFFING","The output is a falsifiable mechanism hypothesis—not merely a feature ranking",section="SYNTHESIS")
        steps=["1  Paired behavioral transitions","2  Same-text sparse representation","3  Conditioned multi-cell screening","4  Meaning + timing + failure mode","5  Directional causal hypothesis","6  Dose-response steering","7  Official paired evaluation","8  Random + matched sham controls","9  New seeds / held-out validation"]
        for i,s in enumerate(steps):
            col=PURPLE if i<3 else ORANGE if i<5 else TEAL if i<7 else RED
            x=.06+(i%3)*.30;y=.62-(i//3)*.19;card(a,x,y,.27,.14,s,"",col,WHITE,8.8,7)
        a.text(.06,.075,"Each stage can reject or refine the hypothesis before a broad mechanistic claim is made.",fontsize=10,color=NAVY,weight="bold")
        p=out/"21_method.png";save(f,p,pdf);imgs.append(p)

        # 22
        f,a=setup("FROM MODEL DIFFERENCES TO TESTABLE CAUSAL MECHANISMS","The partial BBASV results already show both the promise and the discipline of the framework",section="CONCLUSION")
        a.text(.07,.66,"Model differences are not only observable in behavior.",fontsize=23,color=NAVY,weight="bold")
        a.text(.07,.55,"They can be localized in a shared latent representation,",fontsize=20,color=PURPLE,weight="bold")
        a.text(.07,.46,"interpreted in context, and—in selected cases—",fontsize=20,color=ORANGE,weight="bold")
        a.text(.07,.37,"causally manipulated.",fontsize=24,color=TEAL,weight="bold")
        card(a,.07,.16,.82,.12,"THE CONTRIBUTION","A testable path from behavioral contrast → sparse feature → semantic hypothesis → controlled intervention → specificity test",GREEN,WHITE,10,9)
        a.text(.07,.09,"PENDING: complete DeepSeek controls, CodeLlama sham/random controls, and multi-seed validation.",fontsize=9.5,color=RED,weight="bold")
        p=out/"22_conclusion.png";save(f,p,pdf);imgs.append(p)
    return imgs

def main():
    ap=argparse.ArgumentParser();ap.add_argument("--repo",type=Path,default=Path("."));ap.add_argument("--source",type=Path,required=True);ap.add_argument("--dest",type=Path,required=True);ap.add_argument("--output-dir",type=Path,required=True)
    x=ap.parse_args();imgs=build(x.repo,x.output_dir);shutil.copy2(x.source,x.dest);prs=Presentation(x.dest)
    # Replace the clean draft preliminary BBASV slide with the expanded block.
    if len(prs.slides)==11:
        sid=prs.slides._sldIdLst[-1];prs.part.drop_rel(sid.rId);del prs.slides._sldIdLst[-1]
    blank=prs.slide_layouts[6]
    for p in imgs:
        s=prs.slides.add_slide(blank);s.shapes.add_picture(str(p),0,0,width=prs.slide_width,height=prs.slide_height)
    prs.save(x.dest)
    print(x.dest)

if __name__=="__main__":main()
