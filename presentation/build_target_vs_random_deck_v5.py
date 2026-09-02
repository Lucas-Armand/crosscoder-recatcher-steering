#!/usr/bin/env python3
from pathlib import Path
from zipfile import ZipFile
import io, re
import numpy as np
import matplotlib.pyplot as plt
from matplotlib.patches import FancyBboxPatch, Rectangle, Circle, FancyArrowPatch
from pptx import Presentation
from PIL import Image

R=Path('/home/lucas/crosscoder-recatcher-steering')
SRC=R/'MODEL DIFFS_SAME_TEXT_PIPELINE_V4_2026-09-01.pptx'
DEST=R/'MODEL DIFFS_TARGET_VS_RANDOM_V5_2026-09-02.pptx'
PDF=R/'MODEL DIFFS_TARGET_VS_RANDOM_V5_2026-09-02.pdf'
OUT=R/'presentation/generated/target_vs_random_v5_20260902';OUT.mkdir(parents=True,exist_ok=True)
BG='#F7F4EE';NAVY='#10233F';PURPLE='#7257C7';ORANGE='#EE8A2D';TEAL='#159C9C';MUTED='#687386';WHITE='#FFFFFF';GRID='#D8D9D8';RED='#C93948';GREEN='#338A63';GOLD='#D4A72C';LIGHT='#EEEAE2'
A=np.array([1,2,3,4,5])
DS_T={1078:[2,5,5,7,6],14175:[1,2,6,6,5],15235:[2,3,6,5,4],2468:[3,4,7,7,7],2621:[3,4,6,5,8]}
DS_TC={1078:[14,23,26,30,39],14175:[13,17,26,32,33],15235:[14,15,23,25,26],2468:[13,25,34,37,44],2621:[12,22,28,31,40]}
DS_R={605:[0,1,4,3,3],6023:[1,1,1,1,4],15173:[1,2,3,3,4]}
DS_RC={605:[12,20,24,26,31],6023:[12,20,22,29,30],15173:[12,19,22,28,30]}
CL_T={4309:[0,0,1,2,2],5642:[0,0,1,1,1],7692:[0,1,1,1,1],10818:[0,0,1,1,1],11596:[0,1,1,1,1]}
CL_TC={4309:[8,14,21,23,27],5642:[9,13,18,25,28],7692:[5,13,17,24,25],10818:[5,13,18,22,26],11596:[9,16,22,25,29]}
CL_R={2947:[0,0,0,0,0],3698:[0,0,0,0,0],14260:[0,1,1,1,1]}
CL_RC={2947:[7,15,17,21,26],3698:[9,11,14,17,20],14260:[6,13,18,21,21]}

def setup(title,sub='',section='PROSPECTIVE TARGET–RANDOM COMPARISON'):
 f=plt.figure(figsize=(13.333,7.5),facecolor=BG);a=f.add_axes([0,0,1,1]);a.axis('off')
 a.text(.045,.955,section,fontsize=8.4,color=TEAL,weight='bold',va='top');a.text(.045,.885,title,fontsize=22,color=NAVY,weight='bold',va='top')
 if sub:a.text(.045,.815,sub,fontsize=10,color=MUTED,va='top')
 a.text(.95,.035,'CrossCoder model diffing · validated direct readout · 2026-09-02',ha='right',fontsize=7.3,color=MUTED)
 return f,a
def card(a,x,y,w,h,title,body,accent,ts=10,bs=8.2):
 a.add_patch(FancyBboxPatch((x,y),w,h,boxstyle='round,pad=.007,rounding_size=.011',facecolor=WHITE,edgecolor='#E1DED7'))
 a.add_patch(Rectangle((x,y+h-.009),w,.009,color=accent,lw=0));a.text(x+.015,y+h-.037,title,fontsize=ts,color=NAVY,weight='bold',va='top');a.text(x+.015,y+h-.085,body,fontsize=bs,color=MUTED,va='top',linespacing=1.28)
def arrow(a,x1,y1,x2,y2,c=GRID,w=1.4):a.add_patch(FancyArrowPatch((x1,y1),(x2,y2),arrowstyle='-|>',mutation_scale=12,color=c,lw=w))
def save(f,n):p=OUT/f'{n:02d}.png';f.savefig(p,dpi=180,facecolor=BG);plt.close(f);return p
def vals(d):return np.array(list(d.values()),dtype=float)
def dose_panel(fig,rect,title,target,random,col,denom,ymax=None):
 ax=fig.add_axes(rect);t=vals(target);r=vals(random)
 for row in t:ax.plot(A,row,color=col,alpha=.20,lw=1.2)
 for row in r:ax.plot(A,row,color=GOLD,alpha=.28,lw=1.2,ls='--')
 ax.plot(A,np.median(t,axis=0),'-o',color=col,lw=3,label='selected features · median')
 ax.plot(A,np.median(r,axis=0),'--o',color=GOLD,lw=2.5,label='random features · median')
 ax.set_title(title,loc='left',fontsize=10,color=NAVY,weight='bold');ax.set_xlabel('|α|');ax.set_ylabel(f'tasks / {denom}');ax.set_xticks(A);ax.set_ylim(bottom=-.15,top=ymax);ax.grid(color=GRID,alpha=.6);ax.spines[['top','right']].set_visible(False);ax.legend(frameon=False,fontsize=7,loc='upper left')
def changed_panel(fig,rect,title,target,random,col,denom):
 ax=fig.add_axes(rect);t=vals(target);r=vals(random)
 ax.plot(A,t.mean(0),'-o',color=col,lw=3,label='selected · mean')
 ax.fill_between(A,t.min(0),t.max(0),color=col,alpha=.13)
 ax.plot(A,r.mean(0),'--o',color=GOLD,lw=2.5,label='random · mean')
 ax.fill_between(A,r.min(0),r.max(0),color=GOLD,alpha=.13)
 ax.set_title(title,loc='left',fontsize=10,color=NAVY,weight='bold');ax.set_xlabel('|α|');ax.set_ylabel(f'evaluated codes changed / {denom}');ax.set_xticks(A);ax.grid(color=GRID,alpha=.6);ax.spines[['top','right']].set_visible(False);ax.legend(frameon=False,fontsize=7)

imgs={}
# 9: update the funnel rationale for the final direct comparison.
f,a=setup('SCREEN FIRST; TEST CAUSALLY AT A COMMON OPERATING POINT','The statistical shortlist reduces the intervention search by ~99.8%',section='ERROR-FOCUSED FEATURE SELECTION')
for x,big,small,col in [(.10,'16,384','latents / CrossCoder',NAVY),(.40,'≈30–40','candidates / model',TEAL),(.70,'≈0.2%','of latent space',RED)]:a.text(x,.63,big,fontsize=36,color=col,weight='bold',ha='center');a.text(x,.56,small,fontsize=9,color=MUTED,ha='center')
arrow(a,.18,.61,.30,.61);arrow(a,.48,.61,.60,.61)
card(a,.07,.20,.39,.25,'WHY |α| = 3?','Earlier dose-response work identified a practical exploratory point: visible causal changes below the most disruptive high-dose regime.',GOLD)
card(a,.54,.20,.39,.25,'WHY STANDARDIZE?','One common magnitude makes dozens of candidates comparable before selecting five features for complete dose curves.',TEAL)
a.text(.07,.11,'|α|=3 is a common operating point—not an independently validated optimum.',fontsize=9.5,color=NAVY,weight='bold');imgs[9]=save(f,9)

# 16: retain task concentration and connect it to random-feature comparison.
f,a=setup('DEEPSEEK: 116 TRANSITIONS COLLAPSE TO 16 SUSCEPTIBLE TASKS','Feature–task outcomes reveal both latent leverage and task-level susceptibility',section='CANONICAL ERROR-FOCUSED α=3 SWEEP')
tasks=[('/316',17),('/435',15),('/166',13),('/496',11),('/937',10),('/349',7),('/459',7),('/7',7),('/634',6),('/695',6),('/940',5),('/188',5),('/666',3),('/1',2),('/576',1),('/823',1)]
ax=f.add_axes([.075,.17,.56,.59]);y=list(range(len(tasks)))[::-1];ax.barh(y,[n for _,n in tasks],color=PURPLE);ax.set_yticks(y,[t for t,_ in tasks],fontsize=7.5);ax.set_xlabel('screened features that corrected the task');ax.spines[['top','right']].set_visible(False);ax.grid(axis='x',color=GRID,alpha=.45)
card(a,.68,.46,.27,.28,'CONCENTRATION','/316   17 features\n/435   15 features\n/166   13 features\n/496   11 features\n/937   10 features',PURPLE,10,8.8)
card(a,.68,.20,.27,.19,'NEXT COMPARISON','The complete dose experiment tests whether the selected pool reaches corrections beyond prospectively sampled random CrossCoder features.',TEAL,9.5,7.8)
a.text(.075,.10,'16 unique corrections across 80 contamination-focused improvements.',fontsize=9,color=NAVY,weight='bold');imgs[16]=save(f,16)

# 20 alpha=3 context
f,a=setup('THE α=3 SWEEP REVEALS BOTH MODEL AND TASK SENSITIVITY','Complete error-focused sweep · official BigCodeBench 0.1.5 evaluation',section='CAUSAL CANDIDATE SWEEP')
ds={0:1,1:5,2:7,3:8,4:3,5:6,6:4,7:1};cl={0:26,1:6,2:0}
for pos,title,hist,col in [([.065,.30,.38,.43],'DeepSeek · 35 screened features',ds,PURPLE),([.535,.30,.38,.43],'CodeLlama · 32 screened features',cl,ORANGE)]:
 ax=f.add_axes(pos);xs=list(hist);ax.bar(xs,[hist[x] for x in xs],color=col,width=.72);ax.set_title(title,loc='left',color=NAVY,weight='bold',fontsize=10);ax.set_xlabel('official fail→pass tasks per feature',fontsize=8);ax.set_ylabel('features',fontsize=8);ax.set_xticks(xs);ax.grid(axis='y',color=GRID,alpha=.65);ax.spines[['top','right']].set_visible(False);ax.tick_params(labelsize=8)
a.text(.065,.205,'116 feature–task transitions  →  16/80 unique tasks',fontsize=9,color=PURPLE,weight='bold');a.text(.535,.205,'6 feature–task transitions  →  4/50 unique tasks',fontsize=9,color=ORANGE,weight='bold')
a.add_patch(FancyBboxPatch((.06,.075),.86,.065,boxstyle='round,pad=.007',facecolor=WHITE,edgecolor=GRID));a.text(.49,.108,'The sweep defines five target features per model for prospective dose-response comparison with random CrossCoder features.',ha='center',va='center',fontsize=9,color=NAVY,weight='bold');imgs[20]=save(f,20)

# 21 protocol
f,a=setup('FROM SCREENING TO A PROSPECTIVE TARGET–RANDOM TEST','Every arm shares the same tasks, canonical seed rule, generation settings, and evaluator',section='CONTROLLED CAUSAL DESIGN')
steps=[('1','α=0 GATE','byte-exact canonical reproduction',GREEN),('2','TARGETS','five screen-selected features',PURPLE),('3','RANDOM','three prospectively sampled CrossCoder latents',GOLD),('4','DOSE','|α| = 1 · 2 · 3 · 4 · 5',TEAL),('5','ENDPOINTS','code changed + official fail→pass',ORANGE)]
for i,(n,t,b,c) in enumerate(steps):card(a,.045+i*.19,.42,.17,.29,n+'  '+t,b,c,9.2,7.8)
card(a,.055,.18,.835,.14,'INTERVENTION','Layer 16 · decoder direction on the intervened model side · last residual position · added at every autoregressive generation step · temperature 0.2 · top-p 0.95 · max 512',NAVY,9.5,8.4)
a.text(.055,.11,'Reproduction gate passed: DeepSeek 80/80 and CodeLlama 50/50 raw completions matched byte for byte; both baselines score 0 passes.',fontsize=9.3,color=GREEN,weight='bold');imgs[21]=save(f,21)

# 22 central outcome comparison
f,a=setup('SCREEN-SELECTED FEATURES OUTPERFORM RANDOM LATENTS','Official fail→pass outcomes across the complete dose range',section='PRIMARY COMPARATIVE RESULT')
dose_panel(f,[.065,.29,.40,.45],'DeepSeek · contamination cohort',DS_T,DS_R,PURPLE,80,9)
dose_panel(f,[.535,.29,.40,.45],'CodeLlama · logic/runtime cohort',CL_T,CL_R,ORANGE,50,2.5)
a.text(.065,.195,'Selected: 16 unique tasks',fontsize=10,color=PURPLE,weight='bold');a.text(.065,.155,'Random: 10 · all overlap selected successes',fontsize=8.5,color=NAVY)
a.text(.535,.195,'Selected: 5 unique tasks',fontsize=10,color=ORANGE,weight='bold');a.text(.535,.155,'Random: 1 · overlaps a selected success',fontsize=8.5,color=NAVY)
a.text(.065,.09,'The selected pools reach six additional DeepSeek tasks and four additional CodeLlama tasks beyond the random-feature pools.',fontsize=9.3,color=GREEN,weight='bold');imgs[22]=save(f,22)

# 23 DeepSeek detailed
f,a=setup('DEEPSEEK: THE SELECTED POOL SHOWS STRONGER AND BROADER CONTROL','Five target features versus three prospectively sampled random CrossCoder features',section='DEEPSEEK · TARGET VS RANDOM')
dose_panel(f,[.065,.30,.40,.43],'Official corrections',DS_T,DS_R,PURPLE,80,9)
changed_panel(f,[.535,.30,.40,.43],'Evaluated-code changes',DS_TC,DS_RC,PURPLE,80)
card(a,.065,.11,.27,.115,'OUTCOME','16 vs 10 unique tasks',PURPLE,9.5,8.3);card(a,.365,.11,.27,.115,'SEPARATION','largest around |α|=2–3',TEAL,9.5,8.3);card(a,.665,.11,.27,.115,'TRAJECTORY','selected mean > random mean at every dose',GREEN,9.5,7.8);imgs[23]=save(f,23)

# 24 feature 2468
f,a=setup('FEATURE 2468: A STABLE DEEPSEEK MECHANISTIC CANDIDATE','Screening rank, timing, and causal response converge around post-solution contamination',section='DEEPSEEK · FEATURE DEEP DIVE')
ax=f.add_axes([.075,.30,.44,.43]);ax.plot(A,DS_T[2468],'-o',lw=3,color=PURPLE,label='official corrections');ax.plot(A,DS_TC[2468],'-s',lw=2.4,color='#A995E4',label='evaluated code changed');ax.set_xticks(A);ax.set_xlabel('|α| · behavior-aligned sign');ax.set_ylabel('tasks / 80');ax.grid(color=GRID,alpha=.6);ax.spines[['top','right']].set_visible(False);ax.legend(frameon=False)
card(a,.57,.49,.36,.24,'SCREENING','paired local · rank #1\nbase-enriched · max\nE/V = −7.29\nsupport: 62/80',PURPLE,10,8.4)
card(a,.57,.25,.36,.19,'CAUSAL READOUT','3 · 4 · 7 · 7 · 7 corrections\n13 → 44 evaluated codes changed',GREEN,10,8.4)
a.text(.57,.145,'Meaning ↔ failure mode ↔ activation timing',fontsize=9.8,color=NAVY,weight='bold');imgs[24]=save(f,24)

# 25 clean case
f,a=setup('DEEPSEEK /7: FEATURE 2468 REMOVES POST-SOLUTION CONTAMINATION','Behavior-aligned steering · |α|=4 · official fail→pass',section='DEEPSEEK · CAUSAL EXAMPLE')
card(a,.055,.30,.42,.43,'BASE · FAIL','return top_product\n\n# test_task.py\nimport unittest\nfrom task import task_func\n\nclass TestTask(unittest.TestCase):\n    ...',RED,11,8.7)
card(a,.52,.30,.42,.43,'STEERED · PASS','return top_product\n\n# task_func("path/to/sales.csv")\n# task_func("sales.csv")\n\n[no executable test harness]',GREEN,11,8.7)
card(a,.055,.13,.885,.105,'MECHANISTIC READING','The useful function is preserved while executable post-solution tests and imports are removed near the targeted late-generation region.',PURPLE,9.5,8.4);imgs[25]=save(f,25)

# 26 non-pass
f,a=setup('THE SAME FEATURE MAKES COHERENT CHANGES WITHOUT ALWAYS REPAIRING THE TASK','Trajectory control and functional correctness are complementary endpoints',section='DEEPSEEK · SEMANTIC BOUNDARY')
card(a,.055,.39,.27,.34,'/97 · STILL FAILS','Expected floating-point constants inside generated tests change.\n\nThe trajectory changes late, but the irrelevant test harness remains.',PURPLE,10,8.2)
card(a,.365,.39,.27,.34,'/823 · STILL FAILS','Generated test precision and comments change near the end.\n\nPost-solution material changes without repairing the evaluated function.',PURPLE,10,8.2)
card(a,.675,.39,.27,.34,'INTERPRETATION','The semantic locus remains coherent even when the intervention is insufficient for an official correction.',ORANGE,10,8.2)
a.text(.055,.17,'Feature 2468 controls late generated content; benchmark success additionally depends on preserving or repairing the functional body.',fontsize=9.2,color=NAVY,weight='bold');imgs[26]=save(f,26)

# 27 CodeLlama aggregate
f,a=setup('CODELLAMA: SELECTED FEATURES PRODUCE SMALLER BUT BROADER OUTCOME EFFECTS','All five targets repair at least one task; the random pool reaches one recurrent task',section='CODELLAMA · TARGET VS RANDOM')
dose_panel(f,[.065,.30,.40,.43],'Official corrections',CL_T,CL_R,ORANGE,50,2.5)
changed_panel(f,[.535,.30,.40,.43],'Evaluated-code changes',CL_TC,CL_RC,ORANGE,50)
card(a,.065,.11,.27,.115,'OUTCOME','5 vs 1 unique tasks',ORANGE,9.5,8.3);card(a,.365,.11,.27,.115,'CONSISTENCY','5/5 targets repair by |α|=3',TEAL,9.5,8.3);card(a,.665,.11,.27,.115,'TRAJECTORY','selected mean > random mean from |α|=2',GREEN,9.5,7.8);imgs[27]=save(f,27)

# 28 best CL
f,a=setup('FEATURE 4309: CODELLAMA’S STRONGEST OUTCOME CANDIDATE','Code changes grow with dose; official repairs remain selective',section='CODELLAMA · FEATURE DEEP DIVE')
ax=f.add_axes([.075,.30,.44,.43]);ax.plot(A,CL_T[4309],'-o',lw=3,color=ORANGE,label='official corrections');ax.plot(A,CL_TC[4309],'-s',lw=2.4,color='#F4BE84',label='evaluated code changed');ax.set_xticks(A);ax.set_xlabel('|α| · behavior-aligned sign');ax.set_ylabel('tasks / 50');ax.grid(color=GRID,alpha=.6);ax.spines[['top','right']].set_visible(False);ax.legend(frameon=False)
card(a,.57,.49,.36,.24,'CAUSAL READOUT','0 · 0 · 1 · 2 · 2 corrections\n8 → 27 evaluated codes changed',ORANGE,10,8.7)
card(a,.57,.25,.36,.19,'MODEL CONTRAST','Outcome control is smaller than DeepSeek, while trajectory control remains clearly dose-responsive.',TEAL,10,8.2)
a.text(.57,.145,'Interpretation: a more distributed logic/runtime regression.',fontsize=9,color=NAVY,weight='bold');imgs[28]=save(f,28)

# 29 CL example
f,a=setup('CODELLAMA /490: STEERING REPAIRS PARSING, WRITING, AND RETURN','Feature 4309 · |α|=5 · official fail→pass',section='CODELLAMA · CAUSAL EXAMPLE')
card(a,.055,.30,.42,.43,'MERGED · FAIL','return json.loads(s)\n\n[does not parse XML correctly]\n[does not write the requested JSON]\n[one generated function returns nothing]',RED,11,8.7)
card(a,.52,.30,.42,.43,'STEERED · PASS','result = xmltodict.parse(s)\nwith open(file_path, "w") as f:\n    json.dump(result, f)\nreturn result',GREEN,11,8.7)
card(a,.055,.13,.885,.105,'MECHANISTIC READING','The intervention produces a coherent multi-part repair in a selected task. Across the cohort, five target features reach five distinct successful tasks.',ORANGE,9.5,8.4);imgs[29]=save(f,29)

# 30 evidence chain
f,a=setup('THE EVIDENCE CHAIN CONNECTS MODEL DIFFERENCES TO CAUSAL CONTROL','Each stage adds a distinct, auditable form of evidence',section='SYNTHESIS')
labels=['BEHAVIORAL\nTRANSITION','FOCUSED\nSCREEN','TOKENS +\nTIMING','DIRECTIONAL\nHYPOTHESIS','α=0\nGATE','DOSE\nCURVE','CODE\nCHANGE','FAIL→PASS','RANDOM-FEATURE\nCOMPARISON'];cols=[MUTED,PURPLE,PURPLE,TEAL,GREEN,TEAL,ORANGE,GREEN,GOLD]
for i,(lab,col) in enumerate(zip(labels,cols)):
 x=.025+i*.106;a.add_patch(FancyBboxPatch((x,.46),.093,.17,boxstyle='round,pad=.005',facecolor=WHITE,edgecolor=col,lw=2));a.text(x+.0465,.545,lab,ha='center',va='center',fontsize=7.1,color=NAVY,weight='bold')
 if i<8:arrow(a,x+.095,.545,x+.108,.545,MUTED,1.1)
card(a,.055,.20,.40,.16,'ASSOCIATIVE + INTERPRETATIVE','transition · screening · tokens · timing',PURPLE,9.5,8.3);card(a,.50,.20,.40,.16,'CAUSAL + COMPARATIVE','reproduction · intervention · code · outcome · random features',GREEN,9.5,8.1);imgs[30]=save(f,30)

# 31 method
f,a=setup('A REUSABLE METHODOLOGY FOR CAUSAL MODEL DIFFING','The output is a testable mechanism hypothesis—not merely a feature ranking',section='METHODOLOGICAL CONTRIBUTION')
items=['Paired behavioral transitions','Error taxonomy + focal cohort','Same-text sparse representation','Conditioned E/V screening','Meaning + timing + failure mode','Behavior-aligned direction','Canonical α=0 gate','Dose-response steering','Random CrossCoder features','Official code + outcome evaluation']
for i,t in enumerate(items):
 c=i%5;r=i//5;x=.035+c*.19;y=.55-r*.25;card(a,x,y,.17,.17,str(i+1),t,[PURPLE,TEAL,ORANGE,GREEN,GOLD][c],10,7.7)
a.text(.055,.13,'The comparison asks whether statistically selected latents provide stronger causal leverage than features sampled without access to screening outcomes.',fontsize=9.3,color=NAVY,weight='bold');imgs[31]=save(f,31)

# 32 conclusion
f,a=setup('MODEL DIFFERENCES CAN BECOME TESTABLE CAUSAL MECHANISMS','A shared sparse representation connects behavioral contrast to controlled intervention',section='CONCLUSION')
a.text(.07,.65,'Behavioral differences',fontsize=18,color=NAVY,weight='bold');a.text(.36,.65,'→',fontsize=24,color=TEAL,weight='bold');a.text(.43,.65,'screen-selected features',fontsize=18,color=PURPLE,weight='bold');a.text(.77,.65,'→',fontsize=24,color=TEAL,weight='bold')
a.text(.07,.50,'interpretable timing + semantics',fontsize=18,color=ORANGE,weight='bold');a.text(.49,.50,'→',fontsize=24,color=TEAL,weight='bold');a.text(.57,.50,'causal steering',fontsize=18,color=GREEN,weight='bold')
card(a,.07,.23,.84,.17,'SUPPORTED RESULT','Screen-selected CrossCoder features produce stronger and broader causal effects than prospectively sampled random CrossCoder features, including official fail→pass transitions.',NAVY,11,10)
a.text(.07,.13,'DeepSeek: stronger, localized outcome control · CodeLlama: smaller outcome effects with dose-dependent trajectory control.',fontsize=9.5,color=MUTED);imgs[32]=save(f,32)

prs=Presentation(SRC)
for n,p in imgs.items():
 s=prs.slides[n-1]
 for sh in list(s.shapes):s.shapes._spTree.remove(sh._element)
 s.shapes.add_picture(str(p),0,0,width=prs.slide_width,height=prs.slide_height)
prs.save(DEST)

with ZipFile(DEST) as z:
 pages=[]
 for n in range(1,len(prs.slides)+1):
  if n in imgs:pages.append(Image.open(imgs[n]).convert('RGB'));continue
  rel=z.read(f'ppt/slides/_rels/slide{n}.xml.rels').decode();c=[]
  for target in re.findall(r'Target="([^"]+)"',rel):
   if target.startswith('../media/'):
    b=z.read('ppt/'+target[3:]);c.append((len(b),b))
  if not c:raise RuntimeError(f'No raster found for slide {n}')
  pages.append(Image.open(io.BytesIO(max(c)[1])).convert('RGB'))
pages[0].save(PDF,save_all=True,append_images=pages[1:],resolution=180)
print(DEST);print(PDF);print('slides',len(prs.slides),'pdf_pages',len(pages));print('generated',sorted(imgs))
