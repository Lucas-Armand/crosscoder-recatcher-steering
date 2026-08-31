#!/usr/bin/env python3
from pathlib import Path
from zipfile import ZipFile
import io, re, shutil
import matplotlib.pyplot as plt
from matplotlib.backends.backend_pdf import PdfPages
from matplotlib.patches import FancyBboxPatch, Rectangle
from pptx import Presentation
from PIL import Image

R=Path('/home/lucas/crosscoder-recatcher-steering')
SRC=R/'MODEL DIFFS_CAUSAL_STEERING_RESULTS_ALPHA3_COMPLETE_2026-08-31.pptx'
DEST=R/'MODEL DIFFS_CAUSAL_STEERING_RESULTS_ALPHA3_COMPLETE_V2_2026-08-31.pptx'
PDF=R/'MODEL DIFFS_CAUSAL_STEERING_RESULTS_ALPHA3_COMPLETE_V2_2026-08-31.pdf'
OUT=R/'presentation/generated/alpha3_complete_v2_slides_20260831'; OUT.mkdir(parents=True,exist_ok=True)
BG='#F7F4EE'; NAVY='#10233F'; PURPLE='#7257C7'; ORANGE='#EE8A2D'; TEAL='#159C9C'; MUTED='#687386'; WHITE='#FFFFFF'; GRID='#D8D9D8'; RED='#C93948'; GREEN='#338A63'

def setup(title,sub):
 f=plt.figure(figsize=(13.333,7.5),facecolor=BG);a=f.add_axes([0,0,1,1]);a.axis('off')
 a.text(.045,.955,'CANONICAL ERROR-FOCUSED α=3 SWEEP · COMPLETE READOUT',fontsize=8.5,color=TEAL,weight='bold',va='top')
 a.text(.045,.885,title,fontsize=22,color=NAVY,weight='bold',va='top');a.text(.045,.815,sub,fontsize=10,color=MUTED,va='top')
 a.text(.95,.035,'CrossCoder model diffing · canonical seed · 2026-08-31',ha='right',fontsize=7.5,color=MUTED)
 return f,a
def card(a,x,y,w,h,title,body,accent,ts=10,bs=8):
 a.add_patch(FancyBboxPatch((x,y),w,h,boxstyle='round,pad=.007,rounding_size=.011',facecolor=WHITE,edgecolor='#E1DED7'))
 a.add_patch(Rectangle((x,y+h-.01),w,.01,color=accent,lw=0));a.text(x+.016,y+h-.04,title,fontsize=ts,color=NAVY,weight='bold',va='top');a.text(x+.016,y+h-.085,body,fontsize=bs,color=MUTED,va='top',linespacing=1.3)
def save(f,n):
 p=OUT/f'{n}.png';f.savefig(p,dpi=175,facecolor=BG);plt.close(f);return p

# Slide 18: final top-five selection after the complete alpha=3 sweep.
f,a=setup('THE COMPLETE α=3 SWEEP DEFINES THE NEXT TARGETS','Ordering: official passes first; |E/V| breaks ties')
ds=[('2468','7/80','paired local · #1','7.29'),('2621','6/80','paired global · #8','10.53'),('1078','6/80','association global · #3','5.72'),('15235','6/80','association global · #8','4.61'),('14175','6/80','association local · #1','4.09')]
cl=[('7692','1/50','association global · #1','5.90'),('10818','1/50','association global · #2','5.90'),('5642','1/50','association global · #3','5.74'),('11596','1/50','association global · #5','5.61'),('4309','1/50','association local · #9','4.86')]
for x,title,data,col in [(.055,'DEEPSEEK',ds,PURPLE),(.525,'CODELLAMA',cl,ORANGE)]:
 a.text(x,.75,title,fontsize=12,color=col,weight='bold')
 cols=[('RANK',.055),('FEATURE',.13),('PASS',.235),('BEST SCREEN',.31),('|E/V|',.445)]
 for label,dx in cols:a.text(x+dx,.69,label,fontsize=7,color=MUTED,weight='bold')
 for i,(fid,pas,scr,ev) in enumerate(data):
  y=.625-i*.085;a.add_patch(Rectangle((x,y-.025),.41,.067,facecolor=WHITE,edgecolor=GRID))
  vals=[str(i+1),fid,pas,scr,ev]
  for val,(_,dx) in zip(vals,cols):a.text(x+dx,y,val,fontsize=7.8,color=NAVY,weight='bold' if i==0 else 'normal')
a.text(.055,.11,'Complete ranking: causal outcome first; screening strength only resolves equal pass counts. Full dose curves and controls test robustness and specificity.',fontsize=9.2,color=RED,weight='bold')
p18=save(f,18)

# Slide 19: complete sweep overview.
f,a=setup('THE α=3 SWEEP REVEALS BOTH MODEL AND TASK SENSITIVITY','Error-focused candidates · continuous last-token steering · official BigCodeBench 0.1.5')
card(a,.055,.37,.40,.36,'DEEPSEEK · CONTAMINATION IMPROVEMENTS','35/35 features evaluated\n116 feature-task transitions\n16/80 unique tasks corrected\n\nSeveral features correct multiple tasks.',PURPLE,12,9.2)
card(a,.49,.37,.40,.36,'CODELLAMA · LOGIC/RUNTIME REGRESSIONS','32/32 features evaluated\n6 feature-task transitions\n4/50 unique tasks corrected\n\nSingle-feature corrections remain sparse.',ORANGE,12,9.2)
card(a,.055,.16,.835,.13,'READING','DeepSeek contains repeated multi-task causal candidates; CodeLlama remains sparse and task-specific. Concentration on susceptible tasks reinforces—rather than replaces—the need for random and sham controls.',TEAL,10,8.8)
a.text(.055,.105,'COMPLETE · Both sweeps used the canonical seed rule and passed their α=0 reproduction gates.',fontsize=9.5,color=GREEN,weight='bold')
p19=save(f,19)

prs=Presentation(SRC)
for n,p in [(18,p18),(19,p19)]:
 s=prs.slides[n-1]
 for sh in list(s.shapes):s.shapes._spTree.remove(sh._element)
 s.shapes.add_picture(str(p),0,0,width=prs.slide_width,height=prs.slide_height)
prs.save(DEST)

# Preserve all other PDF pages using the largest slide image embedded in the source deck.
with ZipFile(SRC) as z:
 pages=[]
 for n in range(1,len(prs.slides)+1):
  if n in (18,19): pages.append(Image.open(OUT/f'{n}.png').convert('RGB'));continue
  rel=z.read(f'ppt/slides/_rels/slide{n}.xml.rels').decode()
  candidates=[]
  for target in re.findall(r'Target="([^\"]+)"',rel):
   if target.startswith('../media/'):
    member='ppt/'+target[3:];data=z.read(member);candidates.append((len(data),data))
  pages.append(Image.open(io.BytesIO(max(candidates)[1])).convert('RGB'))
pages[0].save(PDF,save_all=True,append_images=pages[1:],resolution=175)
print(DEST);print(PDF);print('slides',len(prs.slides),'pdf_pages',len(pages))
