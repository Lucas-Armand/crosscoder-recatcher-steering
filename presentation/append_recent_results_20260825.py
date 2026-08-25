#!/usr/bin/env python3
"""Append the 2026-08-25 screening/steering audit to the user-edited deck."""
from pathlib import Path
import argparse, csv, shutil, subprocess
import numpy as np
import matplotlib.pyplot as plt
from matplotlib.backends.backend_pdf import PdfPages
from matplotlib.patches import FancyBboxPatch, Rectangle
from pptx import Presentation

BG="#F6F3EC"; NAVY="#10233F"; TEAL="#169C9C"; CORAL="#E5655B"; GOLD="#E3A62F"; MUTED="#687386"; WHITE="#FFFFFF"; GRID="#D7D9D8"
ALPHAS=np.array([.5,1,2,3,4,5,6])

def setup(title, subtitle=""):
    fig=plt.figure(figsize=(13.333,7.5),facecolor=BG); ax=fig.add_axes([0,0,1,1]);ax.axis("off")
    ax.text(.045,.955,"RECENT EXTENSION · SCREENING → SEMANTICS → CAUSAL STRESS TEST",fontsize=8.5,color=TEAL,weight="bold",va="top")
    ax.text(.045,.885,title,fontsize=25,color=NAVY,weight="bold",va="top")
    if subtitle: ax.text(.045,.815,subtitle,fontsize=11,color=MUTED,va="top")
    ax.text(.95,.04,"CrossCoder model diffing · exploratory audit · 2026-08-25",ha="right",fontsize=7.5,color=MUTED)
    return fig,ax

def card(ax,x,y,w,h,title,body="",accent=TEAL):
    ax.add_patch(FancyBboxPatch((x,y),w,h,boxstyle="round,pad=0.008,rounding_size=.012",facecolor=WHITE,edgecolor="#E2E0D9",linewidth=1))
    ax.add_patch(Rectangle((x,y+h-.012),w,.012,color=accent,linewidth=0))
    ax.text(x+.025,y+h-.06,title,fontsize=12,color=NAVY,weight="bold",va="top")
    if body: ax.text(x+.025,y+h-.115,body,fontsize=9.2,color=MUTED,va="top",linespacing=1.45)

def save(fig,path,pdf):
    fig.savefig(path,dpi=150,facecolor=BG);pdf.savefig(fig,facecolor=BG);plt.close(fig)

def line_chart(fig,rect,title,target,shams,controls,ylabel="Official passes"):
    a=fig.add_axes(rect,facecolor=WHITE)
    sm=np.mean(shams,axis=0); slo=np.min(shams,axis=0);shi=np.max(shams,axis=0)
    cm=np.mean(controls,axis=0);clo=np.min(controls,axis=0);chi=np.max(controls,axis=0)
    a.fill_between(ALPHAS,slo,shi,color=MUTED,alpha=.16);a.plot(ALPHAS,sm,"--o",color=MUTED,label="3 orthogonal shams")
    a.fill_between(ALPHAS,clo,chi,color=GOLD,alpha=.14);a.plot(ALPHAS,cm,"--o",color=GOLD,label="3 latent controls")
    a.plot(ALPHAS,target,"-o",lw=3,color=TEAL,label="selected feature")
    a.set_title(title,loc="left",fontsize=11,color=NAVY,weight="bold");a.set_xlabel("|α|",fontsize=9);a.set_ylabel(ylabel,fontsize=9)
    a.grid(axis="y",color=GRID,alpha=.7);a.spines[["top","right"]].set_visible(False);a.legend(frameon=False,fontsize=8,loc="upper left")
    return a

def build(outdir):
    outdir.mkdir(parents=True,exist_ok=True); images=[]
    pdf_path=outdir/"recent_results_appendix.pdf"
    with PdfPages(pdf_path) as pdf:
        fig,ax=setup("EIGHT COMPLEMENTARY SCREENS", "Association, paired model difference, and error-local timing answer different questions")
        labels=[("REGRESSION","association · global","failures vs retained successes"),("REGRESSION","association · local","same contrast near first divergence"),("REGRESSION","paired · global","specialized text − base text"),("REGRESSION","paired · local","paired difference near divergence"),("IMPROVEMENT","association · global","improvements vs remaining base failures"),("IMPROVEMENT","association · local","same contrast near first divergence"),("IMPROVEMENT","paired · global","specialized text − base text"),("IMPROVEMENT","paired · local","paired difference near divergence")]
        for i,(head,mid,body) in enumerate(labels):
            x=.05+(i%4)*.235;y=.49-(i//4)*.28
            card(ax,x,y,.205,.21,head,mid+"\n"+body,TEAL if head=="IMPROVEMENT" else CORAL)
        ax.text(.055,.16,"KNOWN FEATURES RE-APPEAR, BUT IN DIFFERENT CELLS",fontsize=10,color=NAVY,weight="bold")
        ax.text(.055,.105,"DSTK100 10168: rank #2 in improvement association (active fraction, E/V −7.67).     CodeLlama 27: rank #45 in local paired regression.",fontsize=9.5,color=MUTED)
        ax.text(.055,.07,"New candidate 13147: rank #1 in global regression association (mean, E/V 5.24, support 95).",fontsize=9.5,color=CORAL,weight="bold")
        p=outdir/"slide_23_eight_screens.png";save(fig,p,pdf);images.append(p)

        fig,ax=setup("α=3 SWEEP: MANY DIRECTIONS MOVE THE SAME TASKS", "A broad causal screen exposed task sensitivity—and the need for placebo directions")
        card(ax,.05,.25,.42,.50,"DSTK100 · 30 selected features","Baseline 15 / 80\n25 / 30 features had positive net change\nBest arm: 3048, +4 net\n29 / 30 corrected task /1030",TEAL)
        card(ax,.53,.25,.42,.50,"CodeLlama · 37 selected features","Baseline 0 / 50\n11 / 37 produced one official pass\n10 / 11 corrected the same task: /119\n13147 alone corrected /490",CORAL)
        ax.text(.05,.16,"INTERPRETATION",fontsize=10,color=NAVY,weight="bold")
        ax.text(.05,.105,"A pass transition is causal evidence that the residual perturbation mattered—but not yet that the selected latent was specific.",fontsize=10,color=MUTED)
        ax.text(.05,.07,"Repeated correction of the same tasks motivated orthogonal shams, matched latent controls, bidirectional steering, and dose curves.",fontsize=10,color=NAVY,weight="bold")
        p=outdir/"slide_24_alpha3_sweep.png";save(fig,p,pdf);images.append(p)

        fig,ax=setup("TWO CANDIDATES, TWO SEMANTIC PROFILES", "Natural TopK activations and CrossCoder geometry suggest different mechanisms")
        card(ax,.05,.22,.43,.56,"DSTK100 FEATURE 3048","TOP TOKENS\n_ · test · Func · from · Task · import · def\n\nCONTEXTS\nhelper functions · test blocks · imports\ntransitions into additional code units\n\nGEOMETRY\ndecoder cosine base↔FT: 0.961\nactive in 1,427 / 2,608 texts",TEAL)
        card(ax,.52,.22,.43,.56,"CODELLAMA FEATURE 13147","TOP TOKENS\nindentation · newline · # · ... · TODO · pass · return\n\nCONTEXTS\n# Hint: Use ... · placeholders\nincomplete or scaffolded implementations\n\nGEOMETRY\ndecoder cosine base↔merged: 0.275\nactive in 718 / 2,608 texts",CORAL)
        ax.text(.05,.13,"3048",fontsize=11,color=TEAL,weight="bold");ax.text(.105,.13,"shared program-structure direction; semantically broad",fontsize=10,color=MUTED)
        ax.text(.52,.13,"13147",fontsize=11,color=CORAL,weight="bold");ax.text(.58,.13,"model-differential incomplete-implementation direction",fontsize=10,color=MUTED)
        p=outdir/"slide_25_semantics.png";save(fig,p,pdf);images.append(p)

        fig,ax=setup("FEATURE 3048: THE EFFECT IS NOT SPECIFIC", "Seed 50000 · DeepSeek base · 80 selected contamination improvements")
        target=np.array([1,1,1,3,2,2,2]);sh=np.array([[0,3,2,3,4,5,5],[0,1,2,2,5,7,8],[0,0,2,3,3,4,5]]);co=np.array([[2,1,4,3,4,4,6],[0,1,2,5,7,7,9],[2,1,1,5,6,9,8]])
        line_chart(fig,[.07,.24,.54,.49],"FAIL → PASS IN THE BASE MODEL",target,sh,co)
        card(ax,.65,.46,.30,.27,"OUTPUTS CHANGED","3048: 13 → 37 / 80 across dose\nAt |α|=3: 27 changed\n8 test-marker changes\n4 import changes\n8 other logic/text changes",TEAL)
        card(ax,.65,.18,.30,.22,"REVERSE DIRECTION","Adding 3048 to FT:\n60 → 55 passes by +5/+6\n6 → 40 outputs changed\nNo reverse sham control yet",CORAL)
        ax.text(.07,.13,"At |α|=3: 3048=3 passes · shams=2–3 · latent controls=3–5",fontsize=10,color=NAVY,weight="bold")
        ax.text(.07,.085,"The selected direction changes behavior, but matched and random perturbations are equally or more effective.",fontsize=9.5,color=CORAL)
        p=outdir/"slide_26_f3048_controls.png";save(fig,p,pdf);images.append(p)

        fig,ax=setup("FEATURE 13147: COHERENT, BUT PLACEBO-SENSITIVE", "Seed 50000 · CodeLlama merged · 50 selected regressions")
        target=np.array([0,0,0,0,1,1,1]);sh=np.array([[0,0,0,0,0,0,0],[0,0,0,0,1,1,2],[0,0,0,1,1,1,1]]);co=np.zeros((3,7))
        line_chart(fig,[.07,.24,.54,.49],"FAIL → PASS IN THE MERGED MODEL",target,sh,co)
        card(ax,.65,.46,.30,.27,"SEMANTIC CHANGE","13147 changes 8 → 35 / 50 outputs\nReturn-line changes: 3 → 18\nAt −4/−5/−6, /490 gains\nparse → dump JSON → return result",CORAL)
        card(ax,.65,.18,.30,.22,"REVERSE DIRECTION","Adding 13147 to base:\n50 → 47 passes at +6\n7 → 25 outputs changed\nLosses emerge from +2",TEAL)
        ax.text(.07,.13,"Matched latent controls: 0 passes at every dose · Shams: up to 2 passes",fontsize=10,color=NAVY,weight="bold")
        ax.text(.07,.085,"The semantic match is strong, but one sham also repairs /490 at −6; specificity remains exploratory.",fontsize=9.5,color=CORAL)
        p=outdir/"slide_27_f13147_controls.png";save(fig,p,pdf);images.append(p)

        fig,ax=setup("THE 116 ARMS ALTER MORE THAN PASS / FAIL", "Exact extraction-v4 evaluated code reveals broad, dose-dependent trajectory changes")
        card(ax,.05,.23,.43,.54,"3048 FAMILY · 80 TASKS","SELECTED FEATURE, −0.5 → −6\nchanged outputs: 13 → 37\ntest-marker changes: 3 → 9\nimport changes: 4 → 9\nother logic/text: 3 → 7\n\nSHAMS AT −6\n29–39 outputs changed · 5–8 passes\n\nLATENT CONTROLS AT −6\n42–47 changed · 6–9 passes",TEAL)
        card(ax,.52,.23,.43,.54,"13147 FAMILY · 50 TASKS","SELECTED FEATURE, −0.5 → −6\nchanged outputs: 8 → 35\nreturn changes: 3 → 18\nother logic/text: 2 → 11\n\nSHAMS AT −6\n26–30 changed · return changes 14–18\n\nLATENT CONTROLS AT −6\n19–24 changed · 0 passes",CORAL)
        ax.text(.05,.14,"RULE-BASED CHANGE TAXONOMY",fontsize=10,color=NAVY,weight="bold")
        ax.text(.05,.095,"test markers → imports → returns → function structure → comments → other logic/text (mutually exclusive, first matching category)",fontsize=9.5,color=MUTED)
        ax.text(.05,.06,"Steering effects are real and structured, but high-dose output change is not itself evidence of latent specificity.",fontsize=9.5,color=NAVY,weight="bold")
        p=outdir/"slide_28_change_taxonomy.png";save(fig,p,pdf);images.append(p)

        fig,ax=setup("REVISED EVIDENCE LADDER", "Controls separate trajectory control from feature-specific causal explanation")
        card(ax,.05,.49,.28,.27,"SUPPORTED","CrossCoder screens recover stable behavioral contrasts.\n\nSelected directions causally alter code trajectories and official outcomes.",TEAL)
        card(ax,.36,.49,.28,.27,"EXPLORATORY","13147 has a coherent meaning↔failure-mode link and a stable /490 repair across −4 to −6 in seed 50000.",GOLD)
        card(ax,.67,.49,.28,.27,"NOT SUPPORTED","3048 does not outperform shams or latent controls.\n\nA pass transition alone does not establish semantic specificity.",CORAL)
        ax.text(.05,.38,"CRITICAL REPLICATION NOTE",fontsize=11,color=CORAL,weight="bold")
        ax.text(.05,.325,"The discovery α=−3 run used seed 1000+task_idx (1490 for /490). The focused sweep used input seed 50000.",fontsize=10,color=NAVY)
        ax.text(.05,.275,"Seed 1490 repaired /490 at −3; seed 50000 required −4 to −6. These are robustness runs—not exact replications.",fontsize=10,color=NAVY,weight="bold")
        ax.text(.05,.18,"NEXT CONFIRMATORY TEST",fontsize=11,color=TEAL,weight="bold")
        ax.text(.05,.125,"Re-run baseline, target feature, 3 shams, and 3 latent controls with the canonical seed convention and multiple paired seeds per α.",fontsize=10,color=MUTED)
        ax.text(.05,.075,"Feature 6404 remains the strongest controlled causal result in the project.",fontsize=11,color=NAVY,weight="bold")
        p=outdir/"slide_29_revised_claims.png";save(fig,p,pdf);images.append(p)
    return images,pdf_path

def append_pptx(deck,images,archive):
    archive.parent.mkdir(parents=True,exist_ok=True);shutil.copy2(deck,archive)
    prs=Presentation(deck); blank=prs.slide_layouts[6]
    for image in images:
        slide=prs.slides.add_slide(blank);slide.shapes.add_picture(str(image),0,0,width=prs.slide_width,height=prs.slide_height)
    prs.save(deck)

def main():
    ap=argparse.ArgumentParser();ap.add_argument("--deck",type=Path,default=Path("MODEL DIFFS.pptx"));ap.add_argument("--output-dir",type=Path,required=True);ap.add_argument("--archive",type=Path,required=True)
    a=ap.parse_args();images,pdf=build(a.output_dir);append_pptx(a.deck,images,a.archive);print(pdf)
if __name__=="__main__":main()
