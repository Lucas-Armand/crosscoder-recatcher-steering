#!/usr/bin/env python3
from pathlib import Path
from zipfile import ZipFile
import io,re
import matplotlib.pyplot as plt
from matplotlib.patches import FancyBboxPatch,Rectangle,FancyArrowPatch
from pptx import Presentation
from PIL import Image
R=Path('/home/lucas/crosscoder-recatcher-steering');SRC=R/'MODEL DIFFS_CAUSAL_STEERING_RESULTS_STORY_VISUALS_V3_2026-08-31.pptx';DEST=SRC;PDF=R/'MODEL DIFFS_CAUSAL_STEERING_RESULTS_STORY_VISUALS_V3_2026-08-31.pdf';OUT=R/'presentation/generated/story_visuals_v3_20260831';p=OUT/'03.png'
BG='#F7F4EE';NAVY='#10233F';PURPLE='#7257C7';ORANGE='#EE8A2D';TEAL='#159C9C';MUTED='#687386';WHITE='#FFFFFF';GRID='#D8D9D8'
def card(a,x,y,w,h,title,body,c):
 a.add_patch(FancyBboxPatch((x,y),w,h,boxstyle='round,pad=.007',facecolor=WHITE,edgecolor='#E1DED7'));a.add_patch(Rectangle((x,y+h-.009),w,.009,color=c,lw=0));a.text(x+.014,y+h-.035,title,fontsize=8.6,color=NAVY,weight='bold',va='top');a.text(x+.014,y+h-.082,body,fontsize=7,color=MUTED,va='top')
def arrow(a,x1,y1,x2,y2,c):a.add_patch(FancyArrowPatch((x1,y1),(x2,y2),arrowstyle='-|>',mutation_scale=12,color=c,lw=1.5))
f=plt.figure(figsize=(13.333,7.5),facecolor=BG);a=f.add_axes([0,0,1,1]);a.axis('off');a.text(.045,.955,'CROSSCODER MODEL DIFFING',fontsize=8.5,color=TEAL,weight='bold',va='top');a.text(.045,.885,'ONE ANALYSIS FRAMEWORK, TWO VERY DIFFERENT MODEL CHANGES',fontsize=22,color=NAVY,weight='bold',va='top');a.text(.045,.815,'Same token IDs turn each model pair into a shared sparse vocabulary',fontsize=10,color=MUTED,va='top')
for y,name,left,right,col,note in [(.49,'DEEPSEEK','BASE','FINE-TUNED',PURPLE,'specialization through fine-tuning'),(.20,'CODELLAMA','BASE','MERGED',ORANGE,'specialization through model merging')]:
 a.text(.055,y+.13,name,fontsize=11.5,color=col,weight='bold');a.text(.055,y+.085,note,fontsize=8.2,color=MUTED)
 card(a,.19,y,.115,.16,left+' · L16','residual stream',col);card(a,.32,y,.115,.16,right+' · L16','residual stream',col);a.text(.312,y-.027,'identical evaluated token IDs',ha='center',fontsize=7.4,color=col,weight='bold');arrow(a,.307,y+.08,.475,y+.08,col);arrow(a,.437,y+.08,.475,y+.08,col)
 card(a,.48,y-.005,.17,.17,'CROSSCODER','16,384 latents · exact TopK-100',col);arrow(a,.655,y+.08,.705,y+.08,col);a.add_patch(FancyBboxPatch((.71,y-.005),.21,.17,boxstyle='round,pad=.007',facecolor=WHITE,edgecolor='#E1DED7'))
 for j in range(12):a.add_patch(Rectangle((.735+j*.013,y+.092),.008,.028,facecolor=col if j in (0,3,7,10) else '#E4E5E4',edgecolor='none'))
 a.text(.73,y+.052,'100 active features',fontsize=8.8,color=NAVY,weight='bold');a.text(.73,y+.015,'per evaluated token',fontsize=7.6,color=MUTED)
a.text(.055,.085,'Different specialization mechanisms; identical analytical question: which shared latent differences track—and control—behavior?',fontsize=9.4,color=NAVY,weight='bold');a.text(.95,.035,'CrossCoder model diffing · canonical readout · 2026-08-31',ha='right',fontsize=7.4,color=MUTED);f.savefig(p,dpi=180,facecolor=BG);plt.close(f)
prs=Presentation(SRC);s=prs.slides[2]
for sh in list(s.shapes):s.shapes._spTree.remove(sh._element)
s.shapes.add_picture(str(p),0,0,width=prs.slide_width,height=prs.slide_height);prs.save(DEST)
with ZipFile(DEST) as z:
 pages=[]
 for n in range(1,len(prs.slides)+1):
  rel=z.read(f'ppt/slides/_rels/slide{n}.xml.rels').decode();slide=z.read(f'ppt/slides/slide{n}.xml').decode();mp={i:t for i,t in re.findall(r'Id="([^\"]+)"[^>]+Target="([^\"]+)"',rel)};vals=[]
  for rid in re.findall(r'r:embed="([^\"]+)"',slide):
   t=mp.get(rid,'')
   if t.startswith('../media/'):d=z.read('ppt/'+t[3:]);vals.append((len(d),d))
  pages.append(Image.open(io.BytesIO(max(vals)[1])).convert('RGB'))
pages[0].save(PDF,save_all=True,append_images=pages[1:],resolution=180);print('fixed',DEST,PDF)
