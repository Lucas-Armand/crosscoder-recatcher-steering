#!/usr/bin/env python3
"""Insert cohort focus, paired examples, and candidate-pool slides."""
from pathlib import Path
import argparse,csv,shutil,subprocess
import matplotlib.pyplot as plt
from matplotlib.backends.backend_pdf import PdfPages
from matplotlib.patches import FancyBboxPatch,Rectangle
from pptx import Presentation

BG="#F7F4EE";NAVY="#10233F";PURPLE="#7257C7";ORANGE="#EE8A2D";TEAL="#159C9C";MUTED="#687386";WHITE="#FFFFFF";GRID="#D8D9D8";RED="#D85A56";GREEN="#338A63";GOLD="#DDA531"

def setup(title,sub=""):
 f=plt.figure(figsize=(13.333,7.5),facecolor=BG);a=f.add_axes([0,0,1,1]);a.axis("off")
 a.text(.045,.955,"FROM BEHAVIORAL TRANSITIONS TO ERROR-FOCUSED SCREENING",fontsize=8.5,color=TEAL,weight="bold",va="top")
 a.text(.045,.885,title,fontsize=23,color=NAVY,weight="bold",va="top")
 if sub:a.text(.045,.815,sub,fontsize=10.2,color=MUTED,va="top")
 a.text(.95,.035,"CrossCoder model diffing · revised draft · 2026-08-26",ha="right",fontsize=7.5,color=MUTED)
 return f,a

def card(a,x,y,w,h,title,body,accent,ts=10.5,bs=8.2):
 a.add_patch(FancyBboxPatch((x,y),w,h,boxstyle="round,pad=.007,rounding_size=.011",facecolor=WHITE,edgecolor="#E1DED7"));a.add_patch(Rectangle((x,y+h-.01),w,.01,color=accent,lw=0));a.text(x+.016,y+h-.04,title,fontsize=ts,color=NAVY,weight="bold",va="top");a.text(x+.016,y+h-.085,body,fontsize=bs,color=MUTED,va="top",linespacing=1.3)

def save(f,p,pdf):f.savefig(p,dpi=175,facecolor=BG);pdf.savefig(f,facecolor=BG);plt.close(f)

def top(repo,model,tr):
 root=repo/("reports/focused_subtype_screening_dstk100_contamination_v1" if model=="DS" else "reports/focused_subtype_screening_codellama_wrong_logic_v1")
 cells=["association_global","association_local","paired_global","paired_local"]
 out={}
 for c in cells:
  with open(root/f"{c}_absolute.csv",newline="") as f:out[c]=[int(r["feature_id"]) for r in list(csv.DictReader(f))[:10]]
 return out

def build(repo,out):
 out.mkdir(parents=True,exist_ok=True);imgs=[]
 with PdfPages(out/"focus_and_pool_insert.pdf") as pdf:
  f,a=setup("WE FOCUS ON TWO ONE-SIDED BEHAVIORAL TRANSITIONS","Taxonomy turns aggregate model differences into semantically targeted cohorts")
  card(a,.055,.39,.42,.34,"DEEPSEEK · 215 IMPROVEMENTS","base fail → fine-tuned pass\n\n119/215: generated test/import contamination\n48: wrong output or logic\n16: missing name/import\n32: other primary categories\n\nMechanistic focus: 80 contamination cases",PURPLE,11,8.3)
  card(a,.53,.39,.42,.34,"CODELLAMA · 291 REGRESSIONS","base pass → merged fail\n\n120/291: API/type mismatch\n50: wrong logic/other runtime\n43: edge case/exception\n78: generation, import, syntax, commentary\n\nMechanistic focus: 50 wrong-logic/runtime cases\nImprovements: only 4 positives · EXCLUDED",ORANGE,11,7.8)
  card(a,.055,.17,.42,.14,"HOW WE CLASSIFIED","Rule/AI-assisted taxonomy + human spot checks",TEAL,9.3,8)
  card(a,.53,.17,.42,.14,"IMPORTANT QUALIFICATION","Primary = most salient observed error—not a unique root cause. CodeLlama failures are heterogeneous and often multilabel.",RED,9.3,7.8)
  p=out/"05_transition_focus.png";save(f,p,pdf);imgs.append(p)

  f,a=setup("THE TWO FOCUSED FAILURE MODES HAVE DIFFERENT SHAPES","Paired extraction-v4 examples motivate where local screening should look")
  card(a,.055,.39,.42,.35,"DEEPSEEK /5 · CONTAMINATION","BASE · FAIL\nCorrect-looking body continues into:\n`#tests/test_task_func.py`\n`import pytest`\n\nFINE-TUNED · PASS\nCompletes the function and stops after a short example.\n\nFirst textual divergence ≈76% of shorter code;\ntarget boundary is late/post-solution.",PURPLE,10.5,7.9)
  card(a,.53,.39,.42,.35,"CODELLAMA /490 · WRONG LOGIC/API CONTRACT","BASE · PASS\nparses XML → writes JSON file → returns result\n\nMERGED · FAIL\n`return xmltodict.parse(s)`\nomits the required file write side effect\n\nFirst divergence ≈81% of base code;\ntarget region is the late function body.",ORANGE,10.5,7.9)
  a.plot([.08,.44],[.23,.23],color=GRID,lw=7,solid_capstyle="round");a.plot([.08,.35],[.23,.23],color=PURPLE,lw=7,solid_capstyle="round");a.scatter([.35],[.23],s=90,color=RED,zorder=3);a.text(.08,.16,"prompt / shared scaffold",fontsize=8,color=MUTED);a.text(.335,.16,"late divergence",fontsize=8,color=RED,weight="bold")
  a.plot([.56,.92],[.23,.23],color=GRID,lw=7,solid_capstyle="round");a.plot([.56,.85],[.23,.23],color=ORANGE,lw=7,solid_capstyle="round");a.scatter([.85],[.23],s=90,color=RED,zorder=3);a.text(.56,.16,"docstring / shared scaffold",fontsize=8,color=MUTED);a.text(.83,.16,"late body",fontsize=8,color=RED,weight="bold")
  a.text(.055,.085,"Local window = ±10% around first normalized divergence; it is a reproducible temporal proxy, not an oracle for semantic causality.",fontsize=9,color=NAVY,weight="bold")
  p=out/"06_paired_error_examples.png";save(f,p,pdf);imgs.append(p)

  ds=top(repo,"DS","improvement");cl=top(repo,"CL","regression")
  f,a=setup("ERROR-FOCUSED TOP-10 LISTS DEFINE THE NEW α=3 CANDIDATE POOL","Three summaries per cell: max · mean · active fraction; ranking by absolute E/V")
  for x,title,d,col in [(.04,"DEEPSEEK · 80 CONTAMINATION IMPROVEMENTS",ds,PURPLE),(.515,"CODELLAMA · 50 WRONG-LOGIC REGRESSIONS",cl,ORANGE)]:
   a.text(x,.755,title,fontsize=11,color=col,weight="bold")
   short=[("association_global","ASSOC · GLOBAL"),("association_local","ASSOC · LOCAL"),("paired_global","PAIRED · GLOBAL"),("paired_local","PAIRED · LOCAL")]
   for j,(suffix,label) in enumerate(short):
    key=next(k for k in d if k.endswith(suffix));xx=x+j*.112
    a.add_patch(Rectangle((xx,.26),.103,.42,facecolor=WHITE,edgecolor=GRID));a.add_patch(Rectangle((xx,.625),.103,.055,facecolor=col,edgecolor=col));a.text(xx+.0515,.652,label,ha="center",va="center",fontsize=6.5,color=WHITE,weight="bold")
    for i,v in enumerate(d[key]):a.text(xx+.0515,.59-i*.031,str(v),ha="center",fontsize=7.2,color=NAVY,weight="bold" if sum(v in z for z in d.values())>1 else "normal")
   union=set().union(*map(set,d.values()));over=sum(sum(v in z for z in d.values())>1 for v in union)
   a.text(x,.205,f"40 nominations → {len(union)} unique features",fontsize=10,color=col,weight="bold");a.text(x,.165,f"{over} feature(s) appear in more than one list",fontsize=8.5,color=MUTED)
  a.add_patch(FancyBboxPatch((.16,.075),.68,.065,boxstyle="round,pad=.008",facecolor="#EEEAE2",edgecolor=GRID));a.text(.50,.108,"NEW α=3 SWEEP: exactly 35 DeepSeek + 32 CodeLlama unique features · no post-hoc additions",ha="center",va="center",fontsize=10,color=NAVY,weight="bold")
  a.text(.50,.065,"Hard gate: baseline must reproduce the original raw completion byte for byte before any steered arm runs.",ha="center",fontsize=8.5,color=RED,weight="bold")
  p=out/"10_candidate_pool.png";save(f,p,pdf);imgs.append(p)
 return imgs

def insert(prs,index,p):
 blank=prs.slide_layouts[6];s=prs.slides.add_slide(blank);s.shapes.add_picture(str(p),0,0,width=prs.slide_width,height=prs.slide_height);sid=prs.slides._sldIdLst[-1];prs.slides._sldIdLst.remove(sid);prs.slides._sldIdLst.insert(index,sid)

def main():
 ap=argparse.ArgumentParser();ap.add_argument("--repo",type=Path,default=Path("."));ap.add_argument("--source",type=Path,required=True);ap.add_argument("--dest",type=Path,required=True);ap.add_argument("--output-dir",type=Path,required=True);x=ap.parse_args()
 imgs=build(x.repo,x.output_dir);shutil.copy2(x.source,x.dest);prs=Presentation(x.dest)
 # Insert two slides after original slide 4, then pool slide after original slide 7.
 insert(prs,4,imgs[0]);insert(prs,5,imgs[1]);insert(prs,9,imgs[2]);prs.save(x.dest);print(x.dest,len(prs.slides))
if __name__=="__main__":main()
