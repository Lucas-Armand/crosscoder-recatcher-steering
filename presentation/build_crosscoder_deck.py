#!/usr/bin/env python3
"""Build the DSTK100 model-diffing presentation and its analytical charts."""

from __future__ import annotations

import csv
import json
import math
from pathlib import Path

import matplotlib.pyplot as plt
import numpy as np
from matplotlib.ticker import PercentFormatter
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
OUT = ROOT / "presentation" / "generated"
CHARTS = OUT / "charts"
DATA = OUT / "data"
OUT.mkdir(parents=True, exist_ok=True)
CHARTS.mkdir(parents=True, exist_ok=True)
DATA.mkdir(parents=True, exist_ok=True)

EXTERNAL = Path("/tmp/cc_deck_data")
ALL_STATS = Path("/tmp/dstk100_all_feature_statistics.csv")

NAVY = "111827"
INK = "172033"
BLUE = "2563EB"
CYAN = "06B6D4"
GREEN = "10B981"
ORANGE = "F59E0B"
RED = "EF4444"
PURPLE = "8B5CF6"
GRAY = "64748B"
LIGHT = "F8FAFC"
PALE = "E2E8F0"
WHITE = "FFFFFF"


def rgb(value: str) -> RGBColor:
    return RGBColor.from_string(value)


def style_axes(ax):
    ax.spines[["top", "right"]].set_visible(False)
    ax.spines[["left", "bottom"]].set_color("#CBD5E1")
    ax.tick_params(colors="#475569", labelsize=9)
    ax.grid(axis="x", color="#E2E8F0", linewidth=.8, zorder=0)
    ax.set_axisbelow(True)


def savefig(fig, name):
    path = CHARTS / name
    fig.savefig(path, dpi=220, bbox_inches="tight", facecolor="white")
    plt.close(fig)
    return path


def chart_transition_balance():
    fig, ax = plt.subplots(figsize=(8.8, 2.6))
    ax.barh([0], [257], color="#10B981", height=.52, label="Improvements")
    ax.barh([0], [86], left=[257], color="#EF4444", height=.52, label="Regressions")
    ax.text(128.5, 0, "257 improvements", ha="center", va="center", color="white", weight="bold", fontsize=15)
    ax.text(300, 0, "86 regressions", ha="center", va="center", color="white", weight="bold", fontsize=13)
    ax.set_xlim(0,343); ax.set_yticks([]); ax.set_xlabel("One-sided behavioral transitions")
    style_axes(ax); ax.spines["left"].set_visible(False)
    return savefig(fig, "transition_balance.png")


def chart_failure_taxonomy():
    labels = ["Truncation / extraction", "Wrong output / logic", "File / path", "Missing name / import", "Other"]
    vals = [19,18,11,11,20]
    fig, ax = plt.subplots(figsize=(8.7,4.5))
    y=np.arange(len(labels))[::-1]
    ax.barh(y, vals, color=["#F59E0B","#EF4444","#8B5CF6","#06B6D4","#94A3B8"], height=.62)
    for yi,v in zip(y,vals): ax.text(v+.35,yi,str(v),va="center",weight="bold",color="#172033")
    ax.set_yticks(y,labels); ax.set_xlim(0,22); ax.set_xlabel("BigCodeBench regressions (n=79)")
    style_axes(ax)
    return savefig(fig,"failure_taxonomy.png")


TOP10 = [
    (8587,"max",4.723,54),(6404,"max",4.465,47),(12956,"max",3.990,74),
    (8959,"max",3.859,77),(6684,"max",3.856,79),(9716,"max",3.774,72),
    (14818,"max",3.713,74),(15246,"max",3.710,48),(2449,"max",3.659,63),
    (10967,"early_max",3.611,66),
]


def chart_ev_ranking():
    fig, ax = plt.subplots(figsize=(8.7,5.0))
    rows=TOP10[::-1]; y=np.arange(len(rows)); vals=[r[2] for r in rows]
    colors=["#2563EB" if r[0]==6404 else "#94A3B8" for r in rows]
    ax.barh(y,vals,color=colors,height=.66)
    ax.set_yticks(y,[f"{r[0]}  ·  {r[1]}" for r in rows])
    for yi,r in zip(y,rows): ax.text(r[2]+.035,yi,f"{r[2]:.2f}   ({r[3]}/79)",va="center",fontsize=9,weight="bold" if r[0]==6404 else "normal")
    ax.set_xlim(0,5.35);ax.set_xlabel("E/V score   ·   support shown in parentheses")
    style_axes(ax)
    return savefig(fig,"ev_top10.png")


def chart_sensitivity():
    cases=["BCB regression","BCB improvement","HE+ regression","HE+ improvement"]
    values=np.array([[10,10,10],[10,2,1],[10,7,4],[10,10,8]])
    cols=["Canonical\nTop 10","E/V only\nmax + early","All four\naggregations"]
    fig,ax=plt.subplots(figsize=(7.5,3.8))
    im=ax.imshow(values,cmap="Blues",vmin=0,vmax=10,aspect="auto")
    for i in range(values.shape[0]):
        for j in range(values.shape[1]): ax.text(j,i,f"{values[i,j]}/10",ha="center",va="center",weight="bold",color="white" if values[i,j]>6 else "#172033",fontsize=12)
    ax.set_xticks(range(3),cols);ax.set_yticks(range(4),cases);ax.tick_params(length=0)
    for s in ax.spines.values():s.set_visible(False)
    return savefig(fig,"screening_sensitivity.png")


def chart_semantic_timing():
    features=[16383,14481,12956,6404,8587,8294,11785]
    names={16383:"validation",14481:"plot/dataframe",12956:"expected output",6404:"assumptions",8587:"numeric zero",8294:"Python meta",11785:"model fitting"}
    data={f:[] for f in features}
    p=EXTERNAL/"feature_failure_relations.csv"
    if p.exists():
        with p.open(newline="") as fh:
            for r in csv.DictReader(fh):
                f=int(r["feature_id"])
                if f in data and int(r["active_tokens"])>0: data[f].append(float(r["first_percent"]))
    else:
        raise FileNotFoundError(p)
    fig,ax=plt.subplots(figsize=(9.2,5.1))
    vals=[data[f] for f in features]
    parts=ax.violinplot(vals,positions=np.arange(len(features)),vert=False,showextrema=False,widths=.72)
    for body in parts["bodies"]: body.set_facecolor("#93C5FD");body.set_edgecolor("#2563EB");body.set_alpha(.8)
    ax.boxplot(vals,positions=np.arange(len(features)),vert=False,widths=.22,showfliers=False,patch_artist=True,boxprops=dict(facecolor="white",edgecolor="#172033"),medianprops=dict(color="#EF4444",linewidth=2),whiskerprops=dict(color="#64748B"),capprops=dict(color="#64748B"))
    ax.set_yticks(np.arange(len(features)),[f"{f} · {names[f]}\n(n={len(data[f])})" for f in features])
    ax.set_xlim(0,100);ax.xaxis.set_major_formatter(PercentFormatter(100));ax.set_xlabel("First activation position in evaluated code")
    style_axes(ax);ax.grid(axis="x",color="#E2E8F0")
    return savefig(fig,"feature_first_activation.png")


def read_task(task_id="BigCodeBench/668"):
    def lookup(path):
        with path.open() as fh:
            for row in map(json.loads,fh):
                if row["task_id"]==task_id:return row
        raise KeyError(task_id)
    return lookup(EXTERNAL/"bigcodebench__baseline_alpha0_results.jsonl"),lookup(EXTERNAL/"bigcodebench__f6404_alpha_neg2_results.jsonl")


def chart_activation_trace():
    _,steer=read_task()
    trace=steer["topk_gate_trace"]
    total=max(1,steer["topk_gate_total_steps"])
    active=[x for x in trace if x["active"]]
    fig,ax=plt.subplots(figsize=(9.0,1.65))
    ax.hlines(0,0,100,color="#CBD5E1",linewidth=5)
    for x in active:
        pos=100*x["step"]/total; strength=x["feature_activation"]
        ax.vlines(pos,-.17,.17,color="#EF4444",linewidth=2+strength/2)
        ax.scatter([pos],[0],s=35+strength*12,color="#EF4444",edgecolor="white",zorder=3)
    ax.text(0,.28,"start",ha="left",color="#64748B");ax.text(100,.28,"end",ha="right",color="#64748B")
    ax.set_xlim(0,100);ax.set_ylim(-.45,.55);ax.set_yticks([]);ax.set_xlabel("Generated-token position (%)")
    ax.xaxis.set_major_formatter(PercentFormatter(100));
    for s in ax.spines.values():s.set_visible(False)
    ax.grid(False)
    return savefig(fig,"f6404_task668_trace.png")


def chart_controls():
    labels=["Feature 6404","Feature 6757","Feature 9388","Feature 6509","Orthogonal sham"]
    passes=[19,3,2,0,0];clean=[41,6,6,2,0]
    x=np.arange(len(labels));w=.34
    fig,ax=plt.subplots(figsize=(9,4.6))
    ax.bar(x-w/2,passes,w,label="Official passes",color="#2563EB")
    ax.bar(x+w/2,clean,w,label="Contamination removed",color="#F59E0B")
    for i,v in enumerate(passes):ax.text(i-w/2,v+.6,str(v),ha="center",weight="bold",color="#172033")
    for i,v in enumerate(clean):ax.text(i+w/2,v+.6,str(v),ha="center",weight="bold",color="#172033")
    ax.set_xticks(x,labels,rotation=14,ha="right");ax.set_ylabel("Tasks (of 80)");ax.set_ylim(0,46);ax.legend(frameon=False,ncol=2,loc="upper right")
    style_axes(ax);ax.grid(axis="y",color="#E2E8F0");ax.grid(axis="x",visible=False)
    return savefig(fig,"matched_controls.png")


def chart_dense_curve():
    src=ROOT/"reports/dstk100_f6404_dense_curve_v1/curve.csv"
    rows=list(csv.DictReader(src.open()))
    alpha=np.array([float(r["alpha"]) for r in rows]);order=np.argsort(alpha);alpha=alpha[order]
    passed=np.array([int(rows[i]["passed"]) for i in order]);removed=np.array([int(rows[i]["test_import_contamination_removed"]) for i in order])
    fig,ax=plt.subplots(figsize=(9,4.7))
    ax.plot(alpha,passed,"o-",color="#2563EB",linewidth=2.8,label="Official passes")
    ax.plot(alpha,removed,"s--",color="#F59E0B",linewidth=2.5,label="Contamination removed")
    ax.axvline(-2,color="#CBD5E1",linewidth=1.5);ax.annotate("effective window",(-2,41),xytext=(-3.4,43),arrowprops=dict(arrowstyle="->",color="#64748B"),color="#172033",weight="bold")
    ax.set_xlabel("Steering coefficient α");ax.set_ylabel("Tasks (of 80)");ax.set_ylim(-1,46);ax.legend(frameon=False,ncol=2)
    style_axes(ax);ax.grid(axis="both",color="#E2E8F0")
    return savefig(fig,"dense_curve_redrawn.png")


def export_derived_data():
    with (DATA/"bigcodebench_regression_ev_top10.csv").open("w",newline="") as f:
        w=csv.writer(f);w.writerow(["rank","feature_id","aggregation","effect_to_variability","support","positives"])
        for i,r in enumerate(TOP10,1):w.writerow([i,*r,79])
    with (DATA/"screening_sensitivity_top10_overlap.csv").open("w",newline="") as f:
        w=csv.writer(f);w.writerow(["case","canonical","ev_only_max_early","all_four_aggregations"])
        w.writerows([["bigcodebench_regression",10,10,10],["bigcodebench_improvement",10,2,1],["humanevalplus_regression",10,7,4],["humanevalplus_improvement",10,10,8]])


class Deck:
    def __init__(self):
        self.prs=Presentation();self.prs.slide_width=Inches(13.333);self.prs.slide_height=Inches(7.5)
        self.blank=self.prs.slide_layouts[6]

    def add_bg(self,slide,color=LIGHT):
        shape=slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,0,0,self.prs.slide_width,self.prs.slide_height)
        shape.fill.solid();shape.fill.fore_color.rgb=rgb(color);shape.line.fill.background();slide.shapes._spTree.remove(shape._element);slide.shapes._spTree.insert(2,shape._element)

    def text(self,slide,text,x,y,w,h,size=20,color=INK,bold=False,font="Aptos",align=PP_ALIGN.LEFT,margin=.04,anchor=MSO_ANCHOR.TOP):
        box=slide.shapes.add_textbox(Inches(x),Inches(y),Inches(w),Inches(h));tf=box.text_frame;tf.clear();tf.margin_left=tf.margin_right=Inches(margin);tf.margin_top=tf.margin_bottom=Inches(margin);tf.vertical_anchor=anchor
        p=tf.paragraphs[0];p.alignment=align;r=p.add_run();r.text=text;r.font.name=font;r.font.size=Pt(size);r.font.bold=bold;r.font.color.rgb=rgb(color)
        return box

    def title(self,slide,title,kicker=None,dark=False):
        if kicker:self.text(slide,kicker.upper(),.65,.32,12,.3,10,CYAN if dark else BLUE,True)
        self.text(slide,title,.65,.72,12,1.0,27,WHITE if dark else NAVY,True)

    def footer(self,slide,n,dark=False):
        self.text(slide,"DSTK100 · preliminary causal model diffing",.65,7.16,7,.18,8,"94A3B8")
        self.text(slide,str(n),12.25,7.12,.35,.2,8,"94A3B8",align=PP_ALIGN.RIGHT)

    def image(self,slide,path,x,y,w,h=None):
        return slide.shapes.add_picture(str(path),Inches(x),Inches(y),width=Inches(w),height=Inches(h) if h else None)

    def card(self,slide,x,y,w,h,color=WHITE,line=PALE,radius=True):
        sh=slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE,Inches(x),Inches(y),Inches(w),Inches(h));sh.fill.solid();sh.fill.fore_color.rgb=rgb(color);sh.line.color.rgb=rgb(line);return sh

    def bullet(self,slide,text,x,y,w,size=17,color=INK,accent=BLUE):
        self.text(slide,"●",x,y,.25,.3,size-4,accent,True);self.text(slide,text,x+.3,y,w-.3,.55,size,color)

    def notes(self,slide,text):
        try:
            tf=slide.notes_slide.notes_text_frame
            tf.text=text
        except Exception:
            pass

    def new(self,bg=LIGHT):
        s=self.prs.slides.add_slide(self.blank);self.add_bg(s,bg);return s

    def save(self,path):self.prs.save(path)


def build_deck(charts):
    d=Deck();n=0
    def slide(bg=LIGHT):
        nonlocal n;n+=1;s=d.new(bg);d.footer(s,n,dark=bg==NAVY);return s

    s=slide(NAVY)
    d.text(s,"FROM MODEL DIFFS\nTO CAUSAL MECHANISMS",.7,1.05,8.3,1.65,38,WHITE,True)
    d.text(s,"Using a CrossCoder to explain behavioral differences\nbetween base and fine-tuned code models",.75,3.05,7.2,1.0,20,"CBD5E1")
    d.card(s,9.25,.8,3.2,5.45,"172033","334155")
    d.text(s,"The case",9.65,1.2,2.4,.35,12,CYAN,True)
    d.text(s,"The model solved\nthe task—then\nkept writing.",9.65,1.7,2.25,1.65,25,WHITE,True)
    d.text(s,"A single latent feature helped control that continuation.",9.65,4.15,2.25,1.0,16,"CBD5E1")
    d.notes(s,"Open with the behavioral paradox, not the architecture. The scientific question is whether a shared sparse representation can explain why the models diverge.")

    s=slide();d.title(s,"The evaluation tells us what changed. The CrossCoder asks why.","Research question")
    d.card(s,.75,2.0,5.65,3.8,"FFF7ED","FED7AA");d.card(s,6.92,2.0,5.65,3.8,"EFF6FF","BFDBFE")
    d.text(s,"BEHAVIOR",1.15,2.42,2,.3,13,ORANGE,True);d.text(s,"REPRESENTATION",7.32,2.42,2.5,.3,13,BLUE,True)
    for i,t in enumerate(["Base passes / fine-tuned fails","Base fails / fine-tuned passes","What kind of error occurred?"]):d.bullet(s,t,1.15,3.05+i*.67,4.8,17,accent=ORANGE)
    for i,t in enumerate(["Paired layer-16 residuals","Shared sparse features","Decoder directions for intervention"]):d.bullet(s,t,7.32,3.05+i*.67,4.8,17,accent=BLUE)
    d.text(s,"Behavior localizes the difference. Features turn it into a testable hypothesis.",1.0,6.25,11.3,.5,20,NAVY,True,align=PP_ALIGN.CENTER)

    s=slide();d.title(s,"DSTK100 gives both models a shared sparse vocabulary","Method")
    for x,label,sub,c in [(1.0,"BASE","layer-16 residual",ORANGE),(1.0,"FINE-TUNED","layer-16 residual",GREEN)]:
        pass
    d.card(s,.8,2.05,2.25,1.05,"FFF7ED","FDBA74");d.text(s,"BASE",1.05,2.25,1.7,.25,16,ORANGE,True,align=PP_ALIGN.CENTER);d.text(s,"layer 16",1.05,2.58,1.7,.2,11,GRAY,align=PP_ALIGN.CENTER)
    d.card(s,.8,3.55,2.25,1.05,"ECFDF5","6EE7B7");d.text(s,"FINE-TUNED",1.05,3.75,1.7,.25,16,GREEN,True,align=PP_ALIGN.CENTER);d.text(s,"layer 16",1.05,4.08,1.7,.2,11,GRAY,align=PP_ALIGN.CENTER)
    d.text(s,"same tokens",3.35,3.0,1.3,.35,13,GRAY,True,align=PP_ALIGN.CENTER)
    d.card(s,4.7,2.45,3.1,2.15,"EFF6FF","93C5FD");d.text(s,"CROSSCODER",5.05,2.85,2.4,.35,19,BLUE,True,align=PP_ALIGN.CENTER);d.text(s,"16,384 latents\nReLU + TopK 100",5.05,3.35,2.4,.7,14,INK,align=PP_ALIGN.CENTER)
    d.text(s,"→",3.8,3.13,.55,.5,28,BLUE,True,align=PP_ALIGN.CENTER);d.text(s,"→",8.05,3.13,.55,.5,28,BLUE,True,align=PP_ALIGN.CENTER)
    d.card(s,8.7,1.75,3.75,3.5,WHITE,"CBD5E1");
    for i in range(7):
        c=BLUE if i in (1,4) else PALE;d.card(s,9.15+i*.39,2.3,.23,.23,c,c,False)
    d.text(s,"100 active features",9.05,2.85,3.0,.3,17,NAVY,True,align=PP_ALIGN.CENTER);d.text(s,"per evaluated token",9.05,3.3,3.0,.3,14,GRAY,align=PP_ALIGN.CENTER)
    d.text(s,"2,608 code texts   ·   617,959 evaluated tokens   ·   same-text alignment",1.0,6.15,11.3,.45,18,NAVY,True,align=PP_ALIGN.CENTER)

    s=slide();d.title(s,"One-sided transitions expose where the models actually diverge","Behavioral map")
    d.image(s,charts['transitions'],1.05,2.0,11.2)
    d.text(s,"343",1.0,4.85,2,.65,34,NAVY,True,align=PP_ALIGN.CENTER);d.text(s,"one-sided transitions",1.0,5.5,2,.3,12,GRAY,align=PP_ALIGN.CENTER)
    d.text(s,"257",5.1,4.85,2,.65,34,GREEN,True,align=PP_ALIGN.CENTER);d.text(s,"base fail → FT pass",5.1,5.5,2,.3,12,GRAY,align=PP_ALIGN.CENTER)
    d.text(s,"86",9.2,4.85,2,.65,34,RED,True,align=PP_ALIGN.CENTER);d.text(s,"base pass → FT fail",9.2,5.5,2,.3,12,GRAY,align=PP_ALIGN.CENTER)

    s=slide();d.title(s,"“Regression” is an outcome—not a single mechanism","Error taxonomy")
    d.image(s,charts['taxonomy'],.8,1.7,7.7)
    d.card(s,8.8,1.9,3.7,3.8,"FFF7ED","FED7AA");d.text(s,"Why classify first?",9.15,2.3,3,.35,18,ORANGE,True)
    for i,t in enumerate(["Different errors imply different causal hypotheses.","A global fail label can mix incompatible mechanisms.","Taxonomy defines the cohort we should intervene on."]):d.bullet(s,t,9.15,3.0+i*.78,2.9,15,accent=ORANGE)

    s=slide();d.title(s,"We ranked stable contrasts—not raw activation alone","Feature screening")
    stages=[("4 summaries","max · early_max\nmean · active fraction"),("Δ per task","aggregate(FT text)\n− aggregate(base text)"),("Observed effect","mean Δ positives\n− mean Δ controls"),("200 permutations","shuffle labels;\nrecompute effect"),("E/V","effect ÷ null SD")]
    for i,(a,b) in enumerate(stages):
        x=.5+i*2.55;d.card(s,x,2.0,2.05,2.1,WHITE,"CBD5E1");d.text(s,a,x+.15,2.34,1.75,.35,16,BLUE if i<4 else RED,True,align=PP_ALIGN.CENTER);d.text(s,b,x+.15,2.96,1.75,.72,12,INK,align=PP_ALIGN.CENTER)
        if i<4:d.text(s,"→",x+2.1,2.72,.42,.4,23,GRAY,True,align=PP_ALIGN.CENTER)
    d.card(s,1.0,4.75,11.3,1.15,"EFF6FF","BFDBFE");d.text(s,"Minimum support",1.35,5.05,2.2,.3,14,BLUE,True);d.text(s,"max(3 tasks, 10% of positives)",3.0,5.02,3.4,.35,17,NAVY,True);d.text(s,"E/V is a permutation signal-to-noise score—not a z-score or a p-value.",6.35,5.0,5.45,.5,15,INK,True)

    s=slide();d.title(s,"Feature 6404 ranked #2—and the result was specification-robust","Screening result")
    d.image(s,charts['ev'],.6,1.5,8.2)
    d.card(s,9.15,1.8,3.25,3.85,"EFF6FF","93C5FD");d.text(s,"6404",9.55,2.2,2.45,.6,36,BLUE,True,align=PP_ALIGN.CENTER);d.text(s,"E/V 4.47",9.55,3.03,2.45,.35,19,NAVY,True,align=PP_ALIGN.CENTER);d.text(s,"47 / 79 support",9.55,3.55,2.45,.35,16,INK,align=PP_ALIGN.CENTER);d.text(s,"Same Top 10 with all four aggregations or E/V alone",9.45,4.32,2.65,.75,14,GRAY,True,align=PP_ALIGN.CENTER)

    s=slide();d.title(s,"Aggregation choices matter elsewhere—but not for the 6404 shortlist","Sensitivity audit")
    d.image(s,charts['sensitivity'],1.0,1.55,7.1)
    d.card(s,8.55,1.85,3.85,3.9,"F8FAFC","CBD5E1");d.text(s,"Interpretation",8.95,2.25,3,.35,17,NAVY,True)
    for i,t in enumerate(["BCB regressions: identical Top 10.","BCB improvements: persistent features dominate when mean/frequency enter.","Aggregation encodes a mechanistic preference, not a neutral preprocessing choice."]):d.bullet(s,t,8.95,2.95+i*.78,3.0,14,accent=BLUE)

    s=slide();d.title(s,"A good causal candidate aligns meaning, failure mode, and timing","Selection logic")
    items=[("MEANING","What tokens and contexts\nactivate the feature?",PURPLE),("FAILURE MODE","Which behavioral transition\nis it associated with?",RED),("TIMING","Does activation precede\nthe relevant decision?",CYAN)]
    for i,(a,b,c) in enumerate(items):
        x=.65+i*4.18;d.card(s,x,2.0,3.65,2.55,WHITE,c);d.text(s,a,x+.25,2.42,3.15,.35,16,c,True,align=PP_ALIGN.CENTER);d.text(s,b,x+.35,3.12,2.95,.75,17,NAVY,True,align=PP_ALIGN.CENTER)
    d.text(s,"Semantic clarity alone is not enough: a late feature may be a consequence or style marker.",1.0,5.35,11.3,.55,19,NAVY,True,align=PP_ALIGN.CENTER)

    s=slide();d.title(s,"First-activation position separates plausible causes from late markers","Temporal plausibility")
    d.image(s,charts['timing'],.65,1.35,8.7)
    d.card(s,9.65,1.7,2.75,4.45,"FFF7ED","FED7AA");d.text(s,"Example",10.0,2.05,2.0,.3,14,ORANGE,True);d.text(s,"Feature 12956",10.0,2.55,2.0,.3,18,NAVY,True);d.text(s,"Expected output / examples",10.0,3.05,2.0,.6,14,INK);d.text(s,"Semantically clear—often too late to explain the implementation decision.",10.0,4.0,2.0,1.1,15,GRAY,True)

    s=slide();d.title(s,"The pearl: many ‘improvements’ were simply cleaner stopping behavior","Mechanistic discovery")
    d.text(s,"119 / 215",.9,1.65,4.2,.8,39,ORANGE,True,align=PP_ALIGN.CENTER);d.text(s,"BigCodeBench improvements",.9,2.5,4.2,.35,15,GRAY,align=PP_ALIGN.CENTER);d.text(s,"55%",.9,3.2,4.2,.95,54,NAVY,True,align=PP_ALIGN.CENTER);d.text(s,"test/import contamination in the base output",.95,4.23,4.1,.65,17,INK,True,align=PP_ALIGN.CENTER)
    d.card(s,5.65,1.45,6.55,4.55,"172033","334155")
    code=" def task_func(...):\n     ... plausible solution ...\n-\n-# tests\n-from task_func import task_func"
    d.text(s,code,6.05,1.95,5.75,2.55,21,WHITE,False,"Aptos Mono")
    d.text(s,"The fine-tuned model often won by stopping—not by implementing a different algorithm.",6.05,4.85,5.75,.65,17,"CBD5E1",True)

    s=slide();d.title(s,"Feature 6404 matched the semantics of post-solution continuation","Feature interpretation")
    d.card(s,.75,1.65,5.2,4.45,"EFF6FF","93C5FD");d.text(s,"6404",1.15,2.0,1.5,.55,34,BLUE,True);d.text(s,"Dominant activation contexts",1.15,2.75,3.6,.3,14,GRAY,True)
    for i,t in enumerate(["expected","should","assumptions / caveats","comments and boilerplate"]):d.text(s,t,1.35,3.25+i*.52,3.6,.35,18,NAVY,True)
    d.card(s,6.35,1.65,5.95,4.45,WHITE,"CBD5E1");d.text(s,"67%",6.8,2.05,2.0,.8,42,BLUE,True,align=PP_ALIGN.CENTER);d.text(s,"contamination cases",6.8,2.85,2.0,.35,13,GRAY,align=PP_ALIGN.CENTER);d.text(s,"vs",9.0,2.25,.5,.4,18,GRAY,True,align=PP_ALIGN.CENTER);d.text(s,"31%",9.65,2.05,2.0,.8,42,GRAY,True,align=PP_ALIGN.CENTER);d.text(s,"other improvements",9.65,2.85,2.0,.35,13,GRAY,align=PP_ALIGN.CENTER);d.text(s,"Exploratory odds ratio ≈ 4.44",7.1,4.0,4.5,.45,19,NAVY,True,align=PP_ALIGN.CENTER);d.text(s,"Association generated the hypothesis; it did not establish causality.",7.0,4.8,4.7,.65,14,GRAY,align=PP_ALIGN.CENTER)

    s=slide();d.title(s,"A deliberately selected cohort turned the hypothesis into a test","Causal design")
    criteria=[("1","Base failed"),("2","Fine-tuned passed"),("3","Base generated tests/imports"),("4","Feature 6404 was naturally active")]
    for i,(num,txt) in enumerate(criteria):
        x=.7+i*3.05;d.card(s,x,1.85,2.6,2.1,WHITE,"CBD5E1");d.text(s,num,x+.9,2.15,.8,.65,30,BLUE,True,align=PP_ALIGN.CENTER);d.text(s,txt,x+.25,3.0,2.1,.55,15,NAVY,True,align=PP_ALIGN.CENTER)
    d.text(s,"80",4.65,4.55,1.3,.8,45,BLUE,True,align=PP_ALIGN.CENTER);d.text(s,"mechanistic cases",5.9,4.75,2.25,.35,19,NAVY,True);d.text(s,"Selected for mechanism—not an out-of-sample estimate.",3.0,5.75,7.3,.4,16,RED,True,align=PP_ALIGN.CENTER)

    s=slide();d.title(s,"We intervened only when 6404 entered the natural TopK","Online gated steering")
    stages=[("Current token","capture layer 16"),("Paired residuals","RMS normalize"),("DSTK100","ReLU + TopK 100"),("6404 active?","gate intervention"),("Next token","continue generation")]
    for i,(a,b) in enumerate(stages):
        x=.45+i*2.58;d.card(s,x,1.85,2.15,1.55,WHITE,"93C5FD" if i==3 else "CBD5E1");d.text(s,a,x+.15,2.15,1.85,.3,15,BLUE if i==3 else NAVY,True,align=PP_ALIGN.CENTER);d.text(s,b,x+.15,2.65,1.85,.35,11,GRAY,align=PP_ALIGN.CENTER)
        if i<4:d.text(s,"→",x+2.15,2.3,.4,.35,22,GRAY,True,align=PP_ALIGN.CENTER)
    d.card(s,1.4,4.25,10.55,1.3,"172033","334155");d.text(s,"h′ₜ = hₜ + α · z₆₄₀₄,ₜ · RMS(hₜ) · d₆₄₀₄,base",1.8,4.62,9.75,.45,22,WHITE,False,"Aptos Mono",align=PP_ALIGN.CENTER);d.text(s,"The intervention changes the next-token distribution; it does not rewrite earlier tokens.",2.0,5.92,9.3,.35,15,GRAY,align=PP_ALIGN.CENTER)

    s=slide();d.title(s,"One example: the baseline solved the task, then imported its own test","Token-level case")
    base,steer=read_task();
    d.card(s,.6,1.45,5.95,3.9,"FFF7ED","FDBA74");d.text(s,"BASELINE · FAIL",.95,1.78,2.6,.3,14,RED,True);base_code="return []\n\n# tests/test_task_func.py\nimport unittest\nfrom task_func import task_func"
    d.text(s,base_code,.95,2.35,5.1,1.65,18,INK,False,"Aptos Mono");d.text(s,"The plausible function body is followed by test/import contamination.",.95,4.65,5.1,.55,14,GRAY)
    d.card(s,6.8,1.45,5.95,3.9,"ECFDF5","6EE7B7");d.text(s,"6404 SUPPRESSION · PASS",7.15,1.78,3.5,.3,14,GREEN,True);steer_code="return min_subsequence\n\n# task_func({...})\n# task_func({...})"
    d.text(s,steer_code,7.15,2.35,5.1,1.65,18,INK,False,"Aptos Mono");d.text(s,"Official pass; no generated test module import.",7.15,4.65,5.1,.55,14,GRAY)
    d.image(s,charts['trace'],1.1,5.55,11.0,1.28);d.text(s,"Red marks: online steps where feature 6404 entered the active TopK in the steered trajectory.",2.0,6.93,9.3,.25,10,GRAY,align=PP_ALIGN.CENTER)
    d.notes(s,"Representative task: BigCodeBench/668. Show the exact full outputs in the appendix/source data if challenged; the slide uses a compact excerpt. Active gate steps were 8, 113, and 152 of 251 generated-token steps.")

    s=slide();d.title(s,"The intervention recovered official passes—not just prettier text","Main causal result")
    d.text(s,"19 / 80",.9,1.75,4.8,1.0,50,BLUE,True,align=PP_ALIGN.CENTER);d.text(s,"fail → official pass",.9,2.85,4.8,.45,19,NAVY,True,align=PP_ALIGN.CENTER)
    d.text(s,"41 / 80",7.0,1.75,4.8,1.0,50,ORANGE,True,align=PP_ALIGN.CENTER);d.text(s,"test/import contamination removed",7.0,2.85,4.8,.45,19,NAVY,True,align=PP_ALIGN.CENTER)
    d.card(s,1.35,4.15,10.65,1.3,"EFF6FF","BFDBFE");d.text(s,"α = −2   ·   TopK-gated   ·   exact code evaluated with extraction v4",1.75,4.58,9.85,.4,21,NAVY,True,align=PP_ALIGN.CENTER);d.text(s,"Cohort selected for this mechanism; do not read 23.75% as general benchmark gain.",2.0,5.9,9.3,.35,15,RED,True,align=PP_ALIGN.CENTER)

    s=slide();d.title(s,"Matched controls isolate the 6404 direction","Specificity")
    d.image(s,charts['controls'],.55,1.45,8.65)
    d.card(s,9.55,1.75,2.75,4.35,"F8FAFC","CBD5E1");d.text(s,"Controls matched",9.9,2.05,2.05,.3,15,NAVY,True)
    for i,t in enumerate(["decoder norm","activation support","temporal profile","online gate activity"]):d.bullet(s,t,9.9,2.62+i*.55,2.0,13,accent=GRAY)
    d.text(s,"vs sham",9.9,4.95,2.0,.25,12,GRAY,True);d.text(s,"pass p ≈ 3.8×10⁻⁶\ncleanup p ≈ 9.1×10⁻¹³",9.9,5.28,2.05,.6,13,BLUE,True)

    s=slide();d.title(s,"Steering is not a volume knob","Non-monotonicity")
    d.image(s,charts['curve'],.65,1.35,8.8)
    d.card(s,9.75,1.7,2.6,4.55,"FFF7ED","FED7AA");d.text(s,"Why jagged?",10.08,2.02,1.95,.3,16,ORANGE,True)
    for i,t in enumerate(["Token choice is discrete.","One token changes the future context.","Future gates then change.","Pass/fail is discontinuous.","One seed per task adds noise."]):d.bullet(s,t,10.08,2.6+i*.58,1.9,12,accent=ORANGE)

    s=slide();d.title(s,"The evidence supports a specific mechanism—not a universal correction vector","Calibrated conclusion")
    d.card(s,.75,1.55,5.85,4.8,"ECFDF5","6EE7B7");d.text(s,"SUPPORTED",1.15,1.95,2,.3,15,GREEN,True)
    for i,t in enumerate(["Interpretable association with post-solution boilerplate","Causal removal of test/import contamination","Specificity against sham and matched features","Official evaluation changes on exact extracted code"]):d.bullet(s,t,1.15,2.65+i*.72,4.8,15,accent=GREEN)
    d.card(s,6.9,1.55,5.65,4.8,"FEF2F2","FCA5A5");d.text(s,"NOT YET SUPPORTED",7.3,1.95,2.8,.3,15,RED,True)
    for i,t in enumerate(["Generalization beyond the selected cohort","Stable alpha peak across generation seeds","Every 6404 activation means test generation","A general pass/fail steering direction"]):d.bullet(s,t,7.3,2.65+i*.72,4.55,15,accent=RED)

    s=slide(NAVY);d.text(s,"A model difference became\na testable causal mechanism.",.85,1.05,8.1,1.5,38,WHITE,True);d.text(s,"Taxonomy → sparse features → semantics → timing → gated steering → controls",.9,3.0,10.8,.5,20,CYAN,True);d.text(s,"We did not find a universal feature of correctness.\nWe found a representation that helps explain why two models behave differently—and changing it changes behavior.",.9,4.15,10.9,1.1,20,"CBD5E1");d.text(s,"Next: multi-seed replication · negative cohort · held-out selection/evaluation · token-level mechanism audit",.9,6.15,11.2,.4,14,WHITE,True)

    # Appendix
    s=slide();d.title(s,"Appendix · exact screening interpretation","Method detail")
    d.text(s,"E/V",.9,1.65,2.3,.75,42,BLUE,True,align=PP_ALIGN.CENTER);d.text(s,"effect observed",3.3,1.82,2.8,.4,18,NAVY,True,align=PP_ALIGN.CENTER);d.text(s,"÷",6.0,1.78,.6,.45,24,GRAY,True,align=PP_ALIGN.CENTER);d.text(s,"SD of 200 permuted effects",6.75,1.82,4.0,.4,18,NAVY,True,align=PP_ALIGN.CENTER)
    d.card(s,1.0,3.0,11.25,2.5,WHITE,"CBD5E1");d.text(s,"Permutation procedure",1.4,3.35,2.6,.3,16,BLUE,True)
    for i,t in enumerate(["Keep feature deltas fixed; shuffle positive/control labels within benchmark × transition.","Recompute the mean contrast for every feature after each shuffle.","Use the feature-wise null SD as the denominator; compute nominal tail p separately.","With 200 permutations, minimum nominal p = 1/201 ≈ 0.00498; initial screen was not maxT-corrected."]):d.bullet(s,t,1.4,3.88+i*.42,10.1,13,accent=BLUE)

    out=OUT/"crosscoder_model_diffing_causal_steering_2026-08-10.pptx";d.save(out);return out


def main():
    export_derived_data()
    charts={
        'transitions':chart_transition_balance(),'taxonomy':chart_failure_taxonomy(),'ev':chart_ev_ranking(),
        'sensitivity':chart_sensitivity(),'timing':chart_semantic_timing(),'trace':chart_activation_trace(),
        'controls':chart_controls(),'curve':chart_dense_curve(),
    }
    out=build_deck(charts)
    print(out)


if __name__=="__main__":main()
