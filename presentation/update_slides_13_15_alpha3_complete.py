#!/usr/bin/env python3
from pathlib import Path
import io, re, zipfile
import matplotlib.pyplot as plt
from matplotlib.patches import FancyBboxPatch, Rectangle
from pptx import Presentation
from PIL import Image

REPO=Path('/home/lucas/crosscoder-recatcher-steering')
SRC=REPO/'MODEL DIFFS_CAUSAL_STEERING_RESULTS_2026-08-27.pptx'
DEST=REPO/'MODEL DIFFS_CAUSAL_STEERING_RESULTS_ALPHA3_COMPLETE_2026-08-31.pptx'
PDF=REPO/'MODEL DIFFS_CAUSAL_STEERING_RESULTS_ALPHA3_COMPLETE_2026-08-31.pdf'
OUT=REPO/'presentation/generated/alpha3_complete_slides_20260831';OUT.mkdir(parents=True,exist_ok=True)
BG='#F7F4EE';NAVY='#10233F';PURPLE='#7257C7';ORANGE='#EE8A2D';TEAL='#159C9C';MUTED='#687386';WHITE='#FFFFFF';GRID='#D8D9D8';RED='#C93948';GREEN='#338A63'

def setup(title,sub=''):
 f=plt.figure(figsize=(13.333,7.5),facecolor=BG);a=f.add_axes([0,0,1,1]);a.axis('off')
 a.text(.045,.955,'CANONICAL ERROR-FOCUSED α=3 SWEEP · COMPLETE READOUT',fontsize=8.5,color=TEAL,weight='bold',va='top')
 a.text(.045,.885,title,fontsize=23,color=NAVY,weight='bold',va='top')
 if sub:a.text(.045,.815,sub,fontsize=10.2,color=MUTED,va='top')
 a.text(.95,.035,'CrossCoder model diffing · canonical seed · 2026-08-31',ha='right',fontsize=7.5,color=MUTED)
 return f,a
def card(a,x,y,w,h,title,body,accent,ts=10,bs=8):
 a.add_patch(FancyBboxPatch((x,y),w,h,boxstyle='round,pad=.007,rounding_size=.011',facecolor=WHITE,edgecolor='#E1DED7'))
 a.add_patch(Rectangle((x,y+h-.01),w,.01,color=accent,lw=0));a.text(x+.016,y+h-.04,title,fontsize=ts,color=NAVY,weight='bold',va='top');a.text(x+.016,y+h-.085,body,fontsize=bs,color=MUTED,va='top',linespacing=1.3)
def save(f,n):
 p=OUT/f'{n}.png';f.savefig(p,dpi=175,facecolor=BG);plt.close(f);return p

# Slide 13
f,a=setup('THE COMPLETE RESULT SEPARATES A STRONGER AND A WEAKER REGIME','Official BigCodeBench 0.1.5 evaluation · α=3 · continuous last-token steering')
for x,title,body,col in [(.06,'DEEPSEEK', 'Features evaluated      35/35\nBaseline                 0/80\nFeature-task transitions 116\nUnique corrected tasks   16/80\n\nSeveral features produce repeated official corrections.',PURPLE),(.53,'CODELLAMA','Features evaluated      32/32\nBaseline                 0/50\nFeature-task transitions 6\nUnique corrected tasks   4/50\n\nSingle-feature corrections remain sparse.',ORANGE)]:card(a,x,.36,.41,.37,title,body,col,12,9.5)
a.text(.06,.235,'STRONGER / LOCALIZED',fontsize=11,color=PURPLE,weight='bold');a.text(.06,.185,'DeepSeek contains several multi-task causal candidates.',fontsize=9,color=NAVY)
a.text(.53,.235,'WEAKER / DIFFUSE',fontsize=11,color=ORANGE,weight='bold');a.text(.53,.185,'CodeLlama yields sparse, task-specific corrections.',fontsize=9,color=NAVY)
a.text(.06,.105,'COMPLETE · DeepSeek 35/35 and CodeLlama 32/32; canonical α=0 gates passed before steering.',fontsize=10,color=GREEN,weight='bold')
p13=save(f,13)

# Slide 14
rows=[(2468,7,34,-7.29,'paired local · #1'),(2621,6,28,-10.53,'paired global · #8'),(1078,6,26,5.72,'association global · #3'),(15235,6,23,-4.61,'association global · #8'),(14175,6,26,4.09,'association local · #1'),(2913,5,29,-10.95,'paired global · #1'),(6895,5,23,5.73,'association global · #2'),(596,5,26,4.99,'association global · #4'),(14820,5,23,-4.69,'association global · #6'),(1939,5,25,-4.65,'association global · #7')]
f,a=setup('DEEPSEEK: THE COMPLETE α=3 RANKING ADDS TWO TOP-5 FEATURES','Causal outcome first; |E/V| breaks ties among features with equal official passes')
ax=f.add_axes([.075,.18,.42,.57]);ys=list(range(10))[::-1];ax.barh(ys,[r[2] for r in rows],color='#D9D0F1',label='output changed');ax.barh(ys,[r[1] for r in rows],color=PURPLE,label='official pass');ax.set_yticks(ys,[str(r[0]) for r in rows]);ax.set_xlabel('tasks / 80');ax.spines[['top','right']].set_visible(False);ax.legend(frameon=False,fontsize=8)
card(a,.54,.225,.41,.51,'FINAL TOP 5','\n'.join(f'{i+1}. {fid}   {pa}/80 pass   {ch}/80 changed   E/V {ev:+.2f}' for i,(fid,pa,ch,ev,_) in enumerate(rows[:5])),PURPLE,11,8.5)
a.text(.54,.16,'1078 and 14175 enter after completion of all 35 arms.',fontsize=8.8,color=NAVY,weight='bold');a.text(.54,.11,'2913 is #6 at α=3, but later peaks at 11/80 at α=−5.',fontsize=8.8,color=RED,weight='bold')
p14=save(f,14)

# Slide 15
tasks=[('/316',17),('/435',15),('/166',13),('/496',11),('/937',10),('/349',7),('/459',7),('/7',7),('/634',6),('/695',6),('/940',5),('/188',5),('/666',3),('/1',2),('/576',1),('/823',1)]
f,a=setup('DEEPSEEK: 116 TRANSITIONS COLLAPSE TO 16 SUSCEPTIBLE TASKS','Feature-task outcomes are not independent causal successes')
ax=f.add_axes([.075,.17,.56,.59]);y=list(range(len(tasks)))[::-1];ax.barh(y,[n for _,n in tasks],color=PURPLE);ax.set_yticks(y,[t for t,_ in tasks],fontsize=7.5);ax.set_xlabel('features that corrected the task');ax.spines[['top','right']].set_visible(False);ax.grid(axis='x',color=GRID,alpha=.45)
card(a,.68,.46,.27,.28,'CONCENTRATION','/316   17 features\n/435   15 features\n/166   13 features\n/496   11 features\n/937   10 features',PURPLE,10,8.8)
card(a,.68,.20,.27,.19,'INTERPRETATION','The sweep finds causal levers, but a small set of tasks is broadly steering-susceptible. Random and orthogonal-sham controls remain necessary.',RED,9.5,7.8)
a.text(.075,.10,'16 unique corrections across 80 contamination-focused improvements.',fontsize=9,color=NAVY,weight='bold')
p15=save(f,15)

# Replace only slides 13–15.
prs=Presentation(SRC)
for number,p in [(13,p13),(14,p14),(15,p15)]:
 s=prs.slides[number-1]
 for sh in list(s.shapes):s.shapes._spTree.remove(sh._element)
 s.shapes.add_picture(str(p),0,0,width=prs.slide_width,height=prs.slide_height)
prs.save(DEST)

# Reconstruct PDF from the largest raster on every slide.
pages=[]
with zipfile.ZipFile(DEST) as z:
 for n in range(1,len(prs.slides)+1):
  rel=z.read(f'ppt/slides/_rels/slide{n}.xml.rels').decode();c=[]
  for target in re.findall(r'Target="\.\./media/([^"]+)"',rel):
   blob=z.read('ppt/media/'+target);c.append((len(blob),blob))
  if not c:raise RuntimeError(f'no image for slide {n}')
  pages.append(Image.open(io.BytesIO(max(c)[1])).convert('RGB'))
pages[0].save(PDF,save_all=True,append_images=pages[1:],resolution=150)
print(DEST);print(PDF);print('slides',len(prs.slides),'pdf_pages',len(pages))
