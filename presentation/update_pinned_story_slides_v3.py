#!/usr/bin/env python3
from pathlib import Path
from zipfile import ZipFile
import io,re
import matplotlib.pyplot as plt
from matplotlib.patches import FancyBboxPatch,Rectangle,Circle,FancyArrowPatch
from pptx import Presentation
from PIL import Image

R=Path('/home/lucas/crosscoder-recatcher-steering')
SRC=R/'MODEL DIFFS_CAUSAL_STEERING_RESULTS_ALPHA3_COMPLETE_V2_2026-08-31.pptx'
DEST=R/'MODEL DIFFS_CAUSAL_STEERING_RESULTS_STORY_VISUALS_V3_2026-08-31.pptx'
PDF=R/'MODEL DIFFS_CAUSAL_STEERING_RESULTS_STORY_VISUALS_V3_2026-08-31.pdf'
OUT=R/'presentation/generated/story_visuals_v3_20260831';OUT.mkdir(parents=True,exist_ok=True)
BG='#F7F4EE';NAVY='#10233F';PURPLE='#7257C7';ORANGE='#EE8A2D';TEAL='#159C9C';MUTED='#687386';WHITE='#FFFFFF';GRID='#D8D9D8';RED='#C93948';GREEN='#338A63';GOLD='#D4A72C';LIGHT='#EEEAE2'

def setup(title,sub='',section='CROSSCODER MODEL DIFFING'):
 f=plt.figure(figsize=(13.333,7.5),facecolor=BG);a=f.add_axes([0,0,1,1]);a.axis('off')
 a.text(.045,.955,section,fontsize=8.5,color=TEAL,weight='bold',va='top');a.text(.045,.885,title,fontsize=22,color=NAVY,weight='bold',va='top')
 if sub:a.text(.045,.815,sub,fontsize=10,color=MUTED,va='top')
 a.text(.95,.035,'CrossCoder model diffing · canonical readout · 2026-08-31',ha='right',fontsize=7.4,color=MUTED)
 return f,a
def card(a,x,y,w,h,title,body,accent,ts=10,bs=8):
 a.add_patch(FancyBboxPatch((x,y),w,h,boxstyle='round,pad=.007,rounding_size=.011',facecolor=WHITE,edgecolor='#E1DED7'))
 a.add_patch(Rectangle((x,y+h-.009),w,.009,color=accent,lw=0));a.text(x+.014,y+h-.035,title,fontsize=ts,color=NAVY,weight='bold',va='top');a.text(x+.014,y+h-.082,body,fontsize=bs,color=MUTED,va='top',linespacing=1.25)
def arrow(a,x1,y1,x2,y2,c=GRID,w=1.5):a.add_patch(FancyArrowPatch((x1,y1),(x2,y2),arrowstyle='-|>',mutation_scale=12,color=c,lw=w))
def save(f,n):p=OUT/f'{n:02d}.png';f.savefig(p,dpi=180,facecolor=BG);plt.close(f);return p

# 3 — two pairs, same same-text mechanism.
f,a=setup('ONE ANALYSIS FRAMEWORK, TWO VERY DIFFERENT MODEL CHANGES','Same token IDs turn each model pair into a shared sparse vocabulary')
for y,name,left,right,col,note in [(.49,'DEEPSEEK','BASE','FINE-TUNED',PURPLE,'specialization through fine-tuning'),(.20,'CODELLAMA','BASE','MERGED',ORANGE,'specialization through model merging')]:
 a.text(.055,y+.16,name,fontsize=11.5,color=col,weight='bold');a.text(.055,y+.115,note,fontsize=8.2,color=MUTED)
 card(a,.19,y,.115,.18,left+'\nlayer 16','residual stream',col,9,7.2);card(a,.19,y-.205,.115,.18,right+'\nlayer 16','residual stream',col,9,7.2)
 a.text(.326,y-.012,'same\ntoken IDs',ha='center',va='center',fontsize=7.2,color=col,weight='bold');arrow(a,.307,y+.075,.385,y+.015,col);arrow(a,.307,y-.115,.385,y-.015,col)
 card(a,.39,y-.07,.18,.23,'CROSSCODER','16,384 latents\nReLU + exact TopK-100',col,10,8)
 arrow(a,.575,y+.045,.635,y+.045,col)
 a.add_patch(FancyBboxPatch((.64,y-.07),.27,.23,boxstyle='round,pad=.007',facecolor=WHITE,edgecolor='#E1DED7'))
 for j in range(12):a.add_patch(Rectangle((.67+j*.017,y+.065),.010,.035,facecolor=col if j in (0,3,7,10) else '#E4E5E4',edgecolor='none'))
 a.text(.665,y+.025,'100 active features',fontsize=9.5,color=NAVY,weight='bold');a.text(.665,y-.025,'per evaluated token',fontsize=8.2,color=MUTED)
a.text(.055,.085,'Different specialization mechanisms; identical analytical question: which shared latent differences track—and control—behavior?',fontsize=9.4,color=NAVY,weight='bold')
p3=save(f,3)

# 4 — transition matrices, not only aggregate cards.
f,a=setup('THE TWO SPECIALIZATIONS MOVE PERFORMANCE IN OPPOSITE DIRECTIONS','Canonical extraction-v4 labels · 1,140 paired BigCodeBench tasks per model pair')
def matrix(x,title,col,base,var,bp,reg,imp,bf,varlabel):
 a.text(x,.735,title,fontsize=11.5,color=col,weight='bold');a.text(x,.69,f'BASE {base}/1140  →  {varlabel} {var}/1140',fontsize=9.5,color=NAVY,weight='bold')
 a.text(x,.615,'BASE',fontsize=7.5,color=MUTED,weight='bold');a.text(x+.135,.615,varlabel,fontsize=7.5,color=MUTED,weight='bold')
 rows=[('PASS','PASS',bp,'both pass',LIGHT),('PASS','FAIL',reg,'REGRESSION',col),('FAIL','PASS',imp,'IMPROVEMENT',col),('FAIL','FAIL',bf,'both fail',LIGHT)]
 for i,(l,r,n,lab,c) in enumerate(rows):
  y=.555-i*.09;a.add_patch(FancyBboxPatch((x,y),.08,.06,boxstyle='round,pad=.004',facecolor=WHITE,edgecolor=GRID));a.text(x+.04,y+.03,l,ha='center',va='center',fontsize=7.5,color=NAVY,weight='bold')
  arrow(a,x+.082,y+.03,x+.13,y+.03,c if c!=LIGHT else GRID,1.4)
  a.add_patch(FancyBboxPatch((x+.14,y),.08,.06,boxstyle='round,pad=.004',facecolor=WHITE,edgecolor=GRID));a.text(x+.18,y+.03,r,ha='center',va='center',fontsize=7.5,color=NAVY,weight='bold')
  a.add_patch(FancyBboxPatch((x+.24,y),.17,.06,boxstyle='round,pad=.004',facecolor=c,edgecolor=c));a.text(x+.255,y+.03,str(n),va='center',fontsize=10,color=NAVY if c==LIGHT else WHITE,weight='bold');a.text(x+.31,y+.03,lab,va='center',fontsize=7.3,color=NAVY if c==LIGHT else WHITE,weight='bold')
matrix(.06,'DEEPSEEK',PURPLE,268,404,189,79,215,657,'FINE-TUNED')
matrix(.54,'CODELLAMA',ORANGE,314,27,23,291,4,822,'MERGED')
a.add_patch(FancyBboxPatch((.06,.105),.88,.09,boxstyle='round,pad=.008',facecolor=WHITE,edgecolor=GRID));a.text(.08,.157,'RECATCHER CONTEXT',fontsize=8,color=TEAL,weight='bold');a.text(.08,.12,'Specialization can change aggregate capability while creating localized one-sided transitions. ReCatcher used a different 10-generation protocol; this is contextual, not a numerical replication.',fontsize=8.2,color=NAVY)
p4=save(f,4)

# 5 — taxonomy as distributions and explicit cohort selection.
f,a=setup('AN OUTCOME LABEL IS NOT A MECHANISM','Taxonomy separates incompatible causal hypotheses before screening')
def taxonomy(x,title,total,segments,selected,col,foot):
 a.text(x,.72,title,fontsize=11.2,color=col,weight='bold');a.text(x,.675,f'{total} one-sided transitions',fontsize=8.5,color=MUTED)
 y=.56;start=x
 greys=['#D9DBDC','#C7CACE','#B6BBC0','#A6ACB3']
 for i,(label,n) in enumerate(segments):
  w=.39*n/total;c=col if label==selected else greys[i]
  a.add_patch(Rectangle((start,y),w,.075,facecolor=c,edgecolor=BG));
  if w>.045:a.text(start+w/2,y+.037,str(n),ha='center',va='center',fontsize=8,color=WHITE if label==selected else NAVY,weight='bold')
  start+=w
 for i,(label,n) in enumerate(segments):
  yy=.49-i*.045;a.add_patch(Rectangle((x,yy),.012,.018,facecolor=col if label==selected else greys[i],edgecolor='none'));a.text(x+.02,yy+.009,f'{label} · {n}',va='center',fontsize=7.7,color=NAVY,weight='bold' if label==selected else 'normal')
 card(a,x,.18,.39,.13,'SCREENING COHORT',foot,col,9.2,8)
taxonomy(.06,'DEEPSEEK · IMPROVEMENTS',215,[('test/import contamination',119),('wrong output or logic',48),('missing name/import',16),('other',32)],'test/import contamination',PURPLE,'80 contamination cases selected for the focused analysis')
taxonomy(.54,'CODELLAMA · REGRESSIONS',291,[('API/type mismatch',120),('wrong logic/runtime',50),('edge case/exception',43),('generation/import/syntax/other',78)],'wrong logic/runtime',ORANGE,'50 wrong-logic/runtime cases selected; 4 improvements excluded')
a.text(.06,.105,'Primary category = most salient observed error, not a unique root cause. CodeLlama failures are especially heterogeneous and often multilabel.',fontsize=8.8,color=RED,weight='bold')
p5=save(f,5)

# 19 — completed alpha=3 outcome distributions.
f,a=setup('THE α=3 SWEEP REVEALS BOTH MODEL AND TASK SENSITIVITY','Complete error-focused sweep · official BigCodeBench 0.1.5 evaluation')
ds={0:1,1:5,2:7,3:8,4:3,5:6,6:4,7:1};cl={0:26,1:6,2:0,3:0,4:0,5:0,6:0,7:0}
for pos,title,hist,col,nfeat in [([.065,.30,.38,.43],'DeepSeek · 35 features',ds,PURPLE,35),([.535,.30,.38,.43],'CodeLlama · 32 features',cl,ORANGE,32)]:
 ax=f.add_axes(pos);xs=list(range(8));ax.bar(xs,[hist[x] for x in xs],color=col,width=.72);ax.set_title(title,loc='left',color=NAVY,weight='bold',fontsize=10);ax.set_xlabel('official fail→pass tasks per feature',fontsize=8);ax.set_ylabel('features',fontsize=8);ax.set_xticks(xs);ax.grid(axis='y',color=GRID,alpha=.65);ax.spines[['top','right']].set_visible(False);ax.tick_params(labelsize=8)
a.text(.065,.215,'116 feature–task transitions  →  16/80 unique tasks',fontsize=9.2,color=PURPLE,weight='bold');a.text(.535,.215,'6 feature–task transitions  →  4/50 unique tasks',fontsize=9.2,color=ORANGE,weight='bold')
a.text(.065,.165,'Most susceptible: /316 corrected by 17 features · /435 by 15 · /166 by 13',fontsize=7.9,color=NAVY);a.text(.535,.165,'Most susceptible: /119 corrected by 3 of the 6 successful features',fontsize=7.9,color=NAVY)
a.add_patch(FancyBboxPatch((.06,.075),.86,.06,boxstyle='round,pad=.007',facecolor=WHITE,edgecolor=GRID));a.text(.49,.105,'Controls are always required; repeated correction of the same tasks makes susceptibility—and the need for specificity tests—especially visible.',ha='center',va='center',fontsize=8.7,color=RED,weight='bold')
p19=save(f,19)

prs=Presentation(SRC)
for n,p in [(3,p3),(4,p4),(5,p5),(19,p19)]:
 s=prs.slides[n-1]
 for sh in list(s.shapes):s.shapes._spTree.remove(sh._element)
 s.shapes.add_picture(str(p),0,0,width=prs.slide_width,height=prs.slide_height)
prs.save(DEST)
with ZipFile(SRC) as z:
 pages=[]
 for n in range(1,len(prs.slides)+1):
  if n in (3,4,5,19):pages.append(Image.open(OUT/f'{n:02d}.png').convert('RGB'));continue
  rel=z.read(f'ppt/slides/_rels/slide{n}.xml.rels').decode();slide=z.read(f'ppt/slides/slide{n}.xml').decode();mp={i:t for i,t in re.findall(r'Id="([^\"]+)"[^>]+Target="([^\"]+)"',rel)};vals=[]
  for rid in re.findall(r'r:embed="([^\"]+)"',slide):
   t=mp.get(rid,'')
   if t.startswith('../media/'):
    d=z.read('ppt/'+t[3:]);vals.append((len(d),d))
  pages.append(Image.open(io.BytesIO(max(vals)[1])).convert('RGB'))
pages[0].save(PDF,save_all=True,append_images=pages[1:],resolution=180)
print(DEST);print(PDF);print('slides',len(prs.slides),'pdf_pages',len(pages))
