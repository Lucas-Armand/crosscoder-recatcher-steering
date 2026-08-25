#!/usr/bin/env python3
"""Create a versioned BBASV discussion deck from MODEL DIFFS.pptx."""
from pathlib import Path
import argparse,csv,json,re,shutil,subprocess
import numpy as np
import matplotlib.pyplot as plt
from matplotlib.backends.backend_pdf import PdfPages
from matplotlib.patches import FancyBboxPatch,Rectangle
from pptx import Presentation

BG="#F6F3EC"; NAVY="#10233F"; TEAL="#169C9C"; CORAL="#E5655B"; GOLD="#E3A62F"
MUTED="#687386"; WHITE="#FFFFFF"; GRID="#D7D9D8"; GREY="#D6D6D2"

def setup(title,sub=""):
    f=plt.figure(figsize=(13.333,7.5),facecolor=BG); a=f.add_axes([0,0,1,1]); a.axis("off")
    a.text(.045,.955,"SCREENING → STANDARDIZED CAUSAL SWEEP → BBASV",fontsize=8.5,color=TEAL,weight="bold",va="top")
    a.text(.045,.885,title,fontsize=24,color=NAVY,weight="bold",va="top")
    if sub:a.text(.045,.815,sub,fontsize=10.5,color=MUTED,va="top")
    a.text(.95,.035,"CrossCoder model diffing · discussion draft · 2026-08-25",ha="right",fontsize=7.5,color=MUTED)
    return f,a

def card(a,x,y,w,h,title,body="",accent=TEAL,face=WHITE,ts=11,bs=8.6):
    a.add_patch(FancyBboxPatch((x,y),w,h,boxstyle="round,pad=0.007,rounding_size=.011",facecolor=face,edgecolor="#E2E0D9"))
    a.add_patch(Rectangle((x,y+h-.01),w,.01,color=accent,linewidth=0))
    a.text(x+.018,y+h-.045,title,fontsize=ts,color=NAVY,weight="bold",va="top")
    if body:a.text(x+.018,y+h-.092,body,fontsize=bs,color=MUTED,va="top",linespacing=1.35)

def save(f,p,pdf):
    f.savefig(p,dpi=170,facecolor=BG);pdf.savefig(f,facecolor=BG);plt.close(f)

def tab(a,data,heads,widths,x=.05,y=.16,w=.90,h=.57,color=NAVY,hi=None,fs=7.7):
    rh=h/(len(data)+1); xx=[x]
    for q in widths:xx.append(xx[-1]+w*q)
    a.add_patch(Rectangle((x,y+h-rh),w,rh,color=color))
    for j,v in enumerate(heads):a.text(xx[j]+.005,y+h-rh/2,v,color=WHITE,fontsize=fs,weight="bold",va="center")
    for i,row in enumerate(data):
        yy=y+h-(i+2)*rh;a.add_patch(Rectangle((x,yy),w,rh,color=WHITE if i%2==0 else "#F0EEE8"))
        for j,v in enumerate(row):a.text(xx[j]+.005,yy+rh/2,str(v),color=TEAL if j==hi else NAVY,fontsize=fs,weight="bold" if j==hi else "normal",va="center")

def read10(p):
    with open(p,newline="") as f:return list(csv.DictReader(f))[:10]

def counts(p):
    z=[]
    for q in p.glob("evaluations/bigcodebench__f*_eval.json"):
        if re.search(r"__f\d+_(?:pos|neg)3_eval",q.name):z.append(json.load(open(q))["passed"])
    return z

def build(repo,out):
    out.mkdir(parents=True,exist_ok=True);imgs=[]
    with PdfPages(out/"bbasv_story_block.pdf") as pdf:
        f,a=setup("THE SAME EIGHT-CELL DESIGN—WITH UNEQUAL EVIDENCE","Each model crosses transition type × comparison design × temporal scope")
        for px,name,col in [(.045,"DEEPSEEK",TEAL),(.515,"CODELLAMA",CORAL)]:
            a.text(px,.755,name,fontsize=13,color=col,weight="bold")
            for trans,yy in [("REGRESSION",.46),("IMPROVEMENT",.18)]:
                off=name=="CODELLAMA" and trans=="IMPROVEMENT"
                card(a,px,yy,.43,.245,trans,"",GREY if off else col,"#E5E4E0" if off else WHITE,10)
                for j,(u,v) in enumerate([("ASSOCIATION","global · local"),("PAIRED MODEL Δ","global · local")]):
                    x=px+.018+j*.205;a.add_patch(Rectangle((x,yy+.035),.19,.105,facecolor="#EFEDEA" if off else BG,edgecolor=GRID))
                    a.text(x+.012,yy+.105,u,fontsize=8,color=MUTED if off else NAVY,weight="bold");a.text(x+.012,yy+.065,v,fontsize=8,color=MUTED)
                if off:a.text(px+.215,yy+.188,"ONLY 4 POSITIVE TRANSITIONS",ha="center",fontsize=8,color=CORAL,weight="bold")
        a.text(.045,.095,"Local = window around the first divergence—not a semantic error category.",fontsize=9,color=NAVY,weight="bold")
        a.text(.515,.095,"CodeLlama improvement cells were excluded from candidate selection.",fontsize=9,color=CORAL,weight="bold")
        p=out/"slide_23_eight_cells_by_model.png";save(f,p,pdf);imgs.append(p)

        specs=[
          ("A SCREEN PRODUCES A RANKED, AUDITABLE SHORTLIST","DeepSeek · improvement association · global · base-enriched",repo/"reports/eight_cell_screening_dstk100_v1/improvement_association_global_base_enriched.csv",TEAL),
          ("THE SAME OUTPUT FORMAT APPLIES TO CODELLAMA","CodeLlama · regression association · global · absolute ranking",repo/"reports/eight_cell_screening_codellama_base_merged_v1/regression_association_global_absolute.csv",CORAL)]
        for k,(title,sub,path,col) in enumerate(specs):
            f,a=setup(title,sub);data=[]
            for r in read10(path):data.append([r["rank"],r["feature_id"],r["summary"],f'{float(r["effect"]):+.3f}',f'{float(r["ev"]):+.2f}',r["support"]])
            tab(a,data,["Rank","Feature","Summary","Effect","E/V","Support"],[.09,.14,.25,.16,.16,.20],.07,.18,.86,.56,col,4,8.5)
            a.text(.07,.115,"E/V = observed contrast ÷ SD of the label-permutation null",fontsize=10,color=NAVY,weight="bold")
            a.text(.07,.075,"Signal-to-noise score—not a z-score, p-value, corrected significance test, or causal estimate.",fontsize=9,color=MUTED)
            p=out/f"screen_example_{k}.png";save(f,p,pdf);imgs.append(p)

        f,a=setup("SCREEN FIRST; INTERVENE AT A COMMON OPERATING POINT","The associative stage makes causal testing computationally tractable")
        for x,big,small,col in [(.10,"16,384","latents per CrossCoder",NAVY),(.40,"≈30–40","screened candidates/model",TEAL),(.70,"≈0.2%","of the latent space tested",CORAL)]:
            a.text(x,.62,big,fontsize=36,color=col,weight="bold",ha="center");a.text(x,.55,small,fontsize=9,color=MUTED,ha="center")
        for x1,x2 in [(.18,.30),(.48,.60)]:a.annotate("",xy=(x2,.60),xytext=(x1,.60),arrowprops=dict(arrowstyle="->",lw=2,color=TEAL))
        card(a,.07,.20,.39,.25,"WHY |α| = 3?","Earlier dose-response experiments suggested a practical exploratory point: visible trajectory changes below the most disruptive high-dose regime.",GOLD)
        card(a,.54,.20,.39,.25,"WHY ONE COMMON MAGNITUDE?","A standardized intervention compares dozens of candidates before full dose curves, reverse-model tests, and controls.",TEAL)
        a.text(.07,.115,"|α|=3 is an exploratory operating point—not an independently validated optimum.",fontsize=10,color=CORAL,weight="bold")
        p=out/"screen_to_alpha3_funnel.png";save(f,p,pdf);imgs.append(p)

        ds=counts(repo/"runs/semantic_top10_dstk100_alpha3_v1");cl=counts(repo/"runs/semantic_top10_codellama_alpha3_v1")
        f,a=setup("THE α=3 RESPONSE DISTRIBUTIONS LOOK VERY DIFFERENT","Official net pass change relative to each cohort baseline")
        for rect,vals,base,title,col in [([.07,.25,.40,.46],ds,15,"DeepSeek · 80 contamination tasks",TEAL),([.54,.25,.40,.46],cl,0,"CodeLlama · 50 logic/runtime tasks",CORAL)]:
            q=f.add_axes(rect,facecolor=WHITE);d=[v-base for v in vals];lo,hi=min(d),max(d)
            q.hist(d,bins=np.arange(lo-.5,hi+1.5),color=col,edgecolor=WHITE,rwidth=.88);q.axvline(0,color=NAVY,lw=1.5,ls="--")
            q.set_xticks(range(lo,hi+1));q.set_xlabel("Net official pass change");q.set_ylabel("Features");q.set_title(title,loc="left",fontsize=11,color=NAVY,weight="bold");q.spines[["top","right"]].set_visible(False);q.grid(axis="y",color=GRID,alpha=.6)
        a.text(.07,.14,"DeepSeek: best +4; many positive arms",fontsize=10,color=TEAL,weight="bold")
        a.text(.54,.14,"CodeLlama: most arms 0; best +1",fontsize=10,color=CORAL,weight="bold")
        a.text(.07,.085,"Consistent with more heterogeneous/distributed tested CodeLlama failures—not proof of greater network-wide damage.",fontsize=9,color=NAVY)
        p=out/"alpha3_histograms.png";save(f,p,pdf);imgs.append(p)

        f,a=setup("DID WE FIND CAUSAL FEATURES—OR STEERABLE TASKS?","Success concentration makes placebo and alternative-latent controls essential")
        card(a,.05,.46,.42,.28,"DEEPSEEK /1030","Corrected by 29 / 30 tested features.\nA broadly sensitive task can inflate many candidate directions.",TEAL)
        card(a,.53,.46,.42,.28,"CODELLAMA /119","Responsible for 10 / 11 positive arms.\nA fail→pass transition alone cannot establish specificity.",CORAL)
        card(a,.05,.17,.27,.20,"TARGET FEATURE","Behavior-aligned latent selected by screening.",TEAL)
        card(a,.365,.17,.27,.20,"RANDOM LATENTS","Not selected using screening or prior α=3 outcomes.",GOLD)
        card(a,.68,.17,.27,.20,"ORTHOGONAL SHAMS","Random directions with controlled norm and geometry.",MUTED)
        a.text(.05,.095,"BBASV tests superiority to generic perturbations and coherent reversal across models.",fontsize=10,color=NAVY,weight="bold")
        p=out/"task_sensitivity_controls.png";save(f,p,pdf);imgs.append(p)

        f,a=setup("TEN CANDIDATES ADVANCE TO BBASV","Official α=3 outcome first; prior applicable screening rank breaks causal ties")
        data=[
          ["DeepSeek","3048","paired global","active frac.","−9.97","19/80","+4","base − / FT +"],
          ["DeepSeek","13801","paired global","max","−10.73","18/80","+3","base − / FT +"],
          ["DeepSeek","7828","association global","mean","−7.01","18/80","+3","base − / FT +"],
          ["DeepSeek","15669","paired global","max","−10.06","18/80","+3","base − / FT +"],
          ["DeepSeek","13191","paired global","max","−11.03","17/80","+2","base − / FT +"],
          ["CodeLlama","13147","association global","mean","+5.24","1/50","+1","merged − / base +"],
          ["CodeLlama","12253","paired global","max","−12.73","1/50","+1","merged + / base −"],
          ["CodeLlama","10570","association global","mean","+4.60","1/50","+1","merged − / base +"],
          ["CodeLlama","14359","association global","mean","−4.46","1/50","+1","merged + / base −"],
          ["CodeLlama","2310","association global","active frac.","−4.30","1/50","+1","merged + / base −"]]
        tab(a,data,["Model","Feature","Best prior screen","Summary","E/V","α=3","Net Δ","BBASV direction"],[.12,.09,.20,.13,.10,.10,.09,.17],.035,.16,.93,.57,NAVY,None,7.5)
        a.text(.04,.105,"Behavior-aligned direction in the worse model + reversed direction in the better model + random latents + orthogonal shams.",fontsize=9,color=NAVY,weight="bold")
        a.text(.04,.065,"Screen support belongs to the broad transition population; causal outcomes use narrower semantic cohorts.",fontsize=8.5,color=CORAL)
        p=out/"bbasv_top10_selection.png";save(f,p,pdf);imgs.append(p)
    return imgs

def remove(prs,i):
    sid=prs.slides._sldIdLst[i];prs.part.drop_rel(sid.rId);del prs.slides._sldIdLst[i]

def main():
    ap=argparse.ArgumentParser();ap.add_argument("--repo",type=Path,default=Path("."));ap.add_argument("--source",type=Path,default=Path("MODEL DIFFS.pptx"));ap.add_argument("--dest",type=Path,required=True);ap.add_argument("--output-dir",type=Path,required=True)
    a=ap.parse_args();imgs=build(a.repo,a.output_dir);shutil.copy2(a.source,a.dest);prs=Presentation(a.dest)
    blank=prs.slide_layouts[6];ids=[]
    for p in imgs:
        s=prs.slides.add_slide(blank);s.shapes.add_picture(str(p),0,0,width=prs.slide_width,height=prs.slide_height);ids.append(prs.slides._sldIdLst[-1])
    remove(prs,22);remove(prs,22)
    for sid in reversed(ids):prs.slides._sldIdLst.remove(sid);prs.slides._sldIdLst.insert(22,sid)
    prs.save(a.dest)

if __name__=="__main__":main()
